from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from upload_tools import upload_file
import io
import logging
from typing import List, Dict, Any
from template_utils import find_pptx_templates

TITLE_LAYOUT = 2
SECTION_LAYOUT = 7
CONTENT_LAYOUT = 4

# Create a logger
logger = logging.getLogger(__name__)


def load_templates():
    """Resolve presentation templates (4:3, 16:9) from custom/default template dirs.

    Returns: tuple[str|None, str|None] -> (path_4_3, path_16_9)
    """
    t43, t169 = find_pptx_templates()
    if not t43 or not t169:
        logger.info("One or more PPT templates missing; will fall back to PowerPoint defaults where needed")
    return t43, t169


class PowerpointPresentation:

    def __init__(self, slides: List[Dict[str, Any]], format: str):
        """Initialize PowerPoint presentation with slides and format"""

        logger.info(f"Initializing PowerPoint: slides={len(slides)}, format={format}")
        # Validate input
        if not slides:
            raise ValueError("At least one slide is required")

        # Loads templates
        self.template_regular, self.template_wide = load_templates()
        logger.debug(f"Selected templates -> 4:3={self.template_regular}, 16:9={self.template_wide}")

        # Create presentation based on the format used
        try:
            if format == "4:3":
                if self.template_regular:
                    self.presentation = Presentation(self.template_regular)
                else:
                    self.presentation = Presentation()  # Use default template
            elif format == "16:9":
                if self.template_wide:
                    self.presentation = Presentation(self.template_wide)
                else:
                    self.presentation = Presentation()  # Use default template
            else:
                logger.warning(f"Unknown format '{format}', defaulting to 4:3")
                if self.template_regular:
                    self.presentation = Presentation(self.template_regular)
                else:
                    self.presentation = Presentation()  # Use default template
        except Exception as e:
            logger.error(f"Failed to load template: {e}")
            logger.info("Falling back to default PowerPoint template")
            self.presentation = Presentation()  # Fallback to default template

        # Remove default slide if it exists
        if len(self.presentation.slides) > 0:
            logger.debug("Removing default first slide from new presentation")
            slide_to_remove = self.presentation.slides[0]
            rId = self.presentation.slides.element.remove(slide_to_remove.element)

        # Create slides
        self._create_slides(slides)

    def _create_slides(self, slides: List[Dict[str, Any]]):
        """Create all slides from the slides data"""
        logger.info(f"Creating {len(slides)} slides")
        for i, slide in enumerate(slides):
            try:
                slide_type = slide.get("slide_type")
                logger.debug(f"Creating slide index={i}, type={slide_type}, title={slide.get('slide_title','')}")

                if slide_type == "content":
                    self.create_content_slide(slide)
                elif slide_type == "section":
                    self.create_section_slide(slide)
                elif slide_type == "title":
                    self.create_title_slide(slide)
                else:
                    logger.warning(f"Unknown slide type '{slide_type}' for slide {i}, skipping")

            except Exception as e:
                logger.error(f"Failed to create slide {i}: {e}")
                raise ValueError(f"Error creating slide {i}: {str(e)}")

    def create_title_slide(self, slide: Dict[str, Any]):
        """Create a title slide"""
        try:
            title_layout = self.presentation.slide_layouts[TITLE_LAYOUT]
            title_slide = self.presentation.slides.add_slide(title_layout)
            logger.debug("Added title slide")

            # Set title
            if len(title_slide.placeholders) > 0:
                title_slide.placeholders[0].text = slide.get("slide_title", "")

            # Set author
            if len(title_slide.placeholders) > 1:
                title_slide.placeholders[1].text = slide.get("author", "")

        except Exception as e:
            logger.error(f"Failed to create title slide: {e}")
            raise

    def create_section_slide(self, slide: Dict[str, Any]):
        """Create a section slide"""
        try:
            section_layout = self.presentation.slide_layouts[SECTION_LAYOUT]
            section_slide = self.presentation.slides.add_slide(section_layout)
            logger.debug("Added section slide")

            # Set title
            if len(section_slide.placeholders) > 0:
                section_slide.placeholders[0].text = slide.get("slide_title", "")

        except Exception as e:
            logger.error(f"Failed to create section slide: {e}")
            raise

    def create_content_slide(self, slide: Dict[str, Any]):
        """Create a content slide with bullet points"""
        try:
            content_layout = self.presentation.slide_layouts[CONTENT_LAYOUT]
            content_slide = self.presentation.slides.add_slide(content_layout)
            logger.debug("Added content slide")

            # Set title
            if len(content_slide.placeholders) > 0:
                content_slide.placeholders[0].text = slide.get("slide_title", "")

            # Add content
            slide_text = slide.get("slide_text", [])
            if slide_text and len(content_slide.placeholders) > 1:
                # Clear existing text
                content_slide.placeholders[1].text = ""

                # Add first paragraph
                first_item = slide_text[0]
                content_slide.placeholders[1].text_frame.paragraphs[0].text = first_item.get("text", "")
                content_slide.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                content_slide.placeholders[1].text_frame.paragraphs[0].level = max(0, int(first_item.get("indentation_level", 1)) - 1)

                # Add remaining paragraphs
                for paragraph_data in slide_text[1:]:
                    p = content_slide.placeholders[1].text_frame.add_paragraph()
                    p.text = paragraph_data.get("text", "")
                    p.alignment = PP_ALIGN.LEFT
                    p.level = max(0, int(paragraph_data.get("indentation_level", 1)) - 1)

        except Exception as e:
            logger.error(f"Failed to create content slide: {e}")
            raise

    def save(self) -> io.BytesIO:
        """Save presentation to BytesIO object"""
        try:
            logger.info("Saving PowerPoint to memory buffer")
            file_like_object = io.BytesIO()
            self.presentation.save(file_like_object)
            file_like_object.seek(0)
            return file_like_object
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            raise

def create_presentation(slides: List[Dict[str, Any]], format: str = "4:3") -> str:
    """Creates new presentation."""

    try:
        # Validate input
        if not slides:
            raise ValueError("No slides provided")

        logger.info(f"Starting create_presentation: slides={len(slides)}, format={format}")

        # Create presentation
        presentation = PowerpointPresentation(slides, format)

        # Save presentation
        file_object = presentation.save()

        # Upload presentation
        text = upload_file(file_object, "pptx")
        file_object.close()

        logger.info("PowerPoint upload completed")
        # Return presentation link
        return text

    except Exception as e:
        logger.error(f"Failed to create presentation: {e}")
        raise
