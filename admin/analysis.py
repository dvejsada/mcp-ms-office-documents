"""Analyse uploaded template assets to drive the admin UI.

For a ``.docx`` we report the ``{{placeholders}}`` and ``{{#if}}`` conditionals
the renderer will act on, the paragraph styles the document actually defines
(so the style-mapping dropdowns can be populated), and which of the styles the
renderer relies on are missing. For an email ``.html`` we report the Mustache
variables and sections.

The reconciliation helper compares detected placeholders/conditionals against a
template's declared ``args`` so the UI can offer to add missing args and warn
about orphans.

This module is import-light and has no FastHTML dependency so it can be unit
tested on its own.
"""
from __future__ import annotations

import io
import logging
import re
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

from docx import Document as DocxDocument
from docx.oxml.ns import qn
from lxml import etree

from docx_tools.dynamic_docx_tools import PLACEHOLDER_PATTERN
from docx_tools.conditionals import parse_marker

logger = logging.getLogger(__name__)

# Styles the markdown renderer applies by name; warn when a template lacks them.
REQUIRED_DOCX_STYLES = [
    "Heading 1", "Heading 2", "Heading 3", "Heading 4", "Heading 5", "Heading 6",
    "List Bullet", "List Bullet 2", "List Bullet 3",
    "List Number", "List Number 2", "List Number 3",
    "Quote", "Table Grid", "Normal",
]

# Email standard fields are injected automatically by the email tool, so they
# never need a user-declared arg.
EMAIL_RESERVED_VARS = {"subject", "to", "cc", "bcc", "file_name"}

# Mustache section / variable patterns for HTML email templates.
_MUSTACHE_SECTION = re.compile(r"\{\{\s*[#^]\s*([a-zA-Z_][a-zA-Z0-9_]*)\s*\}\}")
_MUSTACHE_VAR = re.compile(r"\{\{\{?\s*([a-zA-Z_][a-zA-Z0-9_]*)\s*\}?\}\}")


@dataclass
class Analysis:
    """Result of analysing a template asset."""
    kind: str
    placeholders: List[str] = field(default_factory=list)
    conditionals: List[str] = field(default_factory=list)
    conditionals_balanced: bool = True
    styles_present: List[str] = field(default_factory=list)
    missing_required_styles: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)


def _ordered_unique(items) -> List[str]:
    """De-duplicate *items* preserving first-seen order."""
    seen = set()
    out = []
    for it in items:
        if it not in seen:
            seen.add(it)
            out.append(it)
    return out


def _iter_docx_paragraph_texts(doc: DocxDocument):
    """Yield combined text of every paragraph in body, tables, headers, footers."""
    def _walk_container(paragraphs, tables):
        for p in paragraphs:
            yield p.text
        for t in tables:
            for row in t.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        yield p.text

    yield from _walk_container(doc.paragraphs, doc.tables)
    for section in doc.sections:
        parts = [section.header, section.footer]
        if section.different_first_page_header_footer:
            parts += [section.first_page_header, section.first_page_footer]
        for part in parts:
            if part is None:
                continue
            yield from _walk_container(part.paragraphs, part.tables)


def _iter_docx_body_marker_texts(doc: DocxDocument):
    """Yield body-level paragraph texts (for conditional balance checking)."""
    p_tag = qn("w:p")
    t_tag = qn("w:t")
    for elem in doc.element.body:
        if elem.tag == p_tag:
            # Concatenate the paragraph's run text straight from the XML, avoiding
            # python-docx's internal Paragraph(elem, None) constructor.
            yield "".join(t.text or "" for t in elem.iter(t_tag))


def analyze_docx(data: bytes) -> Analysis:
    """Analyse a ``.docx`` given its bytes."""
    analysis = Analysis(kind="docx")
    try:
        doc = DocxDocument(io.BytesIO(data))
    except Exception as e:
        analysis.warnings.append(f"Could not open as a Word document: {e}")
        return analysis

    placeholders: List[str] = []
    for text in _iter_docx_paragraph_texts(doc):
        if "{{" in text:
            placeholders.extend(PLACEHOLDER_PATTERN.findall(text))
    analysis.placeholders = _ordered_unique(placeholders)

    # Conditionals + balance (body-level, matching the renderer's scope).
    conditionals: List[str] = []
    depth = 0
    balanced = True
    for text in _iter_docx_body_marker_texts(doc):
        marker = parse_marker(text)
        if marker is None:
            continue
        if marker.kind == "open":
            conditionals.append(marker.name)
            depth += 1
        else:
            depth -= 1
            if depth < 0:
                balanced = False
                depth = 0
    if depth != 0:
        balanced = False
    analysis.conditionals = _ordered_unique(conditionals)
    analysis.conditionals_balanced = balanced
    if not balanced:
        analysis.warnings.append(
            "Unbalanced {{#if}}/{{/if}} markers — every {{#if x}} needs a matching {{/if}}."
        )

    # Styles present vs required.
    try:
        present = {s.name for s in doc.styles if getattr(s, "name", None)}
    except Exception:
        present = set()
    analysis.styles_present = sorted(present)
    analysis.missing_required_styles = [s for s in REQUIRED_DOCX_STYLES if s not in present]

    return analysis


def analyze_html(data: bytes) -> Analysis:
    """Analyse an email ``.html`` template given its bytes."""
    analysis = Analysis(kind="email")
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("utf-8", errors="replace")
        analysis.warnings.append("File was not valid UTF-8; decoded with replacements.")

    sections = set(_MUSTACHE_SECTION.findall(text))
    analysis.conditionals = sorted(sections)

    variables = [
        v for v in _MUSTACHE_VAR.findall(text)
        if v not in sections and v not in EMAIL_RESERVED_VARS
    ]
    analysis.placeholders = _ordered_unique(variables)
    return analysis


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

# The theme lives on a relationship from the slide master. Spelt out rather
# than imported so this module keeps its light import surface.
_THEME_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"


@dataclass
class LayoutInfo:
    """One slide layout, as the admin UI needs to show it."""
    index: int
    name: str
    placeholders: List[str] = field(default_factory=list)
    role: Optional[str] = None          # auto-detected role, None when unrecognised
    has_footer: bool = False


@dataclass
class PptxAnalysis:
    """Result of analysing a ``.pptx`` / ``.potx`` template.

    Deliberately *not* an :class:`Analysis`. That type describes a parameterised
    document — placeholders, conditionals, args to declare. A presentation
    template has none of those: it is a design, and what an admin needs to see
    is which layouts it offers and which slide roles those cover. The only
    field the two share is ``warnings``, which is all the shared view code
    touches.
    """
    kind: str = "pptx"
    aspect: Optional[str] = None
    slide_size: Optional[str] = None
    layouts: List[LayoutInfo] = field(default_factory=list)
    role_map: Dict[str, str] = field(default_factory=dict)   # role -> layout name
    missing_roles: List[str] = field(default_factory=list)
    theme_fonts: Dict[str, str] = field(default_factory=dict)
    theme_colors: Dict[str, str] = field(default_factory=dict)
    embedded_slides: int = 0
    warnings: List[str] = field(default_factory=list)

    @property
    def layout_names(self) -> List[str]:
        """Layout names, for populating the role-override dropdowns."""
        return [layout.name for layout in self.layouts]


def _placeholder_label(placeholder) -> str:
    """Human-readable placeholder type, e.g. ``TITLE`` or ``BODY``."""
    try:
        kind = placeholder.placeholder_format.type
    except Exception:
        return "UNKNOWN"
    name = getattr(kind, "name", None)
    return str(name or kind)


def _theme_root(presentation):
    """The parsed theme XML of the first slide master, or None.

    python-pptx has no model class for a theme, so the relationship resolves to
    a generic ``Part`` carrying only raw bytes — there is no ``_element`` to
    read. Reaching for one returns None and yields an empty palette rather than
    an error, so the blob is parsed here instead.
    """
    try:
        master = presentation.slide_masters[0]
    except (IndexError, AttributeError):
        return None
    try:
        theme = master.part.part_related_by(_THEME_RELTYPE)
    except (KeyError, AttributeError):
        return None

    blob = getattr(theme, "blob", None)
    if not blob:
        return None
    try:
        return etree.fromstring(blob)
    except etree.XMLSyntaxError:
        return None


def analyze_pptx(data: bytes) -> PptxAnalysis:
    """Analyse a PowerPoint template given its bytes.

    Reports what actually decides how a generated deck looks: the slide size,
    the layouts on offer with their placeholder signatures, which slide role
    each layout is auto-detected as, and the theme's fonts and accent colours.

    A ``.potx`` is accepted — it is opened through the same
    :func:`pptx_tools.templates.open_template` the tool uses, which rewrites the
    template content type in memory.
    """
    from pptx.util import Emu

    from pptx_tools.layouts import ROLES, classify_layout
    from pptx_tools.templates import aspect_of, open_template

    analysis = PptxAnalysis()

    # open_template takes a path, so the upload is staged to a temp file. The
    # .potx branch keys off the content type rather than the name, so the
    # suffix used here does not decide anything.
    with tempfile.TemporaryDirectory() as tmp:
        staged = Path(tmp) / "upload.pptx"
        staged.write_bytes(data)
        try:
            presentation = open_template(staged)
        except Exception as e:
            analysis.warnings.append(f"Could not open as a PowerPoint template: {e}")
            return analysis

        analysis.aspect = aspect_of(presentation)
        width = Emu(presentation.slide_width).inches
        height = Emu(presentation.slide_height).inches
        analysis.slide_size = f"{width:.2f} × {height:.2f} in"
        analysis.embedded_slides = len(presentation.slides._sldIdLst)

        seen_roles: Dict[str, str] = {}
        for index, layout in enumerate(presentation.slide_layouts):
            placeholders = [_placeholder_label(p) for p in layout.placeholders]
            role = classify_layout(layout)
            analysis.layouts.append(LayoutInfo(
                index=index,
                name=layout.name,
                placeholders=placeholders,
                role=role,
                has_footer=any(label == "FOOTER" for label in placeholders),
            ))
            # First layout detected for a role is the one the resolver picks.
            if role and role not in seen_roles:
                seen_roles[role] = layout.name

        analysis.role_map = seen_roles
        analysis.missing_roles = [role for role in ROLES if role not in seen_roles]
        analysis.theme_fonts, analysis.theme_colors = _read_theme(presentation)

    _add_pptx_warnings(analysis)
    return analysis


def _read_theme(presentation) -> tuple:
    """Return ``(fonts, colors)`` from the first slide master's theme."""
    # python-pptx's namespace map, not python-docx's. Both happen to resolve
    # the DrawingML "a:" prefix identically, but reading a PowerPoint theme
    # through the Word nsmap is an accident waiting to be broken by either side.
    from pptx.oxml.ns import qn as pptx_qn

    fonts: Dict[str, str] = {}
    colors: Dict[str, str] = {}
    theme = _theme_root(presentation)
    if theme is None:
        return fonts, colors

    scheme = theme.find(pptx_qn("a:themeElements"))
    if scheme is None:
        return fonts, colors

    font_scheme = scheme.find(pptx_qn("a:fontScheme"))
    if font_scheme is not None:
        for tag, label in (("a:majorFont", "headings"), ("a:minorFont", "body")):
            node = font_scheme.find(pptx_qn(tag))
            latin = node.find(pptx_qn("a:latin")) if node is not None else None
            typeface = latin.get("typeface") if latin is not None else None
            if typeface:
                fonts[label] = typeface

    color_scheme = scheme.find(pptx_qn("a:clrScheme"))
    if color_scheme is not None:
        for name in ("dk1", "lt1", "dk2", "lt2",
                     "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
                     "hlink", "folHlink"):
            node = color_scheme.find(pptx_qn(f"a:{name}"))
            if node is None:
                continue
            srgb = node.find(pptx_qn("a:srgbClr"))
            if srgb is not None and srgb.get("val"):
                colors[name] = f"#{srgb.get('val').upper()}"
                continue
            # A system colour (windowText / window) carries its resolved value
            # in lastClr; without this dk1/lt1 read as blank on Office themes.
            sys_clr = node.find(pptx_qn("a:sysClr"))
            if sys_clr is not None and sys_clr.get("lastClr"):
                colors[name] = f"#{sys_clr.get('lastClr').upper()}"
                continue
            # OOXML also permits scrgbClr, hslClr, prstClr and schemeClr here.
            # Office writes only srgbClr and sysClr, so the others are not worth
            # converting — but say so rather than dropping the entry in silence,
            # which is the failure mode this whole reader was written to avoid.
            if len(node):
                logger.info(
                    "[template-analysis] Theme colour %r uses <%s>, which is not "
                    "one of srgbClr/sysClr; omitted from the reported palette.",
                    name, etree.QName(node[0]).localname,
                )
    return fonts, colors


def _add_pptx_warnings(analysis: PptxAnalysis) -> None:
    """Attach the warnings an admin can act on."""
    if not analysis.layouts:
        analysis.warnings.append("This template defines no slide layouts.")
        return

    # Roles that fall back by position when absent, which is what produces a
    # deck laid out on the wrong layouts.
    if analysis.missing_roles:
        analysis.warnings.append(
            "No layout was detected for: " + ", ".join(analysis.missing_roles) +
            ". Slides needing those roles fall back by position — set an override "
            "below, or ignore this if the template is not meant to cover them."
        )

    without_footer = [layout.name for layout in analysis.layouts if not layout.has_footer]
    if without_footer and len(without_footer) < len(analysis.layouts):
        analysis.warnings.append(
            f"No footer placeholder on {len(without_footer)} of {len(analysis.layouts)} "
            "layouts; footer text and slide numbers will not show on those slides."
        )
    elif len(without_footer) == len(analysis.layouts):
        analysis.warnings.append(
            "No layout has a footer placeholder, so footer text and slide numbers "
            "will not appear anywhere in a generated deck."
        )

    if analysis.embedded_slides:
        analysis.warnings.append(
            f"The file contains {analysis.embedded_slides} slide(s) of its own. "
            "They are stripped from generated decks by default (strip_slides)."
        )

    unrecognised = [layout.name for layout in analysis.layouts if layout.role is None]
    if unrecognised:
        analysis.warnings.append(
            f"{len(unrecognised)} layout(s) match no known role and can only be used "
            "by naming them explicitly: " + ", ".join(unrecognised[:6])
            + ("…" if len(unrecognised) > 6 else "")
        )


def analyze(kind: str, data: bytes):
    """Dispatch analysis by template *kind* (``docx``, ``email`` or ``pptx``)."""
    if kind == "docx":
        return analyze_docx(data)
    if kind == "email":
        return analyze_html(data)
    if kind == "pptx":
        return analyze_pptx(data)
    raise ValueError(f"Unknown template kind: {kind!r}")


@dataclass
class Reconciliation:
    """Comparison of detected placeholders/conditionals against declared args."""
    missing_args: List[str] = field(default_factory=list)      # placeholder, no arg
    orphan_args: List[str] = field(default_factory=list)       # arg, no placeholder
    non_bool_conditions: List[str] = field(default_factory=list)  # condition arg not bool


def reconcile(analysis: Analysis, args: List[Dict[str, Any]]) -> Reconciliation:
    """Compare an :class:`Analysis` with a template's declared ``args``."""
    arg_by_name = {a.get("name"): a for a in (args or []) if isinstance(a, dict) and a.get("name")}
    arg_names = set(arg_by_name)

    detected = set(analysis.placeholders) | set(analysis.conditionals)
    rec = Reconciliation()
    rec.missing_args = [p for p in analysis.placeholders if p not in arg_names]
    # Conditionals also need a (bool) arg.
    rec.missing_args += [c for c in analysis.conditionals if c not in arg_names and c not in rec.missing_args]
    rec.orphan_args = [n for n in arg_by_name if n not in detected and n not in EMAIL_RESERVED_VARS]
    rec.non_bool_conditions = [
        c for c in analysis.conditionals
        if c in arg_by_name and str(arg_by_name[c].get("type", "string")).lower() not in ("bool", "boolean")
    ]
    return rec
