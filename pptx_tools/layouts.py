"""Resolve slide layouts by name and placeholder signature, not by index.

The builder used to address layouts positionally — layout 0 is the title
slide, 1 is title-and-content, and so on — which only holds for a template
whose layouts sit in PowerPoint's default order. A corporate template that
reorders, renames or deletes layouts is routine, and the failure was silent:
a template with the section layout moved first built the deck's title slide on
it, and a template trimmed to three layouts raised "slide layout index out of
range".

Resolution order, most specific first:

1. the ``layout`` named on the slide itself,
2. the ``layouts:`` mapping in the template's registry entry,
3. the placeholder signature of each layout (below),
4. the positional index, with a warning.

Signature classification reads the *types* of a layout's placeholders, so it
is independent of both language and order — the shipped Czech-named templates
classify from the same rules as an English one.
"""

from __future__ import annotations

import logging
from typing import Dict, List, Optional, Tuple

from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn

from .constants import (
    TITLE_LAYOUT, SECTION_LAYOUT, CONTENT_LAYOUT,
    TWO_COLUMN_LAYOUT, TWO_COLUMN_TEXT_LAYOUT, TITLE_ONLY_LAYOUT, BLANK_LAYOUT,
)

logger = logging.getLogger(__name__)


# Roles a layout can play. These are what slide types ask for; they are not
# PowerPoint concepts, and a template need not provide all of them.
ROLE_TITLE = "title"
ROLE_SECTION = "section"
ROLE_CONTENT = "content"
ROLE_TWO_COLUMN = "two_column"
ROLE_COMPARISON = "comparison"
ROLE_IMAGE_TEXT = "image_text"
ROLE_TITLE_ONLY = "title_only"
ROLE_BLANK = "blank"

ROLES = (
    ROLE_TITLE, ROLE_SECTION, ROLE_CONTENT, ROLE_TWO_COLUMN,
    ROLE_COMPARISON, ROLE_IMAGE_TEXT, ROLE_TITLE_ONLY, ROLE_BLANK,
)

# Positional fallbacks, i.e. what the builder assumed before this module.
ROLE_FALLBACK_INDEX = {
    ROLE_TITLE: TITLE_LAYOUT,
    ROLE_SECTION: SECTION_LAYOUT,
    ROLE_CONTENT: CONTENT_LAYOUT,
    ROLE_TWO_COLUMN: TWO_COLUMN_LAYOUT,
    ROLE_COMPARISON: TWO_COLUMN_TEXT_LAYOUT,
    ROLE_IMAGE_TEXT: CONTENT_LAYOUT,
    ROLE_TITLE_ONLY: TITLE_ONLY_LAYOUT,
    ROLE_BLANK: BLANK_LAYOUT,
}

# Which role each slide type wants. two_column picks between two roles
# depending on whether its columns carry headings, so it is resolved in
# :func:`role_for_slide` rather than listed here.
SLIDE_TYPE_ROLE = {
    "title": ROLE_TITLE,
    # Slide types that draw their own shapes want a clear canvas with a title.
    "kpi": ROLE_TITLE_ONLY,
    "timeline": ROLE_TITLE_ONLY,
    "agenda": ROLE_CONTENT,
    "closing": ROLE_TITLE,
    "section": ROLE_SECTION,
    "content": ROLE_CONTENT,
    "table": ROLE_CONTENT,
    "chart": ROLE_CONTENT,
    "scatter": ROLE_CONTENT,
    "quote": ROLE_CONTENT,
    "image": ROLE_CONTENT,
}

# Placeholders every layout carries regardless of its purpose; they say nothing
# about the role, so the signature ignores them.
_CHROME = {PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER}

_TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
# BODY and OBJECT are both "put content here"; templates use them
# interchangeably, so the signature counts them together and only separates
# them where it actually distinguishes a layout (comparison vs two content).
_CONTENT_TYPES = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}


def _is_vertical(placeholder) -> bool:
    """True for a vertical-text placeholder.

    PowerPoint's "Title and Vertical Text" layout has the same TITLE + BODY
    signature as "Section Header", so without this the two are
    indistinguishable and a template could get its section slides laid out
    sideways.
    """
    try:
        bodyPr = placeholder.text_frame._txBody.find(qn('a:bodyPr'))
    except AttributeError:
        return False
    if bodyPr is None:
        return False
    return (bodyPr.get('vert') or 'horz') != 'horz'


def layout_signature(layout) -> Tuple[List, bool]:
    """Return the meaningful placeholder types of *layout*, and its verticality."""
    types = []
    vertical = False
    for placeholder in layout.placeholders:
        kind = placeholder.placeholder_format.type
        if kind in _CHROME:
            continue
        types.append(kind)
        if kind not in _TITLE_TYPES and _is_vertical(placeholder):
            vertical = True
    return types, vertical


def classify_layout(layout) -> Optional[str]:
    """Return the role *layout* plays, or None when it matches no known shape.

    Deliberately conservative: an unrecognised layout returns None and is left
    for an explicit name to select, rather than being forced into the nearest
    role and producing a subtly wrong deck.
    """
    types, vertical = layout_signature(layout)

    if not types:
        return ROLE_BLANK

    titles = [t for t in types if t in _TITLE_TYPES]
    contents = [t for t in types if t in _CONTENT_TYPES]
    pictures = [t for t in types if t == PP_PLACEHOLDER.PICTURE]

    # A picture placeholder is the distinguishing feature; nothing else has one.
    # It still has to carry a title, like every other titled role: a full-bleed
    # photo layout with no title placeholder would otherwise be picked
    # confidently and then silently drop the slide's title.
    if pictures and titles:
        return ROLE_IMAGE_TEXT

    if vertical:
        # Vertical-text layouts share signatures with horizontal ones. They are
        # a deliberate typographic choice, never a sensible automatic pick.
        return None

    has_center_title = any(t == PP_PLACEHOLDER.CENTER_TITLE for t in types)
    if has_center_title and any(t == PP_PLACEHOLDER.SUBTITLE for t in types):
        return ROLE_TITLE
    if titles and any(t == PP_PLACEHOLDER.SUBTITLE for t in types):
        return ROLE_TITLE

    if not titles:
        return None

    if not contents:
        return ROLE_TITLE_ONLY

    # Comparison is title + heading/content twice over; Two Content is title
    # plus two content placeholders and no headings.
    if len(contents) == 4:
        return ROLE_COMPARISON
    if len(contents) == 2:
        body_count = sum(1 for t in contents if t == PP_PLACEHOLDER.BODY)
        object_count = sum(1 for t in contents if t == PP_PLACEHOLDER.OBJECT)
        if object_count == 2:
            return ROLE_TWO_COLUMN
        if body_count == 1 and object_count == 1:
            # "Content with Caption": one content area plus a caption block.
            # Not a two-column slide, and not a plain content slide either.
            return None
        return ROLE_TWO_COLUMN
    if len(contents) == 1:
        # Section Header and Title and Content differ only by BODY vs OBJECT.
        return ROLE_SECTION if contents[0] == PP_PLACEHOLDER.BODY else ROLE_CONTENT

    return None


def role_for_slide(slide) -> str:
    """The layout role a validated slide model wants."""
    if slide.type == "two_column":
        has_headings = bool(slide.left.heading or slide.right.heading)
        return ROLE_COMPARISON if has_headings else ROLE_TWO_COLUMN
    if slide.type == "quote" and not slide.title:
        # An untitled quote on a titled layout leaves an empty title
        # placeholder sitting above it, which shows as "Click to add title"
        # the moment anyone opens the deck to edit it.
        return ROLE_BLANK
    return SLIDE_TYPE_ROLE.get(slide.type, ROLE_CONTENT)


class LayoutResolver:
    """Chooses a slide layout for each role in one presentation."""

    def __init__(self, presentation, configured: Optional[Dict[str, str]] = None):
        self._presentation = presentation
        self._layouts = list(presentation.slide_layouts)
        self._configured = {k: v for k, v in (configured or {}).items() if v}

        self._by_name: Dict[str, object] = {}
        for layout in self._layouts:
            # First definition wins, matching how PowerPoint resolves a
            # duplicate layout name.
            self._by_name.setdefault(layout.name, layout)
            self._by_name.setdefault(layout.name.strip().lower(), layout)

        self._by_role: Dict[str, object] = {}
        self._roles_of: Dict[str, str] = {}
        for layout in self._layouts:
            role = classify_layout(layout)
            self._roles_of[layout.name] = role
            if role and role not in self._by_role:
                self._by_role[role] = layout

    # -- lookup ------------------------------------------------------------

    def by_name(self, name: str):
        """Return the layout called *name*, case-insensitively, or None."""
        if not name:
            return None
        return self._by_name.get(name) or self._by_name.get(name.strip().lower())

    def resolve(self, role: str, override: Optional[str] = None):
        """Return ``(layout, warning)`` for *role*.

        *warning* is None when the layout was chosen confidently; otherwise it
        explains what the caller got instead of what it asked for.
        """
        if override:
            layout = self.by_name(override)
            if layout is not None:
                return layout, None
            # Fall through to the normal resolution so the slide is still
            # built, but say plainly that the name did not match.
            layout, _ = self._resolve_without_override(role)
            return layout, (
                f"layout {override!r} is not in this template "
                f"(available: {', '.join(sorted(self.layout_names))}); "
                "used the default layout for this slide type."
            )

        return self._resolve_without_override(role)

    def _resolve_without_override(self, role: str):
        configured_name = self._configured.get(role)
        if configured_name:
            layout = self.by_name(configured_name)
            if layout is not None:
                return layout, None
            return self._resolve_detected(role, extra=(
                f"template config maps {role!r} to layout {configured_name!r}, "
                "which this file does not contain"
            ))

        return self._resolve_detected(role)

    def _resolve_detected(self, role: str, extra: Optional[str] = None):
        layout = self._by_role.get(role)
        if layout is not None:
            return layout, (f"{extra}; fell back to detected layout {layout.name!r}." if extra else None)

        index = ROLE_FALLBACK_INDEX.get(role, CONTENT_LAYOUT)
        if index < len(self._layouts):
            layout = self._layouts[index]
            note = (
                f"no layout in this template matches the {role!r} role; "
                f"used layout {index} ({layout.name!r}) by position."
            )
            return layout, (f"{extra}; {note}" if extra else note)

        # Nothing at that index either: use the last layout rather than raising
        # IndexError, which is what a trimmed template used to do.
        layout = self._layouts[-1]
        note = (
            f"this template has only {len(self._layouts)} layout(s) and none "
            f"matches the {role!r} role; used {layout.name!r}."
        )
        return layout, (f"{extra}; {note}" if extra else note)

    # -- reporting ---------------------------------------------------------

    @property
    def layout_names(self) -> List[str]:
        return [layout.name for layout in self._layouts]

    def coverage(self) -> Dict[str, Optional[str]]:
        """Map each role to the layout name serving it, or None if detected none."""
        return {
            role: (self._by_role[role].name if role in self._by_role else None)
            for role in ROLES
        }

    def describe(self) -> List[Dict[str, object]]:
        """Per-layout detail for the template-listing tool."""
        described = []
        for index, layout in enumerate(self._layouts):
            types, vertical = layout_signature(layout)
            described.append({
                "index": index,
                "name": layout.name,
                "role": self._roles_of.get(layout.name),
                "vertical": vertical,
                "placeholders": [
                    {
                        "idx": ph.placeholder_format.idx,
                        "type": str(ph.placeholder_format.type).split()[0],
                    }
                    for ph in layout.placeholders
                ],
            })
        return described

    def missing_roles(self) -> List[str]:
        """Roles the slide types actually use that this template does not provide."""
        used = set(SLIDE_TYPE_ROLE.values()) | {ROLE_TWO_COLUMN, ROLE_COMPARISON}
        return sorted(role for role in used if role not in self._by_role)

    def layouts_without_footer(self) -> List[str]:
        """Layouts with no footer placeholder, where footer_text cannot appear."""
        missing = []
        for layout in self._layouts:
            kinds = {ph.placeholder_format.type for ph in layout.placeholders}
            if PP_PLACEHOLDER.FOOTER not in kinds:
                missing.append(layout.name)
        return missing
