"""Utility functions for creating charts in PowerPoint presentations.

This module provides functionality to create various chart types in PowerPoint slides
using python-pptx's chart capabilities.
"""

import logging
from typing import Dict, Any, Optional

from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Pt

logger = logging.getLogger(__name__)


# Mapping of chart type strings to python-pptx chart types.
#
# Every entry here is category-based and is built from CategoryChartData.
# 'scatter' (XY_SCATTER) is deliberately absent: an XY chart needs XyChartData
# and (x, y) point pairs, so building one from categories+series raised
# "'CategoryWorkbookWriter' object has no attribute 'x_values_ref'" on every
# call. It was advertised in the tool description but could never succeed.
# A dedicated scatter slide type with its own data shape is tracked separately.
CHART_TYPE_MAP = {
    'bar': XL_CHART_TYPE.BAR_CLUSTERED,
    'bar_stacked': XL_CHART_TYPE.BAR_STACKED,
    'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
    'column_stacked': XL_CHART_TYPE.COLUMN_STACKED,
    'line': XL_CHART_TYPE.LINE,
    'line_markers': XL_CHART_TYPE.LINE_MARKERS,
    'pie': XL_CHART_TYPE.PIE,
    'doughnut': XL_CHART_TYPE.DOUGHNUT,
    'area': XL_CHART_TYPE.AREA,
    'area_stacked': XL_CHART_TYPE.AREA_STACKED,
    'radar': XL_CHART_TYPE.RADAR,
}

# Scatter is an XY chart: it needs XyChartData and (x, y) pairs rather than
# categories, so it is built by add_scatter_to_slide() and deliberately absent
# from CHART_TYPE_MAP. Asking for it as a category chart_type is still a
# mistake worth naming precisely.
UNSUPPORTED_CHART_TYPES = {
    'scatter': (
        "Scatter is its own slide type, not a category chart_type: use "
        "{\"type\": \"scatter\", \"series\": [{\"name\": ..., \"points\": [[x, y], ...]}]}."
    ),
}

LEGEND_POSITIONS = {
    'left': XL_LEGEND_POSITION.LEFT,
    'right': XL_LEGEND_POSITION.RIGHT,
    'top': XL_LEGEND_POSITION.TOP,
    'bottom': XL_LEGEND_POSITION.BOTTOM,
}


class ChartDataError(Exception):
    """Exception raised when chart data is invalid."""
    pass


def validate_chart_data(chart_data: Dict[str, Any], chart_type: str) -> None:
    """Validate chart data structure.

    Args:
        chart_data: Dictionary containing categories and series.
        chart_type: Type of chart being created.

    Raises:
        ChartDataError: If data is invalid.
    """
    if not chart_data:
        raise ChartDataError("Chart data is required")

    if chart_type in UNSUPPORTED_CHART_TYPES:
        raise ChartDataError(
            f"{UNSUPPORTED_CHART_TYPES[chart_type]} Available: {', '.join(CHART_TYPE_MAP)}"
        )

    if chart_type not in CHART_TYPE_MAP:
        raise ChartDataError(f"Unknown chart type: {chart_type}. Available: {', '.join(CHART_TYPE_MAP.keys())}")

    if 'categories' not in chart_data:
        raise ChartDataError("Chart data must include 'categories'")
    if not chart_data['categories']:
        raise ChartDataError("Categories list cannot be empty")

    if 'series' not in chart_data:
        raise ChartDataError("Chart data must include 'series'")
    if not chart_data['series']:
        raise ChartDataError("Series list cannot be empty")

    # Validate each series
    for i, series in enumerate(chart_data['series']):
        if not isinstance(series, dict):
            raise ChartDataError(f"Series {i} must be a dictionary")
        if 'name' not in series:
            raise ChartDataError(f"Series {i} must have a 'name'")
        if 'values' not in series:
            raise ChartDataError(f"Series {i} must have 'values'")
        if not series['values']:
            raise ChartDataError(f"Series {i} values cannot be empty")


def create_chart_data(chart_data: Dict[str, Any]) -> CategoryChartData:
    """Create CategoryChartData from structured data.

    Args:
        chart_data: Dictionary with 'categories' and 'series' keys.

    Returns:
        CategoryChartData object ready for chart creation.
    """
    data = CategoryChartData()
    data.categories = chart_data['categories']

    for series in chart_data['series']:
        data.add_series(series['name'], series['values'])

    return data


def add_chart_to_slide(
    slide,
    chart_type: str,
    chart_data: Dict[str, Any],
    left: int,
    top: int,
    width: int,
    height: int,
    has_legend: bool = True,
    legend_position: str = 'right',
    title: Optional[str] = None
) -> None:
    """Add a chart to a slide.

    Args:
        slide: PowerPoint slide object.
        chart_type: Type of chart (bar, line, pie, etc.).
        chart_data: Dictionary with categories and series data.
        left: Left position in EMUs.
        top: Top position in EMUs.
        width: Chart width in EMUs.
        height: Chart height in EMUs.
        has_legend: Whether to show legend.
        legend_position: Legend position (left, right, top, bottom).
        title: Optional chart title.
    """
    logger.debug(f"Adding {chart_type} chart to slide")

    # Validate data
    validate_chart_data(chart_data, chart_type)

    # Get chart type enum
    xl_chart_type = CHART_TYPE_MAP[chart_type]

    # Create chart data
    data = create_chart_data(chart_data)

    # Add chart to slide
    chart_shape = slide.shapes.add_chart(
        xl_chart_type,
        left, top, width, height,
        data
    )

    chart = chart_shape.chart

    _configure_legend(chart, has_legend, legend_position)
    _configure_title(chart, title)

    logger.debug(f"Chart added successfully with {len(chart_data['series'])} series")
    return chart


# ---------------------------------------------------------------------------
# Shared chart configuration
# ---------------------------------------------------------------------------

def _configure_legend(chart, has_legend: bool, legend_position: str) -> None:
    """Show or hide the legend and place it."""
    if not has_legend or legend_position == 'none':
        chart.has_legend = False
        return
    chart.has_legend = True
    chart.legend.position = LEGEND_POSITIONS.get(legend_position, XL_LEGEND_POSITION.RIGHT)
    chart.legend.include_in_layout = False


def _configure_title(chart, title: Optional[str]) -> None:
    """Set or clear the chart's own title."""
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
    else:
        chart.has_title = False


def configure_data_labels(chart, enabled: bool, number_format: Optional[str] = None) -> None:
    """Turn value labels on each point on or off.

    A number format without labels is still applied to the value axis, so
    "#,##0" or "0.0%" formats the tick labels even when the caller did not ask
    for per-point labels.
    """
    plot = chart.plots[0]
    plot.has_data_labels = bool(enabled)
    if enabled:
        labels = plot.data_labels
        labels.font.size = Pt(10)
        if number_format:
            labels.number_format = number_format
            labels.number_format_is_linked = False
        try:
            labels.position = XL_LABEL_POSITION.OUTSIDE_END
        except (ValueError, NotImplementedError):
            # Not valid for every chart type (pie/doughnut/stacked); the
            # default position is fine there.
            pass
    elif number_format:
        try:
            chart.value_axis.tick_labels.number_format = number_format
            chart.value_axis.tick_labels.number_format_is_linked = False
        except (ValueError, NotImplementedError):
            pass


def set_axis_titles(chart, x_title: Optional[str] = None, y_title: Optional[str] = None) -> None:
    """Label the category and value axes, where the chart type has them.

    Pie and doughnut charts have no axes; asking for a title there is a no-op
    rather than an error, because it is a reasonable thing for a caller to send
    without knowing the chart type's geometry.
    """
    for axis_name, title in (("category_axis", x_title), ("value_axis", y_title)):
        if not title:
            continue
        try:
            axis = getattr(chart, axis_name)
            axis.has_title = True
            axis.axis_title.text_frame.paragraphs[0].text = title
        except (ValueError, NotImplementedError, AttributeError):
            logger.debug("Chart type has no %s; skipping its title", axis_name)


def _series_field(entry, field: str):
    """Read *field* from a model attribute or a mapping key.

    Written out rather than expressed as ``getattr(...) or entry.get(...)``:
    that idiom cannot tell "absent" from "present but falsy", so a legal
    empty-string series name resolved to None. It is a narrow case — the
    rendered legend entry came out the same either way — but it is the kind of
    thing duck typing hides until it matters.
    """
    if isinstance(entry, dict):
        return entry.get(field)
    return getattr(entry, field, None)


def add_scatter_to_slide(
    slide,
    series: list,
    left: int,
    top: int,
    width: int,
    height: int,
    legend: str = 'right',
    title: Optional[str] = None,
    x_title: Optional[str] = None,
    y_title: Optional[str] = None,
):
    """Add an XY (scatter) chart built from [x, y] point pairs.

    Separate from :func:`add_chart_to_slide` because an XY chart needs
    ``XyChartData``: feeding categories and series to it raises
    ``'CategoryWorkbookWriter' object has no attribute 'x_values_ref'``, which
    is why the previously advertised 'scatter' chart_type could never build.

    Args:
        series: Objects with ``name`` and ``points`` ([[x, y], ...]).
    """
    if not series:
        raise ChartDataError("Scatter chart needs at least one series")

    data = XyChartData()
    for i, entry in enumerate(series):
        name = _series_field(entry, "name")
        points = _series_field(entry, "points")

        if name is None:
            raise ChartDataError(f"Scatter series {i} has no name")
        if not points:
            raise ChartDataError(f"Scatter series {name!r} has no points")

        series_data = data.add_series(name)
        for x, y in points:
            series_data.add_data_point(float(x), float(y))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, left, top, width, height, data
    )
    chart = chart_shape.chart

    _configure_legend(chart, legend != 'none', legend)
    _configure_title(chart, title)
    set_axis_titles(chart, x_title, y_title)

    logger.debug("Scatter chart added with %d series", len(series))
    return chart


