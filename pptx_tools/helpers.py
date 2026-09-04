"""PowerPoint helper utilities and slide-building mixin.

This module provides a single SlideHelpers mixin class that consolidates
all common slide operations (text, tables, images) and standalone utility
functions for template loading and data parsing.
"""

import logging
import math
import re
from typing import List, Tuple, Optional, Any

from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn

from .constants import (
    CONTENT_LAYOUT,
    DEFAULT_BODY_FONT_SIZE,
    MARGIN_LEFT, MAX_INDENT_LEVEL,
    AVG_CHAR_WIDTH_RATIO, LINE_HEIGHT_RATIO, MIN_AUTOFIT_SCALE,
    TABLE_MIN_FONT_SIZE, TABLE_ROW_HEIGHT_PER_POINT,
    TABLE_HEADER_FILL, TABLE_HEADER_TEXT, TABLE_ALT_ROW_FILL,
)
from .schema import Bullet, THEME_COLORS
from .image_utils import load_image, ImageDownloadError, ImageValidationError
from .inline_formatting import needs_inline_processing, apply_inline_formatting

logger = logging.getLogger(__name__)


# =============================================================================
# Utility Functions
# =============================================================================

# Regex matching a single markdown table separator cell: optional colon, 3+ dashes, optional colon
_SEPARATOR_CELL_RE = re.compile(r'^\s*:?-{3,}:?\s*$')


def cell_to_text(value: Any) -> str:
    """Render one raw table cell as text.

    Table data arrives straight from the model, which sends numbers as numbers
    and omissions as ``null``. Everything downstream (the separator regex, the
    cell writer) needs a string, so normalise here rather than at each use.

    ``None`` becomes an empty cell; every other value is stringified, so a
    numeric ``0`` survives as "0" instead of being dropped by a falsy test.
    """
    if value is None:
        return ""
    if isinstance(value, str):
        return value
    return str(value)


def _is_separator_row(row: List[str]) -> bool:
    """Check if a row is a markdown table separator row (e.g., |:---|:---:|---:|).

    Uses strict per-cell regex to avoid false positives on content containing dashes.
    Expects cells already normalised to text by :func:`cell_to_text`.
    """
    return bool(row) and all(_SEPARATOR_CELL_RE.match(cell) for cell in row)


def _extract_alignments(row: List[str]) -> List[Optional[int]]:
    """Extract column alignments from a markdown separator row.

    Args:
        row: List of separator cells (e.g., [':---', ':---:', '---:']).

    Returns:
        List of PP_ALIGN values (LEFT, CENTER, RIGHT) or None per column.
    """
    alignments = []
    for cell in row:
        cell = cell.strip()
        if cell.startswith(':') and cell.endswith(':'):
            alignments.append(PP_ALIGN.CENTER)
        elif cell.endswith(':'):
            alignments.append(PP_ALIGN.RIGHT)
        else:
            alignments.append(None)  # left/default
    return alignments


def parse_table_data(table_data: List[List[Any]]) -> tuple:
    """Clean table data by removing markdown separator rows and extracting alignments.

    Cells are normalised to text first, so numeric and ``null`` cells are
    accepted; previously they raised ``TypeError`` inside the separator regex
    and failed the whole presentation.

    Args:
        table_data: Raw table data as list of rows. Cells may be any scalar.

    Returns:
        Tuple of (cleaned_rows, column_alignments).
        column_alignments is a list of PP_ALIGN values or None per column,
        or None if no separator row was found.
    """
    if not table_data:
        return [], None

    cleaned = []
    col_alignments = None

    for row in table_data:
        # A row sent as a bare scalar becomes a single-cell row rather than an error.
        cells = [cell_to_text(c) for c in row] if isinstance(row, (list, tuple)) else [cell_to_text(row)]

        if _is_separator_row(cells):
            col_alignments = _extract_alignments(cells)
        else:
            cleaned.append(cells)

    return cleaned, col_alignments


def parse_color(color_hex: Any, default: RGBColor) -> RGBColor:
    """Parse a hex color string to RGBColor.

    Accepts both "4172C4" and "#4172C4" — models routinely send the leading
    hash, which ``RGBColor.from_string`` rejects. An unusable value falls back
    to *default* and is logged at WARNING, because silently substituting a
    colour is the kind of failure nobody notices until the deck is on screen.

    Args:
        color_hex: Hex color string, with or without a leading '#'.
        default: Default color if parsing fails.

    Returns:
        RGBColor object.
    """
    if color_hex is None:
        return default

    try:
        return RGBColor.from_string(str(color_hex).strip().lstrip('#').upper())
    except (ValueError, AttributeError):
        logger.warning(
            "Unrecognised colour %r (expected 6-digit hex such as '4172C4' or '#4172C4'); using %s",
            color_hex, default,
        )
        return default


_BULLET_LINE_RE = re.compile(r'^(\s*)[-*+]\s+(.*)$')


def body_to_bullets(body: Any) -> List[Bullet]:
    """Normalise a slide ``body`` into a list of :class:`Bullet`.

    Accepts what the schema allows: a markdown bullet string, or explicit
    bullet objects. In the string form, nesting comes from indentation and the
    unit is inferred rather than fixed — the distinct indent widths present are
    sorted and mapped onto levels 1, 2, 3…, so two-space, four-space and tab
    indentation all work without the caller declaring which they used. A line
    with no bullet marker is treated as a top-level bullet.
    """
    if body is None:
        return []

    if isinstance(body, list):
        out: List[Bullet] = []
        for item in body:
            if isinstance(item, Bullet):
                out.append(item)
            elif isinstance(item, dict):
                out.append(Bullet(**item))
            else:
                out.append(Bullet(text=cell_to_text(item)))
        return out

    if not isinstance(body, str):
        return [Bullet(text=cell_to_text(body))]

    parsed: List[Tuple[int, str]] = []
    for line in body.splitlines():
        if not line.strip():
            continue
        match = _BULLET_LINE_RE.match(line)
        if match:
            indent = len(match.group(1).replace('\t', '    '))
            text = match.group(2).strip()
        else:
            indent = 0
            text = line.strip()
        if text:
            parsed.append((indent, text))

    if not parsed:
        return []

    # Map the distinct indent widths onto consecutive levels.
    widths = sorted({indent for indent, _ in parsed})
    level_of = {width: min(rank + 1, MAX_INDENT_LEVEL) for rank, width in enumerate(widths)}

    return [Bullet(text=text, level=level_of[indent]) for indent, text in parsed]


def estimate_text_fill(
    bullets: List[Bullet],
    width: int,
    height: int,
    font_size_pt: float = 18.0,
) -> float:
    """Estimate how full a text box will be, as a ratio (1.0 = exactly full).

    Deliberately approximate — see the constants module. Used to decide whether
    to ask PowerPoint to shrink the text and whether to warn the caller.
    """
    width_in = Emu(width).inches
    height_in = Emu(height).inches
    if width_in <= 0 or height_in <= 0 or not bullets:
        return 0.0

    char_width_in = AVG_CHAR_WIDTH_RATIO * font_size_pt / 72.0
    line_height_in = LINE_HEIGHT_RATIO * font_size_pt / 72.0

    total_lines = 0
    for bullet in bullets:
        # Deeper levels are indented, so less width is available for text.
        indent_in = 0.3 * (bullet.level - 1)
        usable_in = max(width_in - indent_in, char_width_in)
        chars_per_line = max(int(usable_in / char_width_in), 1)
        total_lines += max(1, math.ceil(len(bullet.text) / chars_per_line))

    return (total_lines * line_height_in) / height_in


def apply_autofit(text_frame, scale: Optional[float] = None) -> None:
    """Mark a text frame as shrink-to-fit, optionally pinning the shrink factor.

    ``<a:normAutofit/>`` alone tells PowerPoint to shrink the text, but it
    recomputes the factor only when the deck is opened for editing; some
    viewers render the text at full size and overflowing until then. Writing an
    explicit ``fontScale`` gives those viewers something to honour immediately,
    and PowerPoint overwrites it with its own exact value on the first edit.
    """
    bodyPr = text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return

    for tag in ('a:noAutofit', 'a:spAutoFit', 'a:normAutofit'):
        existing = bodyPr.find(qn(tag))
        if existing is not None:
            bodyPr.remove(existing)

    autofit = parse_xml(
        '<a:normAutofit xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
    )
    if scale is not None and scale < 1.0:
        percent = max(scale, MIN_AUTOFIT_SCALE)
        autofit.set('fontScale', str(int(percent * 100000)))
        autofit.set('lnSpcReduction', '10000')
    bodyPr.append(autofit)


def set_runs_language(shape_or_frame, language: str) -> None:
    """Stamp ``lang`` on every run of a text frame.

    Without this the runs inherit the template's language, so a Czech deck is
    proof-read against the template's locale and every word is underlined as a
    misspelling.
    """
    text_frame = getattr(shape_or_frame, "text_frame", shape_or_frame)
    try:
        paragraphs = text_frame.paragraphs
    except AttributeError:
        return

    for paragraph in paragraphs:
        for run in paragraph.runs:
            run._r.get_or_add_rPr().set('lang', language)
        endParaRPr = paragraph._p.find(qn('a:endParaRPr'))
        if endParaRPr is not None:
            endParaRPr.set('lang', language)


def fit_table_font_size(row_count: int, height: int) -> int:
    """Pick a cell font size that keeps *row_count* rows inside *height*.

    PowerPoint grows a table to fit its text, so a tall table silently runs off
    the bottom of the slide rather than shrinking. Choosing the size up front
    from the row budget is what keeps a 20-row table on the slide.
    """
    height_pt = max(Emu(height).inches * 72.0, 1.0)
    per_row_pt = height_pt / max(row_count, 1)
    size = int(per_row_pt / TABLE_ROW_HEIGHT_PER_POINT)
    return max(TABLE_MIN_FONT_SIZE, min(size, int(DEFAULT_BODY_FONT_SIZE.pt)))


def table_overflows(row_count: int, height: int, font_size_pt: int) -> bool:
    """True when the table still cannot fit even at *font_size_pt*."""
    needed_pt = row_count * font_size_pt * TABLE_ROW_HEIGHT_PER_POINT
    return needed_pt > Emu(height).inches * 72.0


def resolve_fill(color: Optional[str], default: RGBColor):
    """Return either an ``RGBColor`` or a theme colour name for *color*.

    The schema accepts both ``#RRGGBB`` and a theme name such as ``accent1``;
    they need different DrawingML (``srgbClr`` vs ``schemeClr``), so the caller
    is told which it got.
    """
    if color is None:
        return default
    text = str(color).strip()
    if text.lower() in THEME_COLORS:
        return text.lower()
    return parse_color(text, default)


def coerce_indent_level(value: Any) -> int:
    """Map a 1-based ``indentation_level`` onto a 0-based pptx paragraph level.

    Tolerant by design: the value reaches us from a model, so a numeric string
    ("2") is accepted, a level past the supported depth is clamped instead of
    producing an invalid paragraph, and an unparseable value falls back to the
    top level with a warning rather than failing the presentation.
    """
    if value is None:
        return 0

    try:
        level = int(value)
    except (TypeError, ValueError):
        logger.warning("Invalid indentation_level %r; using 1", value)
        level = 1

    return max(1, min(level, MAX_INDENT_LEVEL)) - 1


# =============================================================================
# Consolidated Slide Helpers Mixin
# =============================================================================

class SlideHelpers:
    """Mixin providing all common slide helper methods (text, tables, images).

    Expects the consuming class to have a `self.presentation` attribute
    holding a python-pptx Presentation object.
    """

    # Type hint for IDE — actual attribute is set by the consuming class
    presentation: Any

    # -------------------------------------------------------------------------
    # Slide Management
    # -------------------------------------------------------------------------

    def _get_slide_dimensions(self) -> Tuple[int, int]:
        """Get slide width and height."""
        return self.presentation.slide_width, self.presentation.slide_height

    def _add_title_content_slide(self, title: str = ""):
        """Add a Title and Content slide and return slide with content placeholder info.

        Args:
            title: Title text for the slide.

        Returns:
            Tuple of (slide, content_left, content_top, content_width, content_height)
        """
        layout = self.presentation.slide_layouts[CONTENT_LAYOUT]
        slide = self.presentation.slides.add_slide(layout)

        # Set title
        if title and len(slide.placeholders) > 0:
            slide.placeholders[0].text = title

        # Get content placeholder bounds (idx 1)
        content_placeholder = None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 1:
                content_placeholder = placeholder
                break

        if content_placeholder:
            left = content_placeholder.left
            top = content_placeholder.top
            width = content_placeholder.width
            height = content_placeholder.height
            # Remove the placeholder so we can add custom content
            sp = content_placeholder._element
            sp.getparent().remove(sp)
        else:
            # Fallback dimensions
            slide_width, slide_height = self._get_slide_dimensions()
            left = MARGIN_LEFT
            top = Inches(1.5)
            width = slide_width - (2 * MARGIN_LEFT)
            height = slide_height - top - Inches(0.5)

        return slide, left, top, width, height

    def _add_speaker_notes(self, slide, notes_text: Optional[str]) -> None:
        """Add speaker notes to a slide.

        Args:
            slide: PowerPoint slide object.
            notes_text: Text for speaker notes.
        """
        if not notes_text:
            return
        try:
            slide.notes_slide.notes_text_frame.text = notes_text
            logger.debug(f"Added speaker notes: {notes_text[:50]}...")
        except Exception as e:
            logger.warning(f"Could not add speaker notes: {e}")

    # -------------------------------------------------------------------------
    # Text Helpers
    # -------------------------------------------------------------------------

    def _add_text_box(
        self,
        slide,
        text: str,
        left: int,
        top: int,
        width: int,
        height: int,
        font_size: Optional[int] = None,
        bold: bool = False,
        italic: bool = False,
        alignment=PP_ALIGN.LEFT,
        word_wrap: bool = True
    ):
        """Add a simple text box to a slide.

        Args:
            slide: PowerPoint slide object.
            text: Text content.
            left, top, width, height: Position and size.
            font_size: Font size.
            bold: Whether to make text bold.
            italic: Whether to make text italic.
            alignment: Text alignment.
            word_wrap: Whether to wrap text.

        Returns:
            Created textbox shape.
        """
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = word_wrap

        para = tf.paragraphs[0]
        para.text = text
        para.font.size = font_size or DEFAULT_BODY_FONT_SIZE
        para.font.bold = bold
        para.font.italic = italic
        para.alignment = alignment

        return shape

    def _fill_bullets(
        self,
        text_frame,
        items: List[dict],
        font_size: Optional[int] = None
    ) -> None:
        """Fill a text frame with bullet list content.

        This is the single method for rendering bullet lists, used by both
        placeholder-based slides and custom textbox-based slides.

        Supports inline markdown formatting in item text:
        **bold**, *italic*, ***bold italic***, ~~strikethrough~~,
        __underline__, `code`.

        Accepts anything :func:`body_to_bullets` understands: a markdown
        string, explicit :class:`Bullet` objects, dicts, or bare strings.

        Args:
            text_frame: PowerPoint text frame object (from placeholder or textbox).
            items: Slide body — markdown string or bullet objects.
            font_size: Optional font size for items.

        Returns:
            The bullets actually rendered, so the caller can measure them.
        """
        bullets = body_to_bullets(items)
        if not bullets:
            return []

        text_frame.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()

            para.alignment = PP_ALIGN.LEFT
            para.level = max(0, min(bullet.level, MAX_INDENT_LEVEL) - 1)

            # Apply inline markdown formatting when markers OR escapes are present
            if needs_inline_processing(bullet.text):
                apply_inline_formatting(para, bullet.text, font_size=font_size)
            else:
                para.text = bullet.text
                if font_size:
                    para.font.size = font_size

        return bullets

    # -------------------------------------------------------------------------
    # Table Helpers
    # -------------------------------------------------------------------------

    def _set_cell_fill(self, cell, color) -> None:
        """Set the background fill color of a table cell.

        Args:
            cell: Table cell object.
            color: An ``RGBColor``, or a theme colour name such as "accent1".
                A theme name is written as ``schemeClr`` so the fill follows
                the template's palette instead of being pinned to a literal.
        """
        try:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()

            # Remove existing fill
            for child in list(tcPr):
                if child.tag.endswith('}solidFill'):
                    tcPr.remove(child)

            ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            if isinstance(color, str):
                inner = f'<a:schemeClr val="{color}"/>'
            else:
                inner = f'<a:srgbClr val="{color}"/>'

            tcPr.append(parse_xml(f'<a:solidFill {ns}>{inner}</a:solidFill>'))
        except Exception as e:
            logger.debug(f"Could not set cell fill color: {e}")

    def _create_styled_table(
        self,
        slide,
        table_data: List[List[str]],
        left: int,
        top: int,
        width: int,
        height: int,
        header_color=None,
        alternate_rows: bool = True,
        column_alignments: Optional[List] = None,
        font_size: Optional[int] = None,
    ):
        """Create a styled table on a slide.

        Args:
            slide: PowerPoint slide object.
            table_data: List of rows (first row is header).
            left, top, width, height: Position and size.
            header_color: Header background colour (RGBColor or theme name).
            alternate_rows: Whether to use alternating row colors.
            column_alignments: Optional list of PP_ALIGN values per column
                (extracted from markdown separator row).
            font_size: Explicit cell font size in points. When omitted, the
                size is chosen from the row count so a tall table still fits
                the content area instead of running off the slide.

        Returns:
            Tuple of (table shape, chosen font size in points or None).
        """
        num_rows = len(table_data)
        num_cols = max((len(row) for row in table_data), default=0)

        if num_rows == 0 or num_cols == 0:
            return None, None

        shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = shape.table

        header_color = header_color if header_color is not None else TABLE_HEADER_FILL
        points = font_size or fit_table_font_size(num_rows, height)

        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx >= num_cols:
                    continue

                cell = table.cell(row_idx, col_idx)
                # cell_to_text, not a falsy test: `if cell_text` blanked a numeric 0.
                cell.text = cell_to_text(cell_text)

                paragraph = cell.text_frame.paragraphs[0]

                # Apply column alignment
                if column_alignments and col_idx < len(column_alignments):
                    alignment = column_alignments[col_idx]
                    if alignment is not None:
                        paragraph.alignment = alignment

                if points:
                    paragraph.font.size = Pt(points)

                if row_idx == 0:  # Header row
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = TABLE_HEADER_TEXT
                    self._set_cell_fill(cell, header_color)
                elif alternate_rows and row_idx % 2 == 0:
                    self._set_cell_fill(cell, TABLE_ALT_ROW_FILL)

        return shape, points

    # -------------------------------------------------------------------------
    # Image Helpers
    # -------------------------------------------------------------------------

    def _add_image(
        self,
        slide,
        source: str,
        left: int,
        top: int,
        max_width: int,
        max_height: int,
        center_horizontal: bool = True,
        center_vertical: bool = False
    ) -> Tuple[Optional[Any], Optional[str]]:
        """Add an image from an https URL or an inline data URI.

        Args:
            slide: PowerPoint slide object.
            source: Image URL, or a ``data:image/...;base64,...`` URI.
            left, top: Position.
            max_width, max_height: Bounding box the image is scaled into.
            center_horizontal: Whether to center horizontally.
            center_vertical: Whether to center vertically.

        Returns:
            ``(picture, None)`` on success, or ``(None, reason)`` — the reason
            is surfaced to the caller rather than only logged, because an image
            that quietly failed to load used to produce a confident success
            response with a placeholder box on the slide.
        """
        if not source:
            return None, "no image source given"

        try:
            image_data, _ = load_image(source)

            picture = slide.shapes.add_picture(
                image_data, left, top, width=max_width
            )

            # Scale to fit height if needed
            if picture.height > max_height:
                scale = max_height / picture.height
                picture.width = int(picture.width * scale)
                picture.height = max_height

            # Center if requested
            if center_horizontal:
                slide_width = self.presentation.slide_width
                picture.left = int((slide_width - picture.width) / 2)

            if center_vertical:
                picture.top = int(top + (max_height - picture.height) / 2)

            logger.debug("Added image from %s", source[:80])
            return picture, None

        except (ImageDownloadError, ImageValidationError) as e:
            logger.error("Failed to load image: %s", e)
            return None, str(e)
        except Exception as e:
            logger.error("Failed to add image from %r: %s", source[:80], e, exc_info=True)
            return None, str(e)

    def _add_image_from_url(self, slide, image_url: str, left: int, top: int,
                            max_width: int, max_height: int,
                            center_horizontal: bool = True,
                            center_vertical: bool = False) -> Optional[Any]:
        """Backwards-compatible wrapper returning only the picture."""
        picture, _ = self._add_image(
            slide, image_url, left, top, max_width, max_height,
            center_horizontal, center_vertical,
        )
        return picture

    def _add_image_placeholder(self, slide, message: str, left: int, top: int, width: int):
        """Add a placeholder text when image cannot be loaded.

        Args:
            slide: PowerPoint slide object.
            message: Error message to display.
            left, top, width: Position and width.
        """
        self._add_text_box(
            slide, f"[{message}]",
            left, top, width, Inches(1),
            italic=True, alignment=PP_ALIGN.CENTER
        )

