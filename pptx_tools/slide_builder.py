"""PowerPoint slide builder class.

Builds a deck from the typed slide models in :mod:`pptx_tools.schema`, using
the SlideHelpers mixin for text, tables, images and charts.

Anything the builder had to work around — an image that would not download, a
slide whose text will not fit, a layout the template does not provide — is
recorded on ``self.warnings`` and returned to the caller alongside the file.
Those situations used to be logged server-side only, so the model was told the
deck was fine and had no way to correct itself.
"""

import io
import copy
import logging
from typing import Any, List, Optional, Sequence

from pptx import Presentation
from pptx.dml.color import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

from .constants import (
    DEFAULT_BODY_FONT_SIZE, DEFAULT_SUBTITLE_FONT_SIZE,
    DEFAULT_CAPTION_FONT_SIZE, DEFAULT_QUOTE_FONT_SIZE,
    KPI_VALUE_FONT_SIZE, TIMELINE_DETAIL_FONT_SIZE,
    TIMELINE_STEP_HEIGHT, TIMELINE_STEP_MIN_HEIGHT,
    TIMELINE_DETAIL_GAP, TIMELINE_DETAIL_HEIGHT,
    TABLE_HEADER_FILL,
    DEFAULT_SLIDE_FORMAT, VALID_SLIDE_FORMATS,
)
from .helpers import (
    SlideHelpers,
    apply_autofit, body_to_bullets, estimate_text_fill, parse_table_data,
    resolve_fill, set_runs_language, table_overflows,
)
from .inline_formatting import needs_inline_processing, apply_inline_formatting
from .chart_utils import (
    add_chart_to_slide, add_scatter_to_slide, configure_data_labels,
    set_axis_titles, ChartDataError,
)
from .layouts import LayoutResolver, role_for_slide
from .schema import Bullet, coerce_slides
from .templates import TemplateSpec, open_template, select_template

logger = logging.getLogger(__name__)


class PowerpointPresentation(SlideHelpers):
    """Builder class for creating PowerPoint presentations from structured data."""

    def __init__(self, slides: Sequence[Any], format: str,
                 author: Optional[str] = None,
                 footer_text: Optional[str] = None,
                 show_slide_numbers: bool = False,
                 language: Optional[str] = None,
                 template: Optional[str] = None,
                 template_spec: Optional[TemplateSpec] = None):
        """Initialize and build presentation.

        Args:
            slides: Slide models, or plain dicts in either the current or the
                previous key spelling — both are validated through
                :func:`~pptx_tools.schema.coerce_slides`.
            format: Presentation format ("4:3" or "16:9"). Ignored when
                *template* names a registered template.
            author: Author name stored in document metadata/properties.
            footer_text: Optional footer text displayed on all slides.
            show_slide_numbers: Whether to show slide numbers on all slides.
            language: BCP-47 tag (e.g. "cs-CZ") stamped on every text run so
                the deck is proof-read in the right language.
            template: Name of a registered template. Overrides *format*.
            template_spec: An already-resolved template to build on, bypassing
                the registry lookup. Overrides *template*. This exists for the
                admin UI's preview, which must render a template the admin has
                uploaded but not yet saved — going through the registry would
                mean either finding nothing or writing the entry before the
                admin has agreed to it.
        """
        if not slides:
            raise ValueError("At least one slide is required")

        self.warnings: List[str] = []
        self.slides = coerce_slides(slides)

        logger.info(
            "Initializing PowerPoint: slides=%d, format=%s, template=%s",
            len(self.slides), format, template,
        )

        self.spec = None
        self.presentation = self._create_presentation(format, template, template_spec)
        self._layouts = LayoutResolver(
            self.presentation, self.spec.layouts if self.spec else None
        )

        defaults = self.spec.defaults if self.spec else {}
        self._footer_text = footer_text if footer_text is not None else defaults.get("footer_text")
        self._show_slide_numbers = (
            show_slide_numbers if show_slide_numbers else bool(defaults.get("show_slide_numbers"))
        )
        self._language = language if language is not None else defaults.get("language")
        self._table_defaults = defaults.get("table") or {}
        self._chart_defaults = defaults.get("chart") or {}

        self._remove_template_slides()
        self._build_slides(self.slides)
        if self._footer_text or self._show_slide_numbers:
            self._apply_footer_and_slide_numbers()
        if self._language:
            self._apply_language(self._language)
        if author:
            self.presentation.core_properties.author = author

    # -------------------------------------------------------------------------
    # Setup
    # -------------------------------------------------------------------------

    def _warn(self, slide_index: int, message: str) -> None:
        """Record a caller-visible warning about one slide."""
        entry = f"slide {slide_index}: {message}"
        self.warnings.append(entry)
        logger.warning("[pptx] %s", entry)

    def _create_presentation(self, format: str, template: Optional[str] = None,
                             template_spec: Optional[TemplateSpec] = None) -> Presentation:
        """Create the presentation from the selected registered template."""
        if format not in VALID_SLIDE_FORMATS:
            logger.warning(
                "Unknown presentation format %r; using %s. Valid formats: %s",
                format, DEFAULT_SLIDE_FORMAT, ", ".join(VALID_SLIDE_FORMATS),
            )
            self.warnings.append(
                f"Unknown format {format!r}; used {DEFAULT_SLIDE_FORMAT}."
            )
            format = DEFAULT_SLIDE_FORMAT

        if template_spec is not None:
            spec, note = template_spec, None
        else:
            spec, note = select_template(template, None if template else format)
        if note:
            self.warnings.append(note)

        if spec is not None:
            try:
                presentation = open_template(spec.path)
                self.spec = spec
                logger.info(
                    "Using template %r (%s, %s)", spec.name, spec.path.name, spec.aspect
                )
                return presentation
            except Exception as e:
                logger.error("Failed to load template %s: %s", spec.path.name, e)
                self.warnings.append(
                    f"Template {spec.name!r} could not be opened ({e}); "
                    "used the built-in PowerPoint theme."
                )

        logger.warning("Using the built-in PowerPoint theme for %s", format)
        return Presentation()

    def _apply_title(self, slide, text, index: int) -> None:
        """Set the slide title, saying so when the layout cannot hold one.

        Every other dropped element on a slide is reported — a subtitle, the
        bullets, the footer. The title was the one that vanished in silence,
        which is the failure this phase exists to remove, not to introduce
        somewhere new.
        """
        if self._set_title(slide, text):
            return
        if text:
            self._warn(
                index,
                f"layout {slide.slide_layout.name!r} has no title placeholder; "
                f"the title {text!r} was dropped.",
            )

    def _content_slide(self, slide_data, index: int):
        """Resolve, create and title a slide, and return its content rectangle.

        Shared by every slide type that draws into the body area itself
        (table, image, chart, scatter, quote).
        """
        slide = self._new_slide(slide_data, index)
        self._apply_title(slide, slide_data.title, index)
        return self._add_title_content_slide("", slide=slide)

    def _new_slide(self, slide_data, index: int):
        """Add a slide on the layout resolved for its type.

        This is what replaced indexing ``slide_layouts`` by position. A
        template that reorders or omits layouts is now matched by name and by
        placeholder signature, and anything approximate about the match is
        reported to the caller rather than silently producing a deck laid out
        on the wrong layouts.
        """
        role = role_for_slide(slide_data)
        layout, note = self._layouts.resolve(role, slide_data.layout)
        if note:
            self._warn(index, note)
        return self.presentation.slides.add_slide(layout)

    def _remove_template_slides(self) -> None:
        """Remove every slide the template ships with, parts included.

        Dropping only the ``<p:sldId>`` entry leaves the slide's relationship in
        place, and python-pptx serialises parts by walking relationships — so a
        template carrying a sample slide shipped that slide's XML, text and
        images inside every generated file even though PowerPoint showed the
        right slide count. Dropping the relationship removes the part too.

        A template can opt out with ``strip_slides: false`` when its own slides
        are meant to survive — a fixed cover or back page the generated slides
        should follow. Until this check existed the option was parsed, stored
        and offered in the admin UI but never read, so unticking the box
        changed nothing and said nothing.
        """
        if self.spec is not None and not self.spec.strip_slides:
            logger.debug(
                "Template %r sets strip_slides: false; keeping its %d slide(s)",
                self.spec.name, len(self.presentation.slides._sldIdLst),
            )
            return

        sldIdLst = self.presentation.slides._sldIdLst
        prs_part = self.presentation.part

        removed = 0
        for sldId in list(sldIdLst):
            rId = sldId.rId
            try:
                sldIdLst.remove(sldId)
                prs_part.drop_rel(rId)
                removed += 1
            except Exception as e:
                logger.warning("Could not fully remove template slide %s: %s", rId, e)

        if removed:
            logger.debug("Removed %d slide(s) carried by the template", removed)

    def _build_slides(self, slides: Sequence[Any]) -> None:
        """Build all slides from validated models."""
        builders = {
            "title": self._build_title_slide,
            "section": self._build_section_slide,
            "content": self._build_content_slide,
            "table": self._build_table_slide,
            "image": self._build_image_slide,
            "two_column": self._build_two_column_slide,
            "chart": self._build_chart_slide,
            "scatter": self._build_scatter_slide,
            "quote": self._build_quote_slide,
            "kpi": self._build_kpi_slide,
            "agenda": self._build_agenda_slide,
            "closing": self._build_closing_slide,
            "timeline": self._build_timeline_slide,
        }

        logger.info("Building %d slides", len(slides))

        for i, slide in enumerate(slides):
            builder = builders.get(slide.type)
            if builder is None:  # pragma: no cover - schema forbids it
                raise ValueError(f"No builder for slide type {slide.type!r} at slide {i}")

            try:
                logger.debug("Building slide %d: type=%s", i, slide.type)
                builder(slide, i)
            except Exception as e:
                logger.error("Failed to create slide %d: %s", i, e)
                raise ValueError(f"Error creating slide {i} ({slide.type}): {e}")

    # -------------------------------------------------------------------------
    # Slide Builders
    # -------------------------------------------------------------------------

    def _build_title_slide(self, slide_data, index: int) -> None:
        """Build a title slide with title and optional subtitle."""
        slide = self._new_slide(slide_data, index)

        self._apply_title(slide, slide_data.title, index)

        subtitle = self._placeholder_of_type(
            slide, (PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
        )
        if subtitle is not None:
            subtitle.text = slide_data.subtitle or ""
        elif slide_data.subtitle:
            self._warn(index, "this layout has no subtitle placeholder; the subtitle was dropped.")

        self._add_speaker_notes(slide, slide_data.notes)

    def _build_section_slide(self, slide_data, index: int) -> None:
        """Build a section divider slide."""
        slide = self._new_slide(slide_data, index)
        self._apply_title(slide, slide_data.title, index)
        self._add_speaker_notes(slide, slide_data.notes)

    def _build_content_slide(self, slide_data, index: int) -> None:
        """Build a content slide with bullet points."""
        slide = self._new_slide(slide_data, index)
        self._apply_title(slide, slide_data.title, index)

        bullets = body_to_bullets(slide_data.body)
        if bullets:
            placeholders = self._content_placeholders(slide)
            if placeholders:
                placeholder = placeholders[0]
                placeholder.text = ""
                self._fill_bullets(placeholder.text_frame, bullets)
                self._fit_text(placeholder, bullets, index)
            else:
                self._warn(index, "this layout has no body placeholder; the bullets were dropped.")

        self._add_speaker_notes(slide, slide_data.notes)

    def _build_table_slide(self, slide_data, index: int) -> None:
        """Build a table slide with a styled table."""
        slide, left, top, width, height = self._content_slide(slide_data, index)

        rows, col_alignments = parse_table_data(slide_data.rows)
        if not rows:
            self._warn(index, "table has no rows; the slide is empty.")
            return

        # An explicit `align` beats an inline markdown separator row.
        if slide_data.align:
            col_alignments = [
                {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[a]
                for a in slide_data.align
            ]

        header_fill = resolve_fill(slide_data.header_color, TABLE_HEADER_FILL)

        _, points = self._create_styled_table(
            slide,
            rows,
            left=left,
            top=top,
            width=width,
            height=height,
            header_color=header_fill,
            alternate_rows=slide_data.zebra,
            column_alignments=col_alignments,
            font_size=slide_data.font_size,
        )

        if points and table_overflows(len(rows), height, points):
            self._warn(
                index,
                f"table of {len(rows)} rows will not fit the content area even at "
                f"{points}pt; split it across slides.",
            )
        elif slide_data.font_size is None and points and points < int(DEFAULT_BODY_FONT_SIZE.pt):
            self._warn(
                index,
                f"table font reduced to {points}pt to fit {len(rows)} rows.",
            )

        self._add_speaker_notes(slide, slide_data.notes)

    def _build_image_slide(self, slide_data, index: int) -> None:
        """Build a slide with an image from a URL or inline data URI."""
        slide, left, top, width, height = self._content_slide(slide_data, index)

        caption = slide_data.caption
        bullets = body_to_bullets(slide_data.body)

        # With text beside it the picture takes the left 55%, the text the rest.
        image_width = width
        if bullets:
            image_width = int(width * 0.55)
            gutter = Inches(0.3)
            text_left = left + image_width + gutter
            text_width = width - image_width - gutter
            box = slide.shapes.add_textbox(text_left, top, text_width, height)
            self._fill_bullets(box.text_frame, bullets)
            apply_autofit(
                box.text_frame,
                scale=self._fit_scale(bullets, text_width, height),
            )

        max_height = height - (Inches(0.6) if caption else 0)

        picture, error = self._add_image(
            slide, slide_data.source,
            left=left, top=top, max_width=image_width, max_height=max_height,
            center_horizontal=not bullets,
        )

        if picture and caption:
            self._add_text_box(
                slide, caption,
                left=picture.left,
                top=picture.top + picture.height + Inches(0.1),
                width=picture.width,
                height=Inches(0.5),
                font_size=DEFAULT_CAPTION_FONT_SIZE,
                italic=True,
                alignment=PP_ALIGN.CENTER,
            )
        elif not picture:
            self._add_image_placeholder(
                slide, "Image could not be loaded", left, top + Inches(1), width
            )
            self._warn(index, f"image could not be loaded ({error}); a placeholder was drawn.")

        self._add_speaker_notes(slide, slide_data.notes)

    def _build_two_column_slide(self, slide_data, index: int) -> None:
        """Build a slide with two text columns using built-in PowerPoint layouts.

        Uses the Comparison layout when either column has a heading, otherwise
        Two Content.

        Placeholder indices:
        - Two Content (3): idx 0=Title, 1=Left content, 2=Right content
        - Comparison (4): idx 0=Title, 1=Left heading, 2=Left content,
          3=Right heading, 4=Right content
        """
        left_col, right_col = slide_data.left, slide_data.right
        has_headings = bool(left_col.heading or right_col.heading)

        slide = self._new_slide(slide_data, index)
        self._apply_title(slide, slide_data.title, index)

        left_bullets = body_to_bullets(left_col.body)
        right_bullets = body_to_bullets(right_col.body)

        if has_headings:
            content_slots = {2: left_bullets, 4: right_bullets}
            heading_slots = {1: left_col.heading, 3: right_col.heading}
        else:
            content_slots = {1: left_bullets, 2: right_bullets}
            heading_slots = {}

        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx

            if idx == 0:
                continue  # already set through _apply_title
            elif idx in heading_slots:
                if heading_slots[idx]:
                    shape.text = heading_slots[idx]
            elif idx in content_slots:
                bullets = content_slots[idx]
                if bullets:
                    self._fill_bullets(shape.text_frame, bullets)
                    self._fit_text(shape, bullets, index)

        self._add_speaker_notes(slide, slide_data.notes)

    def _build_chart_slide(self, slide_data, index: int) -> None:
        """Build a slide with a category chart."""
        slide, left, top, width, height = self._content_slide(slide_data, index)

        bullets = body_to_bullets(slide_data.body)
        if bullets:
            chart_width = int(width * 0.6)
            gutter = Inches(0.3)
            text_left = left + chart_width + gutter
            text_width = width - chart_width - gutter
            box = slide.shapes.add_textbox(text_left, top, text_width, height)
            self._fill_bullets(box.text_frame, bullets)
            apply_autofit(
                box.text_frame,
                scale=self._fit_scale(bullets, text_width, height),
            )
            width = chart_width

        chart_data = {
            "categories": slide_data.categories,
            "series": [
                {"name": s.name, "values": s.values} for s in slide_data.series
            ],
        }

        try:
            chart = add_chart_to_slide(
                slide,
                chart_type=slide_data.chart_type,
                chart_data=chart_data,
                left=left, top=top, width=width, height=height,
                has_legend=slide_data.legend != "none",
                legend_position=slide_data.legend,
                title=slide_data.chart_title,
            )
            configure_data_labels(chart, slide_data.data_labels, slide_data.number_format)
            set_axis_titles(chart, slide_data.x_title, slide_data.y_title)
        except ChartDataError as e:
            logger.error(f"Chart error: {e}")
            self._add_text_box(
                slide, f"[Chart error: {e}]",
                left, top, width, Inches(1), alignment=PP_ALIGN.CENTER,
            )
            self._warn(index, f"chart could not be built ({e}).")
            self._add_speaker_notes(slide, slide_data.notes)
            return

        # A series whose length disagrees with the categories is accepted by
        # python-pptx but renders with gaps, so say so rather than let it pass.
        for series in slide_data.series:
            if len(series.values) != len(slide_data.categories):
                self._warn(
                    index,
                    f"series {series.name!r} has {len(series.values)} values for "
                    f"{len(slide_data.categories)} categories.",
                )

        self._add_speaker_notes(slide, slide_data.notes)

    def _build_scatter_slide(self, slide_data, index: int) -> None:
        """Build a slide with an XY (scatter) chart."""
        slide, left, top, width, height = self._content_slide(slide_data, index)

        try:
            add_scatter_to_slide(
                slide,
                series=slide_data.series,
                left=left, top=top, width=width, height=height,
                legend=slide_data.legend,
                title=slide_data.chart_title,
                x_title=slide_data.x_title,
                y_title=slide_data.y_title,
            )
        except ChartDataError as e:
            logger.error(f"Scatter chart error: {e}")
            self._add_text_box(
                slide, f"[Chart error: {e}]",
                left, top, width, Inches(1), alignment=PP_ALIGN.CENTER,
            )
            self._warn(index, f"scatter chart could not be built ({e}).")

        self._add_speaker_notes(slide, slide_data.notes)

    def _build_quote_slide(self, slide_data, index: int) -> None:
        """Build a quote/citation slide."""
        slide, left, top, width, height = self._content_slide(slide_data, index)

        quote_box = slide.shapes.add_textbox(left, top, width, height)
        tf = quote_box.text_frame
        tf.word_wrap = True

        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        formatted_quote = f'"{slide_data.text}"'
        if needs_inline_processing(formatted_quote):
            apply_inline_formatting(para, formatted_quote,
                                    font_size=DEFAULT_QUOTE_FONT_SIZE, italic=True)
        else:
            para.text = formatted_quote
            para.font.size = DEFAULT_QUOTE_FONT_SIZE
            para.font.italic = True

        if slide_data.attribution:
            author_para = tf.add_paragraph()
            author_para.text = f"— {slide_data.attribution}"
            author_para.font.size = DEFAULT_SUBTITLE_FONT_SIZE
            author_para.font.bold = True
            author_para.alignment = PP_ALIGN.CENTER
            author_para.space_before = Pt(24)

        self._add_speaker_notes(slide, slide_data.notes)

    # -------------------------------------------------------------------------
    # Slide types that draw their own shapes
    # -------------------------------------------------------------------------

    def _build_kpi_slide(self, slide_data, index: int) -> None:
        """Build a row of headline figures.

        Drawn rather than placed in a body placeholder: the point of a KPI row
        is the typographic contrast between a large figure and a small label,
        which a bullet list cannot express.
        """
        slide, left, top, width, height = self._content_slide(slide_data, index)

        items = slide_data.items
        gutter = Inches(0.25)
        cell_width = int((width - gutter * (len(items) - 1)) / len(items))
        # Sit the row a little above centre; a KPI slide reads as a headline.
        block_height = min(height, Inches(2.4))
        block_top = top + int(max(0, (height - block_height)) / 3)

        for position, item in enumerate(items):
            cell_left = left + position * (cell_width + gutter)
            box = slide.shapes.add_textbox(cell_left, block_top, cell_width, block_height)
            frame = box.text_frame
            frame.word_wrap = True

            value = frame.paragraphs[0]
            value.text = item.value
            value.alignment = PP_ALIGN.CENTER
            value.font.size = KPI_VALUE_FONT_SIZE
            value.font.bold = True

            label = frame.add_paragraph()
            label.text = item.label
            label.alignment = PP_ALIGN.CENTER
            label.font.size = DEFAULT_CAPTION_FONT_SIZE
            label.space_before = Pt(4)

            if item.delta:
                delta = frame.add_paragraph()
                delta.text = item.delta
                delta.alignment = PP_ALIGN.CENTER
                delta.font.size = DEFAULT_CAPTION_FONT_SIZE
                delta.font.italic = True
                delta.space_before = Pt(2)
                # Theme accent, so the figure follows the template's palette.
                delta.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1

        if len(items) > 4:
            self._warn(
                index,
                f"{len(items)} figures on one KPI slide will be cramped; "
                "two to four read best.",
            )

        self._add_speaker_notes(slide, slide_data.notes)

    def _build_agenda_slide(self, slide_data, index: int) -> None:
        """Build a numbered agenda.

        With no explicit items the entries are taken from the deck's own
        ``section`` slides, in order — which is the usual case and cannot drift
        out of step with the deck the way a hand-written list does.
        """
        items = slide_data.items
        derived = False
        if items is None:
            items = [s.title for s in self.slides if s.type == "section" and s.title]
            derived = True

        if not items:
            self._warn(
                index,
                "agenda has no items and the deck has no section slides to derive "
                "them from; the slide is empty.",
            )

        slide = self._new_slide(slide_data, index)
        self._apply_title(slide, slide_data.title or "Agenda", index)

        if items:
            bullets = [Bullet(text=f"{n}.  {text}") for n, text in enumerate(items, 1)]
            placeholders = self._content_placeholders(slide)
            if placeholders:
                placeholders[0].text = ""
                self._fill_bullets(placeholders[0].text_frame, bullets)
                self._fit_text(placeholders[0], bullets, index)
            else:
                self._warn(index, "this layout has no body placeholder; the agenda was dropped.")

        if derived and items:
            logger.debug("Agenda derived from %d section slides", len(items))

        self._add_speaker_notes(slide, slide_data.notes)

    def _build_closing_slide(self, slide_data, index: int) -> None:
        """Build a closing / thank-you slide."""
        slide = self._new_slide(slide_data, index)
        self._apply_title(slide, slide_data.title or "Thank you", index)

        lines = []
        if slide_data.subtitle:
            lines.append(slide_data.subtitle)
        lines.extend(slide_data.contact or [])

        if lines:
            target = self._placeholder_of_type(
                slide, (PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT)
            )
            if target is not None:
                target.text = ""
                self._fill_bullets(target.text_frame, [Bullet(text=line) for line in lines])
            else:
                self._warn(
                    index,
                    "this layout has no subtitle or body placeholder; the closing "
                    "lines were dropped.",
                )

        self._add_speaker_notes(slide, slide_data.notes)

    def _timeline_detail_band(self, steps, height, index: int):
        """Height to reserve under the shapes for step detail, as (gap, height).

        Returns ``(0, 0)`` when no step has detail, or when the content
        rectangle is too short to carry a legible caption — in which case the
        detail is dropped and reported, rather than drawn off the slide.
        """
        if not any(step.detail for step in steps):
            return 0, 0

        wanted = TIMELINE_DETAIL_GAP + TIMELINE_DETAIL_HEIGHT
        if height - wanted >= TIMELINE_STEP_MIN_HEIGHT:
            return TIMELINE_DETAIL_GAP, TIMELINE_DETAIL_HEIGHT

        self._warn(
            index,
            "the content area is too short to fit step detail under the timeline "
            "shapes; the detail lines were dropped. Use a layout with a taller "
            "body area, or move the detail into speaker notes.",
        )
        return 0, 0

    def _build_timeline_slide(self, slide_data, index: int) -> None:
        """Build a row of steps as chevrons or boxes.

        Autoshapes rather than SmartArt: python-pptx cannot create SmartArt, and
        a row of chevrons carries the same "these follow one another" meaning
        without the dependency.
        """
        slide, left, top, width, height = self._content_slide(slide_data, index)

        steps = slide_data.steps
        shape_type = MSO_SHAPE.CHEVRON if slide_data.style == "chevron" else MSO_SHAPE.ROUNDED_RECTANGLE
        # Chevrons interlock, so they overlap slightly; boxes get a gutter.
        gutter = Inches(-0.12) if slide_data.style == "chevron" else Inches(0.15)

        step_width = int((width - gutter * (len(steps) - 1)) / len(steps))

        # Reserve the detail band up front. Sizing the shapes first and then
        # hanging the captions underneath overflowed the content rectangle on
        # any layout with a short body placeholder — an explicit ``layout``
        # override onto a section-header layout put the detail lines 0.54in
        # below the box, and a custom template with a low, short placeholder
        # would put them off the slide entirely.
        gap, detail_height = self._timeline_detail_band(steps, height, index)
        band = gap + detail_height
        step_height = min(TIMELINE_STEP_HEIGHT, max(TIMELINE_STEP_MIN_HEIGHT, height - band))
        step_top = top + int(max(0, (height - (step_height + band))) / 3)

        for position, step in enumerate(steps):
            shape = slide.shapes.add_shape(
                shape_type,
                left + position * (step_width + gutter),
                step_top, step_width, step_height,
            )
            shape.fill.solid()
            shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
            shape.line.fill.background()

            frame = shape.text_frame
            frame.word_wrap = True
            label = frame.paragraphs[0]
            label.text = step.label
            label.alignment = PP_ALIGN.CENTER
            label.font.size = DEFAULT_CAPTION_FONT_SIZE
            label.font.bold = True

            if step.detail and detail_height:
                # Detail below the shape, so a long line cannot burst the chevron.
                caption = slide.shapes.add_textbox(
                    left + position * (step_width + gutter),
                    step_top + step_height + gap,
                    step_width, detail_height,
                )
                caption.text_frame.word_wrap = True
                detail = caption.text_frame.paragraphs[0]
                detail.text = step.detail
                detail.alignment = PP_ALIGN.CENTER
                detail.font.size = TIMELINE_DETAIL_FONT_SIZE

        self._add_speaker_notes(slide, slide_data.notes)

    # -------------------------------------------------------------------------
    # Fit, language, footer
    # -------------------------------------------------------------------------

    def _fit_scale(self, bullets, width, height):
        """Shrink factor for a text box, or None when the text already fits."""
        fill = estimate_text_fill(
            bullets, width, height, font_size_pt=float(DEFAULT_BODY_FONT_SIZE.pt)
        )
        return (1.0 / fill) if fill > 1.0 else None

    def _fit_text(self, placeholder, bullets, index: int) -> None:
        """Ask PowerPoint to shrink overfull body text, and warn when it is far gone."""
        fill = estimate_text_fill(
            bullets, placeholder.width, placeholder.height,
            font_size_pt=float(DEFAULT_BODY_FONT_SIZE.pt),
        )
        apply_autofit(placeholder.text_frame, scale=(1.0 / fill) if fill > 1.0 else None)

        # Below the shrink floor the slide is genuinely overfull; shrinking
        # further would produce text nobody can read from a room.
        if fill > 1.0 / 0.6:
            self._warn(
                index,
                f"body text is about {fill:.1f}x the space available and will be "
                "shrunk to fit; consider splitting it across slides.",
            )

    def _apply_language(self, language: str) -> None:
        """Stamp the proofing language on every run in the deck."""
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    set_runs_language(shape.text_frame, language)
                elif getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            set_runs_language(cell.text_frame, language)
            if slide.has_notes_slide:
                set_runs_language(slide.notes_slide.notes_text_frame, language)

        logger.debug("Applied language %s to all runs", language)

    def _apply_footer_and_slide_numbers(self) -> None:
        """Apply footer text and/or slide numbers to all slides.

        Clones the footer/slide-number placeholder shapes from each slide's
        layout into the slide itself so they become visible. Assigns unique
        shape IDs to avoid PPTX corruption.
        """
        from xml.sax.saxutils import escape as xml_escape

        missing_footer = 0

        for slide in self.presentation.slides:
            layout = slide.slide_layout
            spTree = slide.shapes._spTree

            # Determine next available shape ID on this slide by collecting the
            # id attributes of the tree's direct children.
            existing_ids = set()
            for sp in spTree:
                cNvPr = sp.find('.//' + qn('p:cNvPr'))
                if cNvPr is None:
                    cNvPr = sp.find('.//' + qn('p:nvSpPr') + '/' + qn('p:cNvPr'))
                if cNvPr is not None and cNvPr.get('id'):
                    existing_ids.add(int(cNvPr.get('id')))
            next_id = max(existing_ids, default=0) + 1

            layout_indices = {ph.placeholder_format.idx for ph in layout.placeholders}
            if self._footer_text and 11 not in layout_indices:
                missing_footer += 1

            for ph in layout.placeholders:
                idx = ph.placeholder_format.idx

                if idx == 11 and self._footer_text:  # FOOTER placeholder
                    sp = copy.deepcopy(ph._element)
                    cNvPr = sp.find(qn('p:nvSpPr') + '/' + qn('p:cNvPr'))
                    if cNvPr is not None:
                        cNvPr.set('id', str(next_id))
                        next_id += 1
                    txBody = sp.find(qn('p:txBody'))
                    if txBody is not None:
                        for p in txBody.findall(qn('a:p')):
                            txBody.remove(p)
                        ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
                        safe_text = xml_escape(self._footer_text)
                        p_xml = f'<a:p {ns}><a:r><a:t>{safe_text}</a:t></a:r></a:p>'
                        txBody.append(parse_xml(p_xml))
                    spTree.append(sp)

                elif idx == 12 and self._show_slide_numbers:  # SLIDE_NUMBER placeholder
                    sp = copy.deepcopy(ph._element)
                    cNvPr = sp.find(qn('p:nvSpPr') + '/' + qn('p:cNvPr'))
                    if cNvPr is not None:
                        cNvPr.set('id', str(next_id))
                        next_id += 1
                    spTree.append(sp)

        if missing_footer:
            # Previously this just produced a deck with no footer and no
            # explanation of why the argument appeared to do nothing.
            self.warnings.append(
                f"footer_text was dropped on {missing_footer} slide(s): their layout "
                "has no footer placeholder."
            )

        logger.debug("Applied footer/slide numbers to all slides")

    # -------------------------------------------------------------------------
    # Output
    # -------------------------------------------------------------------------

    def save(self) -> io.BytesIO:
        """Save presentation to a BytesIO object.

        Returns:
            BytesIO containing the presentation.

        Raises:
            RuntimeError: If saving fails.
        """
        logger.info("Saving PowerPoint to memory buffer")
        try:
            buffer = io.BytesIO()
            self.presentation.save(buffer)
            buffer.seek(0)
            return buffer
        except Exception as e:
            logger.error("Failed to save PowerPoint presentation: %s", e, exc_info=True)
            raise RuntimeError(f"Failed to save presentation: {e}") from e
