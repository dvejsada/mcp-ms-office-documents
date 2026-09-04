"""Tests for the template registry and layout resolution (review Phase 2).

The bug this phase exists to remove was silent: layouts were addressed by
position, so a customer template that reordered them built the deck's title
slide on the section layout and said nothing. These tests build deliberately
awkward templates — reordered, trimmed, renamed, .potx — and assert on the
layout each slide actually landed on.
"""

import shutil
import sys
import time
import zipfile
from pathlib import Path

project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

import pytest
import yaml
from pptx import Presentation as PptxReader
from pptx.util import Emu

from pptx_tools import layouts as layouts_mod
from pptx_tools import templates as templates_mod
from pptx_tools.layouts import (
    LayoutResolver, ROLE_COMPARISON, ROLE_CONTENT, ROLE_SECTION, ROLE_TITLE,
    ROLE_TWO_COLUMN, classify_layout,
)
from pptx_tools.slide_builder import PowerpointPresentation
from pptx_tools.templates import (
    TemplateSpec, aspect_of, load_specs, open_template, select_template,
    validate_templates,
)

BASE_16_9 = project_root / "default_templates" / "default_pptx_template_16_9.pptx"
BASE_4_3 = project_root / "default_templates" / "default_pptx_template_4_3.pptx"

DECK = [
    {"type": "title", "title": "Deck title", "subtitle": "sub"},
    {"type": "section", "title": "Part 1"},
    {"type": "content", "title": "Body", "body": "- a"},
    {"type": "two_column", "title": "2c", "left": {"body": "- l"}, "right": {"body": "- r"}},
]


# =============================================================================
# Fixtures
# =============================================================================

def _layout_ids(prs):
    return prs.slide_masters[0]._element.sldLayoutIdLst


def make_template(path: Path, mutate=None) -> Path:
    """Copy the shipped 16:9 template, optionally mutating its layouts."""
    prs = PptxReader(str(BASE_16_9))
    if mutate:
        mutate(prs)
    prs.save(str(path))
    return path


def reorder_layouts(prs):
    """Move the Section Header layout to position 0."""
    ids = _layout_ids(prs)
    section = list(ids)[2]
    ids.remove(section)
    ids.insert(0, section)


def trim_layouts(prs):
    """Keep only the first three layouts (title, content, section)."""
    ids = _layout_ids(prs)
    for lid in list(ids)[3:]:
        ids.remove(lid)


@pytest.fixture
def registry(tmp_path, monkeypatch):
    """Point the registry and template search at a temporary directory."""
    config = tmp_path / "config"
    custom = tmp_path / "custom_templates"
    config.mkdir()
    custom.mkdir()

    monkeypatch.setattr(templates_mod, "_APP_CONFIG_DIR", tmp_path / "nonexistent")
    monkeypatch.setattr(templates_mod, "_LOCAL_CONFIG_DIR", config)

    import template_utils
    monkeypatch.setattr(template_utils, "_candidate_dirs", lambda: [custom, project_root / "default_templates"])

    templates_mod.clear_cache()
    yield {"config": config, "custom": custom}
    templates_mod.clear_cache()


def write_registry(config: Path, entries):
    (config / "pptx_templates.yaml").write_text(
        yaml.safe_dump({"templates": entries}), encoding="utf-8"
    )


def build(slides, fmt="16:9", **kwargs):
    return PowerpointPresentation(slides, fmt, **kwargs)


def layout_names_of(pres):
    return [slide.slide_layout.name for slide in pres.presentation.slides]


# =============================================================================
# Layout classification
# =============================================================================

class TestClassification:

    @pytest.mark.parametrize("path", [BASE_16_9, BASE_4_3])
    def test_shipped_templates_cover_every_role(self, path):
        resolver = LayoutResolver(PptxReader(str(path)))
        assert resolver.missing_roles() == []
        coverage = resolver.coverage()
        assert all(coverage[role] for role in
                   ("title", "section", "content", "two_column", "comparison", "image_text"))

    def test_classification_is_language_independent(self):
        """The shipped templates are Czech-named; detection reads types, not names."""
        prs = PptxReader(str(BASE_16_9))
        roles = {layout.name: classify_layout(layout) for layout in prs.slide_layouts}
        assert roles["Úvodní snímek"] == ROLE_TITLE
        assert roles["Záhlaví oddílu"] == ROLE_SECTION
        assert roles["Nadpis a obsah"] == ROLE_CONTENT
        assert roles["Dva obsahy"] == ROLE_TWO_COLUMN
        assert roles["Porovnání"] == ROLE_COMPARISON

    def test_vertical_layouts_are_never_auto_selected(self):
        """They share a signature with Section Header / Title and Content."""
        prs = PptxReader(str(BASE_16_9))
        roles = {layout.name: classify_layout(layout) for layout in prs.slide_layouts}
        assert roles["Nadpis a svislý text"] is None
        assert roles["Svislý nadpis a text"] is None

    def test_content_with_caption_is_not_mistaken_for_two_column(self):
        prs = PptxReader(str(BASE_16_9))
        roles = {layout.name: classify_layout(layout) for layout in prs.slide_layouts}
        assert roles["Obsah s titulkem"] is None

    def test_picture_layout_without_a_title_is_refused(self):
        """A full-bleed photo layout has nowhere to put the slide's title.

        The picture branch used to return image_text before the "no title,
        refuse" guard every other role goes through, so a titleless photo
        layout was picked confidently and then dropped the title.
        """
        from types import SimpleNamespace
        from pptx.enum.shapes import PP_PLACEHOLDER as PH

        class FakePlaceholder:
            def __init__(self, kind):
                self.placeholder_format = SimpleNamespace(type=kind, idx=0)
                self.text_frame = None

        class FakeLayout:
            def __init__(self, kinds):
                self.placeholders = [FakePlaceholder(k) for k in kinds]
                self.name = "fake"

        assert classify_layout(FakeLayout([PH.PICTURE])) is None
        assert classify_layout(FakeLayout([PH.PICTURE, PH.BODY])) is None
        assert classify_layout(FakeLayout([PH.TITLE, PH.PICTURE])) == "image_text"

    def test_picture_layout_detected(self):
        prs = PptxReader(str(BASE_16_9))
        roles = {layout.name: classify_layout(layout) for layout in prs.slide_layouts}
        assert roles["Obrázek s titulkem"] == "image_text"


# =============================================================================
# The regression this phase exists for
# =============================================================================

class TestAwkwardTemplates:

    def test_reordered_template_still_places_every_slide(self, tmp_path, registry):
        """Was: the title slide got built on the section layout, silently."""
        make_template(registry["custom"] / "custom_pptx_template_16_9.pptx", reorder_layouts)

        pres = build(DECK)

        assert layout_names_of(pres) == [
            "Úvodní snímek", "Záhlaví oddílu", "Nadpis a obsah", "Dva obsahy",
        ]
        assert pres.warnings == []

    def test_trimmed_template_does_not_raise(self, tmp_path, registry):
        """Was: ValueError('slide layout index out of range')."""
        make_template(registry["custom"] / "custom_pptx_template_16_9.pptx", trim_layouts)

        pres = build(DECK)

        assert len(pres.presentation.slides) == len(DECK)
        # The three roles the trimmed file still provides resolve exactly; only
        # two_column, which it no longer has, falls back.
        assert layout_names_of(pres)[:3] == [
            "Úvodní snímek", "Záhlaví oddílu", "Nadpis a obsah",
        ]

    def test_trimmed_template_reports_the_approximation(self, tmp_path, registry):
        make_template(registry["custom"] / "custom_pptx_template_16_9.pptx", trim_layouts)

        pres = build(DECK)

        assert any("two_column" in w and "slide 3" in w for w in pres.warnings)


# =============================================================================
# Resolution order
# =============================================================================

class TestResolutionOrder:

    def test_per_slide_layout_wins(self, registry):
        """An explicit name reaches a layout auto-detection deliberately refuses.

        "Obsah s titulkem" (Content with Caption) is classified as None, so it
        can only be selected by naming it — which is the point of the field.
        """
        make_template(registry["custom"] / "custom_pptx_template_16_9.pptx")

        pres = build([{"type": "content", "title": "C", "body": "- a",
                       "layout": "Obsah s titulkem"}])

        assert layout_names_of(pres) == ["Obsah s titulkem"]
        assert pres.warnings == []

    def test_layout_without_a_body_reports_the_dropped_bullets(self, registry):
        """Choosing Title Only for a slide with bullets loses them; say so."""
        make_template(registry["custom"] / "custom_pptx_template_16_9.pptx")

        pres = build([{"type": "content", "title": "C", "body": "- a",
                       "layout": "Jenom nadpis"}])

        assert layout_names_of(pres) == ["Jenom nadpis"]
        assert any("no body placeholder" in w for w in pres.warnings)

    def test_per_slide_layout_is_case_insensitive(self, registry):
        make_template(registry["custom"] / "custom_pptx_template_16_9.pptx")
        pres = build([{"type": "content", "title": "C", "layout": "jenom nadpis"}])
        assert layout_names_of(pres) == ["Jenom nadpis"]

    def test_unknown_layout_name_warns_and_falls_back(self, registry):
        make_template(registry["custom"] / "custom_pptx_template_16_9.pptx")

        pres = build([{"type": "content", "title": "C", "layout": "No Such Layout"}])

        assert layout_names_of(pres) == ["Nadpis a obsah"]
        assert any("No Such Layout" in w for w in pres.warnings)

    def test_configured_mapping_overrides_detection(self, registry):
        make_template(registry["custom"] / "brand.pptx")
        write_registry(registry["config"], [{
            "name": "brand",
            "pptx_path": "brand.pptx",
            "default": True,
            "layouts": {"content": "Jenom nadpis"},
        }])

        pres = build([{"type": "content", "title": "C"}], template="brand")

        assert layout_names_of(pres) == ["Jenom nadpis"]

    def test_configured_layout_that_does_not_exist_is_reported(self, registry):
        make_template(registry["custom"] / "brand.pptx")
        write_registry(registry["config"], [{
            "name": "brand", "pptx_path": "brand.pptx", "default": True,
            "layouts": {"content": "Missing Layout"},
        }])

        pres = build([{"type": "content", "title": "C"}], template="brand")

        assert any("Missing Layout" in w for w in pres.warnings)
        # Still built, on the detected layout.
        assert layout_names_of(pres) == ["Nadpis a obsah"]

    def test_a_dropped_title_is_reported(self, registry):
        """Every other dropped element warns; the title used to vanish silently."""
        make_template(registry["custom"] / "custom_pptx_template_16_9.pptx")

        pres = build([{"type": "content", "title": "THIS TITLE MATTERS",
                       "body": "- a", "layout": "Prázdný"}])

        slide = pres.presentation.slides[0]
        assert not any(
            "THIS TITLE MATTERS" in shape.text_frame.text
            for shape in slide.shapes if shape.has_text_frame
        )
        assert any("THIS TITLE MATTERS" in w and "no title placeholder" in w
                   for w in pres.warnings)

    @pytest.mark.parametrize("slide_type,extra", [
        ("content", {"body": "- a"}),
        ("table", {"rows": [["A"], ["1"]]}),
        ("quote", {"text": "q"}),
        ("section", {}),
        ("two_column", {"left": {"body": "- l"}, "right": {"body": "- r"}}),
    ])
    def test_every_slide_type_reports_a_dropped_title(self, registry, slide_type, extra):
        make_template(registry["custom"] / "custom_pptx_template_16_9.pptx")

        pres = build([{"type": slide_type, "title": "Gone", "layout": "Prázdný", **extra}])

        assert any("no title placeholder" in w for w in pres.warnings)

    def test_content_goes_in_the_largest_placeholder(self, registry):
        """On Comparison the first content placeholder is the heading strip.

        Taking the first put a table into a 0.9-inch band instead of the
        4-inch content box.
        """
        make_template(registry["custom"] / "custom_pptx_template_16_9.pptx")

        pres = build([{"type": "table", "title": "T", "rows": [["A"], ["1"]],
                       "layout": "Porovnání"}])

        table_shape = [s for s in pres.presentation.slides[0].shapes if s.has_table][0]
        assert Emu(table_shape.height).inches > 3.0

    def test_two_column_picks_comparison_when_headings_present(self, registry):
        make_template(registry["custom"] / "custom_pptx_template_16_9.pptx")

        with_headings = build([{"type": "two_column", "title": "t",
                                "left": {"heading": "L", "body": "- a"},
                                "right": {"heading": "R", "body": "- b"}}])
        without = build([{"type": "two_column", "title": "t",
                          "left": {"body": "- a"}, "right": {"body": "- b"}}])

        assert layout_names_of(with_headings) == ["Porovnání"]
        assert layout_names_of(without) == ["Dva obsahy"]


# =============================================================================
# Registry
# =============================================================================

class TestRegistry:

    def test_no_registry_falls_back_to_the_filename_slots(self, registry):
        """An existing deployment with no YAML must behave exactly as before."""
        specs = load_specs()
        assert [spec.name for spec in specs] == ["16_9", "4_3"]
        assert specs[0].aspect == "16:9"

    def test_custom_file_beats_the_shipped_default(self, registry):
        make_template(registry["custom"] / "custom_pptx_template_16_9.pptx")
        spec, _ = select_template(None, "16:9")
        assert spec.path.parent == registry["custom"]

    def test_registered_templates_are_listed(self, registry):
        make_template(registry["custom"] / "brand.pptx")
        write_registry(registry["config"], [
            {"name": "brand", "pptx_path": "brand.pptx", "description": "Brand deck"},
        ])
        specs = load_specs()
        assert [spec.name for spec in specs] == ["brand"]
        assert specs[0].is_default is True   # the only entry becomes the default
        assert specs[0].description == "Brand deck"

    def test_unknown_template_name_warns_and_uses_the_default(self, registry):
        make_template(registry["custom"] / "brand.pptx")
        write_registry(registry["config"], [{"name": "brand", "pptx_path": "brand.pptx"}])

        pres = build([{"type": "title", "title": "T"}], template="missing")

        assert any("missing" in w and "brand" in w for w in pres.warnings)
        assert len(pres.presentation.slides) == 1

    def test_entry_with_a_path_is_rejected(self, registry, caplog):
        write_registry(registry["config"], [
            {"name": "escape", "pptx_path": "../../etc/passwd"},
        ])
        specs = load_specs()
        # Falls back to the filename slots rather than loading a traversal path.
        assert all(spec.name != "escape" for spec in specs)

    def test_entry_with_a_missing_file_is_skipped(self, registry):
        write_registry(registry["config"], [{"name": "ghost", "pptx_path": "ghost.pptx"}])
        assert all(spec.name != "ghost" for spec in load_specs())

    def test_spec_dir_overrides_the_master_entry(self, registry):
        make_template(registry["custom"] / "brand.pptx")
        make_template(registry["custom"] / "override.pptx")
        write_registry(registry["config"], [{"name": "brand", "pptx_path": "brand.pptx"}])
        spec_dir = registry["config"] / "pptx_templates.d"
        spec_dir.mkdir()
        (spec_dir / "brand.yaml").write_text(
            yaml.safe_dump({"name": "brand", "pptx_path": "override.pptx"}), encoding="utf-8"
        )

        specs = load_specs()
        assert [spec.name for spec in specs] == ["brand"]
        assert specs[0].path.name == "override.pptx"

    def test_defaults_are_applied_when_the_call_omits_them(self, registry):
        make_template(registry["custom"] / "brand.pptx")
        write_registry(registry["config"], [{
            "name": "brand", "pptx_path": "brand.pptx", "default": True,
            "defaults": {"footer_text": "ACME Confidential", "language": "cs-CZ"},
        }])

        pres = build([{"type": "content", "title": "C", "body": "- a"}], template="brand")

        xml = pres.presentation.slides[0].shapes[1]._element.xml
        assert 'lang="cs-CZ"' in xml
        footers = [
            shape.text_frame.text
            for shape in pres.presentation.slides[0].shapes
            if shape.has_text_frame
        ]
        assert any("ACME Confidential" in text for text in footers)

    def test_an_explicit_argument_beats_the_template_default(self, registry):
        make_template(registry["custom"] / "brand.pptx")
        write_registry(registry["config"], [{
            "name": "brand", "pptx_path": "brand.pptx", "default": True,
            "defaults": {"footer_text": "From template"},
        }])

        pres = build([{"type": "content", "title": "C"}],
                     template="brand", footer_text="From the call")

        texts = [
            shape.text_frame.text
            for shape in pres.presentation.slides[0].shapes
            if shape.has_text_frame
        ]
        assert any("From the call" in text for text in texts)
        assert not any("From template" in text for text in texts)


# =============================================================================
# Hot reload
# =============================================================================

class TestHotReload:

    def test_a_new_template_is_picked_up_without_a_restart(self, registry):
        assert [spec.name for spec in load_specs()] == ["16_9", "4_3"]

        make_template(registry["custom"] / "brand.pptx")
        write_registry(registry["config"], [{"name": "brand", "pptx_path": "brand.pptx"}])

        assert [spec.name for spec in load_specs()] == ["brand"]

    def test_a_template_replaced_in_place_is_reread(self, registry):
        """Overwriting brand.pptx must invalidate, not just adding a file.

        A directory's mtime changes when an entry is added or removed, but not
        when an existing file's contents are replaced — which is how a template
        actually gets updated. Without the per-file mtimes in the fingerprint,
        the cached aspect ratio stayed at the old value and selection by aspect
        then picked the wrong template.
        """
        shutil.copyfile(BASE_16_9, registry["custom"] / "brand.pptx")
        write_registry(registry["config"], [{"name": "brand", "pptx_path": "brand.pptx"}])
        assert load_specs()[0].aspect == "16:9"

        time.sleep(0.01)
        shutil.copyfile(BASE_4_3, registry["custom"] / "brand.pptx")

        assert load_specs()[0].aspect == "4:3"

    def test_an_edited_registry_is_reread(self, registry):
        make_template(registry["custom"] / "brand.pptx")
        make_template(registry["custom"] / "second.pptx")
        write_registry(registry["config"], [{"name": "brand", "pptx_path": "brand.pptx"}])
        assert [spec.name for spec in load_specs()] == ["brand"]

        write_registry(registry["config"], [
            {"name": "brand", "pptx_path": "brand.pptx"},
            {"name": "second", "pptx_path": "second.pptx"},
        ])
        assert [spec.name for spec in load_specs()] == ["brand", "second"]


# =============================================================================
# File formats
# =============================================================================

class TestFileFormats:

    def test_potx_opens(self, tmp_path):
        """python-pptx refuses a .potx outright; the content type is rewritten."""
        potx = tmp_path / "brand.potx"
        with zipfile.ZipFile(BASE_16_9) as source:
            with zipfile.ZipFile(potx, "w", zipfile.ZIP_DEFLATED) as target:
                for item in source.infolist():
                    data = source.read(item.filename)
                    if item.filename == "[Content_Types].xml":
                        data = data.decode().replace(
                            "presentationml.presentation.main+xml",
                            "presentationml.template.main+xml",
                        ).encode()
                    target.writestr(item, data)

        presentation = open_template(potx)
        assert len(presentation.slide_layouts) == 11

    def test_potx_builds_a_deck(self, tmp_path, registry):
        potx = registry["custom"] / "brand.potx"
        with zipfile.ZipFile(BASE_16_9) as source:
            with zipfile.ZipFile(potx, "w", zipfile.ZIP_DEFLATED) as target:
                for item in source.infolist():
                    data = source.read(item.filename)
                    if item.filename == "[Content_Types].xml":
                        data = data.decode().replace(
                            "presentationml.presentation.main+xml",
                            "presentationml.template.main+xml",
                        ).encode()
                    target.writestr(item, data)
        write_registry(registry["config"], [{"name": "brand", "pptx_path": "brand.potx"}])

        pres = build(DECK, template="brand")
        assert layout_names_of(pres)[0] == "Úvodní snímek"

    def test_a_file_that_is_not_a_presentation_degrades_to_the_builtin_theme(self, registry):
        broken = registry["custom"] / "broken.pptx"
        broken.write_bytes(b"not a zip at all")
        write_registry(registry["config"], [{"name": "broken", "pptx_path": "broken.pptx"}])

        pres = build([{"type": "title", "title": "T"}], template="broken")

        assert len(pres.presentation.slides) == 1
        assert any("could not be opened" in w for w in pres.warnings)

    def test_aspect_is_read_from_the_file_not_declared(self):
        assert aspect_of(PptxReader(str(BASE_16_9))) == "16:9"
        assert aspect_of(PptxReader(str(BASE_4_3))) == "4:3"


# =============================================================================
# Startup validation
# =============================================================================

class TestValidation:

    def test_report_covers_each_template(self, registry):
        reports = validate_templates()
        assert {report["name"] for report in reports} == {"16_9", "4_3"}
        for report in reports:
            assert report["missing_roles"] == []
            assert report["coverage"]["title"]

    def test_missing_roles_are_reported(self, registry):
        make_template(registry["custom"] / "brand.pptx", trim_layouts)
        write_registry(registry["config"], [{"name": "brand", "pptx_path": "brand.pptx"}])

        report = validate_templates()[0]

        assert "two_column" in report["missing_roles"]
        assert "comparison" in report["missing_roles"]

    def test_unknown_configured_layout_is_reported(self, registry):
        make_template(registry["custom"] / "brand.pptx")
        write_registry(registry["config"], [{
            "name": "brand", "pptx_path": "brand.pptx",
            "layouts": {"content": "Nope"},
        }])

        report = validate_templates()[0]

        assert report["unknown_configured_layouts"] == ["content->Nope"]

    def test_broken_template_reports_an_error_rather_than_raising(self, registry):
        broken = registry["custom"] / "broken.pptx"
        broken.write_bytes(b"nonsense")
        write_registry(registry["config"], [{"name": "broken", "pptx_path": "broken.pptx"}])

        report = validate_templates()[0]

        assert "error" in report
