"""Regression tests for PowerPoint input handling (review Phase 0).

Each test here pins down a failure that reached a user: a crash on input a
model plausibly sends, or — worse — a deck that came back wrong with a success
response. Unlike test_pptx_creation.py these assert on the built presentation,
not merely that a file appeared on disk.
"""

import io
import re
import sys
import zipfile
from pathlib import Path

# Add project root to path for imports
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

import pytest
from pptx import Presentation as PptxReader
from pptx.dml.color import RGBColor

from pptx_tools import slide_builder
from pptx_tools.chart_utils import CHART_TYPE_MAP, ChartDataError, validate_chart_data
from pptx_tools.constants import DEFAULT_SLIDE_FORMAT, TABLE_HEADER_FILL
from pptx_tools.helpers import cell_to_text, coerce_indent_level, parse_color, parse_table_data
from pptx_tools.inline_formatting import needs_inline_processing
from pptx_tools.slide_builder import PowerpointPresentation

_SLIDE_PART_RE = re.compile(r"^ppt/slides/slide\d+\.xml$")


def build(slides, fmt="16:9", **kwargs) -> PowerpointPresentation:
    """Build a presentation in memory."""
    return PowerpointPresentation(slides, fmt, **kwargs)


def reload_presentation(pres: PowerpointPresentation):
    """Round-trip through the saved bytes, as PowerPoint would see it."""
    buffer = pres.save()
    return PptxReader(buffer)


def first_table(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            return shape.table
    raise AssertionError("No table on slide")


def body_paragraphs(slide, exclude_title=True):
    """Paragraphs of the first non-title text frame carrying text."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if exclude_title and shape.is_placeholder and shape.placeholder_format.idx == 0:
            continue
        if shape.text_frame.text.strip():
            return shape.text_frame.paragraphs
    raise AssertionError("No body text frame on slide")


# =============================================================================
# Table cell coercion
# =============================================================================

class TestTableCellCoercion:
    """Numbers and nulls in table_data used to fail the whole presentation."""

    def test_numeric_cells_are_rendered(self):
        pres = build([{
            "slide_type": "table",
            "slide_title": "Numbers",
            "table_data": [["Item", "Qty", "Price"], ["Widget", 12, 9.99]],
        }])
        table = first_table(reload_presentation(pres).slides[0])
        assert table.cell(1, 1).text == "12"
        assert table.cell(1, 2).text == "9.99"

    def test_zero_is_not_blanked(self):
        """`if cell_text` treated a numeric 0 as empty and dropped it."""
        pres = build([{
            "slide_type": "table",
            "slide_title": "Zero",
            "table_data": [["Region", "Sales"], ["North", 0], ["South", "0"]],
        }])
        table = first_table(reload_presentation(pres).slides[0])
        assert table.cell(1, 1).text == "0"
        assert table.cell(2, 1).text == "0"

    def test_none_cell_becomes_empty_string(self):
        pres = build([{
            "slide_type": "table",
            "slide_title": "Missing",
            "table_data": [["A", "B"], [None, "y"]],
        }])
        table = first_table(reload_presentation(pres).slides[0])
        assert table.cell(1, 0).text == ""
        assert table.cell(1, 1).text == "y"

    def test_separator_row_still_sets_alignment_with_mixed_cells(self):
        rows, alignments = parse_table_data([["A", "B"], [":---", "---:"], [1, None]])
        assert rows == [["A", "B"], ["1", ""]]
        assert alignments is not None and len(alignments) == 2

    def test_cell_to_text(self):
        assert cell_to_text(None) == ""
        assert cell_to_text(0) == "0"
        assert cell_to_text(False) == "False"
        assert cell_to_text("x") == "x"


# =============================================================================
# Unknown slide types
# =============================================================================

class TestUnknownSlideType:
    """An unknown type used to be dropped while the caller was told it built."""

    def test_unknown_type_raises(self):
        with pytest.raises(ValueError) as excinfo:
            build([
                {"slide_type": "title", "slide_title": "Deck"},
                {"slide_type": "bullets", "slide_title": "Oops"},
            ])
        message = str(excinfo.value)
        assert "bullets" in message
        assert "slide 1" in message
        # The message must list what the model should have used instead.
        assert "content" in message and "two_column" in message

    def test_no_slides_are_silently_dropped(self):
        slides = [
            {"slide_type": "title", "slide_title": "Deck"},
            {"slide_type": "section", "slide_title": "Part 1"},
            {"slide_type": "content", "slide_title": "Body",
             "slide_text": [{"text": "a", "indentation_level": 1}]},
        ]
        pres = build(slides)
        assert len(reload_presentation(pres).slides) == len(slides)

    def test_text_slide_is_read_not_rejected(self):
        """A client that can only send strings still gets its deck."""
        pres = build(["## A heading\n- a bullet"])
        assert pres.slides[0].type == "content"
        assert pres.slides[0].title == "A heading"

    def test_non_dict_slide_raises(self):
        with pytest.raises(ValueError, match="valid dictionary or object"):
            build([42])


# =============================================================================
# Colours
# =============================================================================

class TestHeaderColor:

    def test_hash_prefixed_hex_is_accepted(self):
        assert parse_color("#C00000", TABLE_HEADER_FILL) == RGBColor.from_string("C00000")

    def test_bare_hex_still_accepted(self):
        assert parse_color("c00000", TABLE_HEADER_FILL) == RGBColor.from_string("C00000")

    def test_unparseable_falls_back_with_warning(self, caplog):
        assert parse_color("cornflower", TABLE_HEADER_FILL) == TABLE_HEADER_FILL
        assert "Unrecognised colour" in caplog.text

    def test_hash_prefixed_header_reaches_the_table(self):
        pres = build([{
            "slide_type": "table",
            "slide_title": "Branded",
            "table_data": [["A", "B"], ["1", "2"]],
            "header_color": "#C00000",
        }])
        table = first_table(reload_presentation(pres).slides[0])
        fill_xml = table.cell(0, 0)._tc.get_or_add_tcPr().xml
        assert "C00000" in fill_xml


# =============================================================================
# Inline formatting
# =============================================================================

class TestInlineFormattingPrecision:

    def test_escape_without_other_markers_is_honoured(self):
        """`price \\* qty` used to keep its backslash: the escape pass was skipped."""
        pres = build([{
            "slide_type": "content",
            "slide_title": "Escapes",
            "slide_text": [{"text": r"price \* qty", "indentation_level": 1}],
        }])
        paragraphs = body_paragraphs(reload_presentation(pres).slides[0])
        assert paragraphs[0].text == "price * qty"
        assert not any(run.font.italic for run in paragraphs[0].runs)

    def test_arithmetic_is_not_italicised(self):
        pres = build([{
            "slide_type": "content",
            "slide_title": "Math",
            "slide_text": [{"text": "5 * 3 * 2 = 30", "indentation_level": 1}],
        }])
        paragraphs = body_paragraphs(reload_presentation(pres).slides[0])
        assert paragraphs[0].text == "5 * 3 * 2 = 30"
        assert not any(run.font.italic for run in paragraphs[0].runs)

    def test_spaced_markers_do_not_format(self):
        assert not needs_inline_processing("10 ~ 20 ~ 30")
        assert not needs_inline_processing("a * b")

    def test_hugging_markers_still_format(self):
        pres = build([{
            "slide_type": "content",
            "slide_title": "Formatting",
            "slide_text": [{"text": "**bold** and *italic* and `code`", "indentation_level": 1}],
        }])
        runs = body_paragraphs(reload_presentation(pres).slides[0])[0].runs
        assert any(run.font.bold and run.text == "bold" for run in runs)
        assert any(run.font.italic and run.text == "italic" for run in runs)
        assert any(run.font.name == "Courier New" and run.text == "code" for run in runs)

    def test_italic_span_ending_in_bold(self):
        """An italic whose last word is bold, with the closers touching.

        The flanking rules first shipped here closed the italic branch on a
        lookbehind, which a nested bold unit fails because it ends in '*'. The
        outer italic was dropped and its asterisks rendered literally on the
        slide. Reported in review of #101.
        """
        pres = build([{
            "slide_type": "content",
            "slide_title": "Nested",
            "slide_text": [{"text": "*Remember, always **backup your data***", "indentation_level": 1}],
        }])
        paragraph = body_paragraphs(reload_presentation(pres).slides[0])[0]

        assert "*" not in paragraph.text, f"stray asterisks rendered: {paragraph.text!r}"
        assert paragraph.text == "Remember, always backup your data"
        # Every run italic; the trailing phrase additionally bold.
        assert all(run.font.italic for run in paragraph.runs)
        assert any(run.font.bold and run.text == "backup your data" for run in paragraph.runs)

    def test_italic_span_with_bold_not_touching_the_close(self):
        """The neighbouring case that already worked must keep working."""
        pres = build([{
            "slide_type": "content",
            "slide_title": "Nested",
            "slide_text": [{"text": "*italic with **bold** inside*", "indentation_level": 1}],
        }])
        paragraph = body_paragraphs(reload_presentation(pres).slides[0])[0]

        assert paragraph.text == "italic with bold inside"
        assert all(run.font.italic for run in paragraph.runs)
        assert any(run.font.bold and run.text == "bold" for run in paragraph.runs)

    def test_backticks_protect_a_dunder(self):
        """`__init__` in backticks renders verbatim rather than underlined."""
        pres = build([{
            "slide_type": "content",
            "slide_title": "Code",
            "slide_text": [{"text": "call `__init__` first", "indentation_level": 1}],
        }])
        runs = body_paragraphs(reload_presentation(pres).slides[0])[0].runs
        dunder = [run for run in runs if run.text == "__init__"]
        assert dunder, "expected a run carrying the literal dunder"
        assert not dunder[0].font.underline


# =============================================================================
# Indentation levels
# =============================================================================

class TestIndentationLevel:

    @pytest.mark.parametrize("value,expected", [
        (1, 0), (2, 1), ("2", 1), (None, 0),
        (7, 4),      # clamped to MAX_INDENT_LEVEL
        (0, 0), (-3, 0),
        ("two", 0),  # unparseable -> top level
    ])
    def test_coercion(self, value, expected):
        assert coerce_indent_level(value) == expected

    def test_string_level_does_not_fail_the_deck(self):
        pres = build([{
            "slide_type": "content",
            "slide_title": "Levels",
            "slide_text": [
                {"text": "top", "indentation_level": "1"},
                {"text": "nested", "indentation_level": "2"},
                {"text": "very deep", "indentation_level": 9},
                {"text": "no level given"},
            ],
        }])
        paragraphs = body_paragraphs(reload_presentation(pres).slides[0])
        assert [p.level for p in paragraphs[:4]] == [0, 1, 4, 0]

    def test_bare_string_bullet_is_accepted(self):
        pres = build([{
            "slide_type": "content",
            "slide_title": "Plain list",
            "slide_text": ["first", "second"],
        }])
        paragraphs = body_paragraphs(reload_presentation(pres).slides[0])
        assert [p.text for p in paragraphs[:2]] == ["first", "second"]


# =============================================================================
# Charts
# =============================================================================

class TestChartTypes:

    def test_scatter_as_a_category_chart_type_points_at_the_slide_type(self):
        """Scatter is a slide type now, so asking for it here names the fix."""
        with pytest.raises(ChartDataError) as excinfo:
            validate_chart_data(
                {"categories": ["a"], "series": [{"name": "s", "values": [1]}]}, "scatter"
            )
        message = str(excinfo.value)
        assert "own slide type" in message
        assert '"type": "scatter"' in message

    def test_scatter_is_not_advertised(self):
        assert "scatter" not in CHART_TYPE_MAP

    def test_supported_chart_still_builds(self):
        pres = build([{
            "slide_type": "chart",
            "slide_title": "Revenue",
            "chart_type": "column",
            "chart_data": {"categories": ["Q1", "Q2"], "series": [{"name": "2026", "values": [1, 2]}]},
        }])
        slide = reload_presentation(pres).slides[0]
        assert any(shape.has_chart for shape in slide.shapes)


# =============================================================================
# Template handling
# =============================================================================

@pytest.fixture
def template_with_sample_slide(tmp_path):
    """A template carrying two sample slides, as a designer would hand over."""
    source = PptxReader(str(project_root / "default_templates" / "default_pptx_template_16_9.pptx"))
    for index, text in ((1, "SAMPLE SLIDE ONE"), (2, "SAMPLE SLIDE TWO")):
        slide = source.slides.add_slide(source.slide_layouts[index])
        slide.shapes.title.text = text

    path = tmp_path / "custom_pptx_template_16_9.pptx"
    source.save(str(path))
    return path


@pytest.fixture
def use_template(monkeypatch):
    """Point the builder at a specific template file.

    Phase 2 replaced the builder's process-lifetime path cache with the
    registry's mtime-checked one, so this now redirects the template search
    directories and drops the registry cache.
    """
    from pptx_tools import templates as templates_mod
    import template_utils

    def _use(path):
        path = Path(path)
        monkeypatch.setattr(
            template_utils, "_candidate_dirs",
            lambda: [path.parent, project_root / "default_templates"],
        )
        templates_mod.clear_cache()

    templates_mod.clear_cache()
    yield _use
    templates_mod.clear_cache()


class TestTemplateSlideRemoval:

    def test_sample_slides_do_not_ship_in_the_output(self, template_with_sample_slide, use_template):
        use_template(template_with_sample_slide)

        pres = build([{"slide_type": "title", "slide_title": "Real deck"}])
        buffer = pres.save()

        package = zipfile.ZipFile(buffer)
        slide_parts = [name for name in package.namelist() if _SLIDE_PART_RE.match(name)]
        assert len(slide_parts) == 1, f"orphan slide parts left in package: {slide_parts}"

        blob = b"".join(package.read(name) for name in slide_parts)
        assert b"SAMPLE SLIDE" not in blob

    def test_slide_count_matches_the_request(self, template_with_sample_slide, use_template):
        use_template(template_with_sample_slide)

        pres = build([
            {"slide_type": "title", "slide_title": "Real deck"},
            {"slide_type": "section", "slide_title": "Part 1"},
        ])
        assert len(reload_presentation(pres).slides) == 2


# =============================================================================
# Format defaults
# =============================================================================

class TestFormatDefault:

    def test_entry_points_share_one_default(self):
        """main.py defaulted to 16:9 while the buffer helper defaulted to 4:3."""
        import inspect

        from pptx_tools.base_pptx_tool import _create_presentation_buffer, create_presentation

        for func in (_create_presentation_buffer, create_presentation):
            assert inspect.signature(func).parameters["format"].default == DEFAULT_SLIDE_FORMAT

    def test_unknown_format_falls_back_to_default(self, caplog):
        pres = build([{"slide_type": "title", "slide_title": "T"}], fmt="widescreen")
        default_pres = build([{"slide_type": "title", "slide_title": "T"}], fmt=DEFAULT_SLIDE_FORMAT)
        assert pres.presentation.slide_width == default_pres.presentation.slide_width
        assert "Unknown presentation format" in caplog.text


def test_buffer_is_a_readable_pptx():
    """The saved buffer opens as a presentation and is positioned at the start."""
    pres = build([{"slide_type": "title", "slide_title": "Round trip"}])
    buffer = pres.save()
    assert isinstance(buffer, io.BytesIO)
    assert buffer.tell() == 0
    assert len(PptxReader(buffer).slides) == 1
