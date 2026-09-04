"""PowerPoint outline-pane sections, derived from ``section`` slides.

python-pptx has no API for sections, so the builder writes the p14 extension
itself. These tests read the XML back. They cannot prove PowerPoint accepts
the file — nothing here can open it — so they check the invariants PowerPoint
is known to enforce: every slide in exactly one section, well-formed names,
GUID ids, and the existing extensions untouched.
"""
import io
import re
import sys
import zipfile
from pathlib import Path

project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

from pptx import Presentation as PptxReader
from pptx.util import Inches

from pptx_tools.slide_builder import PowerpointPresentation
from pptx_tools.templates import TemplateSpec, aspect_of, open_template

TEMPLATE_16_9 = project_root / "default_templates" / "default_pptx_template_16_9.pptx"

SECTION_RE = re.compile(
    r'<p14:section name="([^"]*)" id="(\{[0-9A-F-]{36}\})"><p14:sldIdLst>(.*?)</p14:sldIdLst>'
)


def presentation_xml(pres) -> str:
    return zipfile.ZipFile(pres.save()).read("ppt/presentation.xml").decode("utf-8")


def sections_of(xml):
    block = re.search(r"<p14:sectionLst.*?</p14:sectionLst>", xml, re.S)
    if not block:
        return None
    return [(m.group(1), m.group(2), re.findall(r'id="(\d+)"', m.group(3)))
            for m in SECTION_RE.finditer(block.group(0))]


def build(slides, **kwargs):
    return PowerpointPresentation(slides, "16:9", **kwargs)


class TestSections:

    def test_slides_are_grouped_under_each_section_slide(self):
        pres = build([
            {"type": "title", "title": "Deck"},
            {"type": "section", "title": "Part one"},
            {"type": "content", "title": "A", "body": "- a"},
            {"type": "section", "title": "Part two"},
            {"type": "content", "title": "B", "body": "- b"},
            {"type": "closing", "title": "Thanks"},
        ])
        names = [(name, len(ids)) for name, _id, ids in sections_of(presentation_xml(pres))]
        assert names == [("Default Section", 1), ("Part one", 2), ("Part two", 3)]

    def test_every_slide_is_in_exactly_one_section(self):
        """PowerPoint repairs a file whose section list does not cover every slide."""
        pres = build([
            {"type": "title", "title": "Deck"},
            {"type": "section", "title": "S"},
            {"type": "content", "title": "A", "body": "- a"},
        ])
        xml = presentation_xml(pres)
        slide_ids = re.findall(r'<p:sldId id="(\d+)"', xml)
        section_ids = [sid for _n, _i, ids in sections_of(xml) for sid in ids]
        assert sorted(section_ids) == sorted(slide_ids)
        assert len(section_ids) == len(set(section_ids))

    def test_no_section_list_without_section_slides(self):
        pres = build([{"type": "content", "title": "only", "body": "- x"}])
        assert sections_of(presentation_xml(pres)) is None

    def test_names_are_escaped_and_ids_are_guids(self):
        pres = build([
            {"type": "section", "title": 'Where <we> "stand" & why'},
            {"type": "content", "title": "A", "body": "- a"},
        ])
        (name, guid, _ids), = sections_of(presentation_xml(pres))
        assert name == "Where &lt;we&gt; &quot;stand&quot; &amp; why"
        assert re.fullmatch(r"\{[0-9A-F]{8}(-[0-9A-F]{4}){3}-[0-9A-F]{12}\}", guid)

    def test_an_untitled_section_slide_gets_a_numbered_name(self):
        pres = build([
            {"type": "section", "title": ""},
            {"type": "content", "title": "A", "body": "- a"},
            {"type": "section", "title": "Named"},
        ])
        assert [n for n, _i, _s in sections_of(presentation_xml(pres))] == ["Section 1", "Named"]

    def test_the_templates_own_extensions_survive(self):
        """The shipped template carries a p15 guide list in the same extLst."""
        pres = build([{"type": "section", "title": "S"}, {"type": "content", "title": "A"}])
        xml = presentation_xml(pres)
        assert "sldGuideLst" in xml
        assert xml.count("<p:extLst>") == 1

    def test_the_file_still_opens_and_is_unchanged_in_slide_count(self):
        pres = build([{"type": "section", "title": "S"}, {"type": "content", "title": "A"}])
        assert len(PptxReader(pres.save()).slides) == 2

    def test_template_slides_kept_by_strip_slides_off_join_the_leading_section(self, tmp_path):
        """Grouping is computed over the real slide list, not the models."""
        prs = PptxReader(str(TEMPLATE_16_9))
        prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = "COVER"
        path = tmp_path / "with_cover.pptx"
        prs.save(str(path))
        spec = TemplateSpec(name="c", path=path, strip_slides=False,
                            aspect=aspect_of(open_template(path)))

        pres = build([
            {"type": "section", "title": "S"},
            {"type": "content", "title": "A", "body": "- a"},
        ], template_spec=spec)
        xml = presentation_xml(pres)
        groups = sections_of(xml)
        assert [(n, len(ids)) for n, _i, ids in groups] == [("Default Section", 1), ("S", 2)]
        assert sorted(sid for _n, _i, ids in groups for sid in ids) == \
            sorted(re.findall(r'<p:sldId id="(\d+)"', xml))
