"""Tests for the typed slide schema (review Phase 1).

Covers the contract itself — discrimination, validation, the legacy-key shim,
markdown bodies — and the behaviour it unlocked: real scatter charts, chart
options, inline images, autofit and the warnings channel.
"""

import base64
import io
import json
import sys
from pathlib import Path

project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

import pytest
from pptx import Presentation as PptxReader
from pptx.oxml.ns import qn
from pptx.util import Emu

from pptx_tools.helpers import body_to_bullets, estimate_text_fill, fit_table_font_size
from pptx_tools.image_utils import ImageValidationError, decode_data_uri, is_data_uri
from pptx_tools.schema import (
    Bullet, ContentSlide, SLIDE_TYPES, _SLIDES_ADAPTER, coerce_slides,
)
from pptx_tools.slide_builder import PowerpointPresentation

# A 1x1 transparent PNG.
PNG_1PX = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)
PNG_DATA_URI = "data:image/png;base64," + base64.b64encode(PNG_1PX).decode()


def build(slides, fmt="16:9", **kwargs):
    return PowerpointPresentation(slides, fmt, **kwargs)


def reload_presentation(pres):
    return PptxReader(pres.save())


def body_paragraphs(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            continue
        if shape.text_frame.text.strip():
            return shape.text_frame.paragraphs
    raise AssertionError("No body text frame on slide")


# =============================================================================
# The contract
# =============================================================================

class TestSchemaShape:

    def test_json_schema_is_discriminated(self):
        schema = _SLIDES_ADAPTER.json_schema()
        items = schema.get("items", schema)
        assert "oneOf" in items
        assert "discriminator" in items
        # One definition per slide type, so the model sees each type's fields.
        assert len(schema["$defs"]) >= len(SLIDE_TYPES)

    def test_every_slide_type_is_reachable(self):
        assert set(SLIDE_TYPES) == {
            "title", "section", "content", "table",
            "chart", "scatter", "image", "two_column", "quote",
        }

    def test_unknown_type_names_the_valid_ones(self):
        with pytest.raises(ValueError) as excinfo:
            coerce_slides([{"type": "bullets", "title": "x"}])
        message = str(excinfo.value)
        assert "bullets" in message
        assert "Valid slide types" in message
        assert "two_column" in message

    def test_misspelled_field_is_rejected_not_ignored(self):
        """The whole point of extra='forbid': a typo used to make an empty slide."""
        with pytest.raises(ValueError) as excinfo:
            coerce_slides([{"type": "quote", "txt": "hello"}])
        assert "txt" in str(excinfo.value)

    def test_error_names_the_slide_and_field(self):
        with pytest.raises(ValueError) as excinfo:
            coerce_slides([
                {"type": "section", "title": "ok"},
                {"type": "chart", "title": "c", "chart_type": "pie",
                 "categories": ["a"], "series": [{"name": "s", "values": ["nope"]}]},
            ])
        message = str(excinfo.value)
        assert "slide 1" in message
        assert "series.0.values.0" in message

    def test_bad_colour_is_rejected_with_guidance(self):
        with pytest.raises(ValueError) as excinfo:
            coerce_slides([{"type": "table", "rows": [["a"]], "header_color": "cornflower"}])
        assert "6-digit hex" in str(excinfo.value)

    @pytest.mark.parametrize("value", ["#C00000", "c00000", "accent1", "dark2"])
    def test_accepted_colours(self, value):
        slides = coerce_slides([{"type": "table", "rows": [["a"]], "header_color": value}])
        assert slides[0].header_color == value

    def test_scatter_points_must_be_pairs(self):
        with pytest.raises(ValueError) as excinfo:
            coerce_slides([{"type": "scatter", "series": [{"name": "s", "points": [[1, 2, 3]]}]}])
        assert "exactly [x, y]" in str(excinfo.value)


# =============================================================================
# Legacy compatibility
# =============================================================================

class TestLegacyKeys:
    """Prompts written against the previous key names must keep working."""

    def test_legacy_deck_still_builds(self):
        legacy = [
            {"slide_type": "title", "slide_title": "Deck", "subtitle": "sub"},
            {"slide_type": "section", "slide_title": "Part 1"},
            {"slide_type": "content", "slide_title": "Body",
             "slide_text": [{"text": "a", "indentation_level": 1},
                            {"text": "b", "indentation_level": 2}],
             "speaker_notes": "notes"},
            {"slide_type": "table", "slide_title": "T",
             "table_data": [["A", "B"], ["1", "2"]], "alternate_rows": False},
            {"slide_type": "two_column", "slide_title": "2c",
             "left_heading": "L", "left_column": [{"text": "l", "indentation_level": 1}],
             "right_heading": "R", "right_column": [{"text": "r", "indentation_level": 1}]},
            {"slide_type": "chart", "slide_title": "Ch", "chart_type": "column",
             "chart_data": {"categories": ["a"], "series": [{"name": "s", "values": [1]}]},
             "has_legend": False},
            {"slide_type": "quote", "quote_text": "q", "quote_author": "me"},
        ]
        pres = build(legacy)
        assert len(reload_presentation(pres).slides) == len(legacy)

    def test_legacy_keys_map_onto_new_fields(self):
        slides = coerce_slides([{
            "slide_type": "content", "slide_title": "T",
            "slide_text": [{"text": "x", "indentation_level": 3}],
            "speaker_notes": "n",
        }])
        slide = slides[0]
        assert slide.type == "content"
        assert slide.title == "T"
        assert slide.notes == "n"
        assert slide.body[0].level == 3

    def test_legacy_two_column_becomes_nested_columns(self):
        slide = coerce_slides([{
            "slide_type": "two_column",
            "left_heading": "L", "left_column": [{"text": "a"}],
            "right_column": [{"text": "b"}],
        }])[0]
        assert slide.left.heading == "L"
        assert slide.left.body[0].text == "a"
        assert slide.right.heading is None
        assert slide.right.body[0].text == "b"

    def test_legacy_has_legend_false_becomes_none(self):
        slide = coerce_slides([{
            "slide_type": "chart", "chart_type": "pie",
            "chart_data": {"categories": ["a"], "series": [{"name": "s", "values": [1]}]},
            "has_legend": False,
        }])[0]
        assert slide.legend == "none"

    def test_out_of_range_level_is_clamped_not_rejected(self):
        """Phase 0 clamped these; the typed schema must not turn it into an error."""
        slide = coerce_slides([{
            "type": "content",
            "body": [{"text": "a", "level": 9}, {"text": "b", "level": "2"},
                     {"text": "c", "level": "deep"}],
        }])[0]
        assert [b.level for b in slide.body] == [5, 2, 1]


# =============================================================================
# Markdown bodies
# =============================================================================

class TestMarkdownBody:

    def test_two_space_indent(self):
        bullets = body_to_bullets("- one\n  - nested\n- two")
        assert [(b.text, b.level) for b in bullets] == [
            ("one", 1), ("nested", 2), ("two", 1)
        ]

    def test_four_space_and_tab_indent_agree(self):
        four = body_to_bullets("- a\n    - b\n        - c")
        tabs = body_to_bullets("- a\n\t- b\n\t\t- c")
        assert [b.level for b in four] == [1, 2, 3]
        assert [b.level for b in tabs] == [1, 2, 3]

    @pytest.mark.parametrize("marker", ["-", "*", "+"])
    def test_all_bullet_markers(self, marker):
        assert body_to_bullets(f"{marker} item")[0].text == "item"

    def test_line_without_marker_is_top_level(self):
        bullets = body_to_bullets("Plain line\n- bullet")
        assert [(b.text, b.level) for b in bullets] == [("Plain line", 1), ("bullet", 1)]

    def test_blank_lines_are_skipped(self):
        assert len(body_to_bullets("- a\n\n\n- b")) == 2

    def test_explicit_bullets_pass_through(self):
        bullets = body_to_bullets([Bullet(text="x", level=2)])
        assert bullets[0].level == 2

    def test_markdown_body_renders(self):
        pres = build([{"type": "content", "title": "C", "body": "- top\n  - child\n- next"}])
        paragraphs = body_paragraphs(reload_presentation(pres).slides[0])
        assert [(p.text, p.level) for p in paragraphs[:3]] == [
            ("top", 0), ("child", 1), ("next", 0)
        ]

    def test_markdown_body_keeps_inline_formatting(self):
        pres = build([{"type": "content", "title": "C", "body": "- Revenue **up 12%**"}])
        runs = body_paragraphs(reload_presentation(pres).slides[0])[0].runs
        assert any(run.font.bold and run.text == "up 12%" for run in runs)

    def test_markdown_body_in_two_column(self):
        pres = build([{
            "type": "two_column", "title": "2c",
            "left": {"heading": "L", "body": "- a\n  - a2"},
            "right": {"body": "- b"},
        }])
        texts = [
            shape.text_frame.text
            for shape in reload_presentation(pres).slides[0].shapes
            if shape.has_text_frame
        ]
        assert any("a2" in text for text in texts)
        assert any("b" in text for text in texts)


# =============================================================================
# Charts
# =============================================================================

class TestCharts:

    def test_scatter_builds(self):
        """The chart type that could never build before."""
        pres = build([{
            "type": "scatter", "title": "XY",
            "series": [{"name": "trial", "points": [[1, 2.5], [2, 3.5], [3, 9]]}],
            "x_title": "dose", "y_title": "response",
        }])
        slide = reload_presentation(pres).slides[0]
        charts = [shape.chart for shape in slide.shapes if shape.has_chart]
        assert len(charts) == 1
        assert len(charts[0].plots[0].series) == 1

    def test_scatter_with_multiple_series(self):
        pres = build([{
            "type": "scatter", "title": "XY",
            "series": [
                {"name": "a", "points": [[1, 1], [2, 2]]},
                {"name": "b", "points": [[1, 3], [2, 4]]},
            ],
        }])
        chart = [s.chart for s in reload_presentation(pres).slides[0].shapes if s.has_chart][0]
        assert len(chart.plots[0].series) == 2

    def test_data_labels_and_number_format(self):
        pres = build([{
            "type": "chart", "title": "C", "chart_type": "column",
            "categories": ["a", "b"], "series": [{"name": "s", "values": [1000, 2000]}],
            "data_labels": True, "number_format": "#,##0",
        }])
        chart = [s.chart for s in reload_presentation(pres).slides[0].shapes if s.has_chart][0]
        plot = chart.plots[0]
        assert plot.has_data_labels
        assert plot.data_labels.number_format == "#,##0"

    def test_legend_none_hides_it(self):
        pres = build([{
            "type": "chart", "title": "C", "chart_type": "pie",
            "categories": ["a"], "series": [{"name": "s", "values": [1]}],
            "legend": "none",
        }])
        chart = [s.chart for s in reload_presentation(pres).slides[0].shapes if s.has_chart][0]
        assert chart.has_legend is False

    def test_axis_titles(self):
        pres = build([{
            "type": "chart", "title": "C", "chart_type": "column",
            "categories": ["a"], "series": [{"name": "s", "values": [1]}],
            "x_title": "Quarter", "y_title": "Revenue",
        }])
        chart = [s.chart for s in reload_presentation(pres).slides[0].shapes if s.has_chart][0]
        assert chart.value_axis.axis_title.text_frame.text == "Revenue"
        assert chart.category_axis.axis_title.text_frame.text == "Quarter"

    def test_axis_title_on_pie_is_ignored_not_fatal(self):
        """A pie has no axes; asking for a title should not fail the deck."""
        pres = build([{
            "type": "chart", "title": "C", "chart_type": "pie",
            "categories": ["a"], "series": [{"name": "s", "values": [1]}],
            "y_title": "Revenue",
        }])
        assert len(reload_presentation(pres).slides) == 1

    def test_chart_title(self):
        pres = build([{
            "type": "chart", "title": "Slide", "chart_type": "column",
            "categories": ["a"], "series": [{"name": "s", "values": [1]}],
            "chart_title": "Inside the chart",
        }])
        chart = [s.chart for s in reload_presentation(pres).slides[0].shapes if s.has_chart][0]
        assert chart.chart_title.text_frame.text == "Inside the chart"

    def test_series_length_mismatch_warns(self):
        pres = build([{
            "type": "chart", "title": "C", "chart_type": "column",
            "categories": ["a", "b", "c"], "series": [{"name": "short", "values": [1]}],
        }])
        assert any("short" in w and "categories" in w for w in pres.warnings)

    def test_null_value_is_allowed(self):
        pres = build([{
            "type": "chart", "title": "C", "chart_type": "line",
            "categories": ["a", "b", "c"],
            "series": [{"name": "s", "values": [1, None, 3]}],
        }])
        assert len(reload_presentation(pres).slides) == 1


# =============================================================================
# Images
# =============================================================================

class TestInlineImages:

    def test_is_data_uri(self):
        assert is_data_uri(PNG_DATA_URI)
        assert not is_data_uri("https://example.com/a.png")

    def test_decode_data_uri(self):
        stream, ext = decode_data_uri(PNG_DATA_URI)
        assert ext == "png"
        assert stream.getvalue() == PNG_1PX

    def test_inline_image_is_placed(self):
        """No network, no public URL: the bytes are already in the payload."""
        pres = build([{"type": "image", "title": "Chart", "source": PNG_DATA_URI}])
        slide = reload_presentation(pres).slides[0]
        assert any(shape.shape_type == 13 for shape in slide.shapes)  # PICTURE
        assert pres.warnings == []

    def test_inline_image_with_caption(self):
        pres = build([{
            "type": "image", "title": "Chart", "source": PNG_DATA_URI, "caption": "Fig 1",
        }])
        texts = [
            shape.text_frame.text
            for shape in reload_presentation(pres).slides[0].shapes
            if shape.has_text_frame
        ]
        assert any("Fig 1" in text for text in texts)

    def test_non_image_data_uri_is_rejected(self):
        with pytest.raises(ImageValidationError, match="Invalid image type"):
            decode_data_uri("data:text/html;base64," + base64.b64encode(b"<h1>x</h1>").decode())

    def test_non_base64_data_uri_is_rejected(self):
        with pytest.raises(ImageValidationError, match="base64"):
            decode_data_uri("data:image/png,notbase64")

    def test_corrupt_base64_is_rejected(self):
        with pytest.raises(ImageValidationError, match="decode"):
            decode_data_uri("data:image/png;base64,!!!!not-base64!!!!")

    def test_image_slide_without_a_source_is_rejected(self):
        """It used to build an empty slide and report success."""
        with pytest.raises(ValueError) as excinfo:
            coerce_slides([{"type": "image", "title": "No picture"}])
        assert "source" in str(excinfo.value)

    def test_failed_image_warns_instead_of_silently_placeholdering(self):
        pres = build([{"type": "image", "title": "X", "source": "data:image/png;base64,!!!"}])
        assert any("image could not be loaded" in w for w in pres.warnings)


# =============================================================================
# Warnings channel
# =============================================================================

class TestWarnings:

    def test_clean_deck_has_no_warnings(self):
        pres = build([
            {"type": "title", "title": "T"},
            {"type": "content", "title": "C", "body": "- short"},
        ])
        assert pres.warnings == []

    def test_layout_override_is_reported_as_ignored(self):
        """Accepted by the schema, not yet honoured — say so rather than drop it."""
        pres = build([{"type": "content", "title": "C", "body": "- a", "layout": "Brand Body"}])
        assert any("Brand Body" in w and "ignored" in w for w in pres.warnings)

    def test_overfull_body_warns(self):
        long_bullets = [{"text": "A fairly long bullet line that will wrap once or twice " * 3}
                        for _ in range(24)]
        pres = build([{"type": "content", "title": "Too much", "body": long_bullets}])
        assert any("shrunk to fit" in w for w in pres.warnings)

    def test_tall_table_warns_and_shrinks(self):
        rows = [["Col A", "Col B"]] + [[f"row {i}", str(i)] for i in range(40)]
        pres = build([{"type": "table", "title": "Big", "rows": rows}])
        assert any("table" in w for w in pres.warnings)

    def test_empty_table_warns(self):
        pres = build([{"type": "table", "title": "Empty", "rows": []}])
        assert any("no rows" in w for w in pres.warnings)

    def test_warnings_name_the_slide_index(self):
        pres = build([
            {"type": "title", "title": "T"},
            {"type": "image", "title": "X", "source": "data:image/png;base64,!!!"},
        ])
        assert any(w.startswith("slide 1:") for w in pres.warnings)


# =============================================================================
# Autofit and table sizing
# =============================================================================

class TestFit:

    def test_estimate_grows_with_text(self):
        small = estimate_text_fill([Bullet(text="short")], Emu(6000000), Emu(3000000))
        large = estimate_text_fill([Bullet(text="x" * 4000)], Emu(6000000), Emu(3000000))
        assert large > small
        assert small < 1.0 < large

    def test_autofit_element_is_written(self):
        pres = build([{
            "type": "content", "title": "C",
            "body": [{"text": "a line that wraps " * 30} for _ in range(20)],
        }])
        slide = reload_presentation(pres).slides[0]
        body = [s for s in slide.shapes if s.is_placeholder and s.placeholder_format.idx == 1][0]
        bodyPr = body.text_frame._txBody.find(qn('a:bodyPr'))
        autofit = bodyPr.find(qn('a:normAutofit'))
        assert autofit is not None
        assert int(autofit.get('fontScale')) <= 100000

    def test_table_font_shrinks_with_row_count(self):
        tall = fit_table_font_size(40, Emu(4000000))
        short = fit_table_font_size(4, Emu(4000000))
        assert tall < short

    def test_table_font_never_below_floor(self):
        assert fit_table_font_size(500, Emu(4000000)) >= 9

    def test_explicit_table_font_size_is_honoured(self):
        rows = [["A", "B"]] + [[str(i), "x"] for i in range(30)]
        pres = build([{"type": "table", "title": "T", "rows": rows, "font_size": 8}])
        table = [s.table for s in reload_presentation(pres).slides[0].shapes if s.has_table][0]
        assert table.cell(1, 0).text_frame.paragraphs[0].font.size.pt == 8


# =============================================================================
# Language
# =============================================================================

class TestLanguage:

    def test_language_is_stamped_on_runs(self):
        pres = build(
            [{"type": "content", "title": "Nadpis", "body": "- Text s **tučným** slovem"}],
            language="cs-CZ",
        )
        xml = reload_presentation(pres).slides[0].shapes[1]._element.xml
        assert 'lang="cs-CZ"' in xml

    def test_language_reaches_table_cells(self):
        pres = build(
            [{"type": "table", "title": "T", "rows": [["Sloupec"], ["Hodnota"]]}],
            language="cs-CZ",
        )
        table = [s for s in reload_presentation(pres).slides[0].shapes if s.has_table][0]
        assert 'lang="cs-CZ"' in table._element.xml

    def test_no_language_leaves_runs_untouched(self):
        pres = build([{"type": "content", "title": "T", "body": "- a"}])
        xml = reload_presentation(pres).slides[0].shapes[1]._element.xml
        assert 'lang=' not in xml


# =============================================================================
# Buffer contract
# =============================================================================

def test_buffer_function_returns_warnings():
    from pptx_tools import _create_presentation_buffer

    buffer, warnings = _create_presentation_buffer(
        [{"type": "content", "title": "C", "body": "- a", "layout": "Nope"}], "16:9"
    )
    assert isinstance(buffer, io.BytesIO)
    assert any("Nope" in w for w in warnings)
    assert len(PptxReader(buffer).slides) == 1
