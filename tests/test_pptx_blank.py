"""The ``blank`` slide: positioned elements, the escape hatch from #100.

Positions are resolved against the real slide size and validated against it.
Asserting on where shapes actually land, in EMU, and on what the warnings
channel says when they would not fit.
"""
import base64
import sys
from pathlib import Path

project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

import pytest
from pptx import Presentation as PptxReader
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches

from pptx_tools.schema import coerce_slides
from pptx_tools.slide_builder import PowerpointPresentation

PNG_DATA_URI = "data:image/png;base64," + base64.b64encode(base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)).decode()


def build(elements, fmt="16:9", **slide):
    pres = PowerpointPresentation([{"type": "blank", "elements": elements, **slide}], fmt)
    return pres, PptxReader(pres.save()).slides[0]


def shapes_of(slide, kind=None):
    return [s for s in slide.shapes if kind is None or s.shape_type == kind]


class TestSchema:

    def test_needs_at_least_one_element(self):
        with pytest.raises(ValueError):
            coerce_slides([{"type": "blank", "elements": []}])

    @pytest.mark.parametrize("bad", [-1, "abc", "50 percent", True, "1.5cm"])
    def test_rejects_positions_it_cannot_resolve(self, bad):
        with pytest.raises(ValueError):
            coerce_slides([{"type": "blank", "elements": [
                {"kind": "text", "text": "x", "x": bad, "y": 0, "w": 1}]}])

    @pytest.mark.parametrize("ok", [0, 1.5, "2", "1.5in", "40%", " 10 % "])
    def test_accepts_inches_and_percentages(self, ok):
        coerce_slides([{"type": "blank", "elements": [
            {"kind": "text", "text": "x", "x": ok, "y": 0, "w": 1}]}])

    def test_element_kind_is_discriminated(self):
        with pytest.raises(ValueError, match="kind"):
            coerce_slides([{"type": "blank", "elements": [{"kind": "video", "x": 0, "y": 0, "w": 1}]}])


class TestPlacement:

    def test_inches_land_where_asked(self):
        _, slide = build([{"kind": "text", "text": "hi", "x": 1, "y": 2, "w": 3, "h": 0.5}])
        box = shapes_of(slide, MSO_SHAPE_TYPE.TEXT_BOX)[0]
        assert (box.left, box.top, box.width, box.height) == (Inches(1), Inches(2), Inches(3), Inches(0.5))

    def test_percentages_resolve_against_the_real_slide_size(self):
        """The same 50% is a different length on 16:9 and 4:3 — as it should be."""
        wide = build([{"kind": "shape", "x": "0%", "y": "0%", "w": "50%", "h": "50%"}], "16:9")[1]
        narrow = build([{"kind": "shape", "x": "0%", "y": "0%", "w": "50%", "h": "50%"}], "4:3")[1]
        w_shape = shapes_of(wide, MSO_SHAPE_TYPE.AUTO_SHAPE)[0]
        n_shape = shapes_of(narrow, MSO_SHAPE_TYPE.AUTO_SHAPE)[0]
        assert w_shape.width == Inches(13.333 / 2) or abs(w_shape.width - Inches(13.333 / 2)) < Emu(20000)
        assert abs(n_shape.width - Inches(5)) < Emu(20000)
        assert w_shape.height == n_shape.height   # both slides are 7.5in tall

    def test_elements_draw_in_order_so_later_ones_sit_on_top(self):
        _, slide = build([
            {"kind": "shape", "x": 0, "y": 0, "w": 2, "h": 2, "fill": "FF0000"},
            {"kind": "text", "text": "over", "x": 0, "y": 0, "w": 2, "h": 1},
        ])
        kinds = [s.shape_type for s in slide.shapes]
        assert kinds.index(MSO_SHAPE_TYPE.AUTO_SHAPE) < kinds.index(MSO_SHAPE_TYPE.TEXT_BOX)

    def test_a_titled_blank_slide_reports_the_dropped_title(self):
        """The blank layout has no title placeholder; saying so beats losing it."""
        pres, _ = build([{"kind": "text", "text": "x", "x": 0, "y": 0, "w": 1}], title="Lost")
        assert any("title" in w.lower() for w in pres.warnings)


class TestValidationAgainstSlideSize:

    def test_an_element_running_past_the_edge_is_shrunk_and_reported(self):
        pres, slide = build([{"kind": "shape", "x": 10, "y": 1, "w": 8, "h": 1}], "16:9")
        shape = shapes_of(slide, MSO_SHAPE_TYPE.AUTO_SHAPE)[0]
        assert shape.left + shape.width <= PptxReader(pres.save()).slide_width
        assert any("ran past the slide edge" in w and "width" in w for w in pres.warnings)

    def test_an_element_starting_off_the_slide_is_skipped_and_reported(self):
        pres, slide = build([
            {"kind": "shape", "x": 40, "y": 1, "w": 1, "h": 1},
            {"kind": "text", "text": "kept", "x": 1, "y": 1, "w": 1},
        ])
        assert shapes_of(slide, MSO_SHAPE_TYPE.AUTO_SHAPE) == []
        assert len(shapes_of(slide, MSO_SHAPE_TYPE.TEXT_BOX)) == 1
        assert any("starts off the slide; skipped" in w for w in pres.warnings)

    def test_nothing_is_reported_when_everything_fits(self):
        pres, _ = build([{"kind": "shape", "x": "10%", "y": "10%", "w": "80%", "h": "80%"}])
        assert pres.warnings == []


class TestKinds:

    def test_text_takes_inline_markdown_size_and_alignment(self):
        _, slide = build([{"kind": "text", "text": "**bold** and *it*", "x": 0, "y": 0, "w": 4,
                           "font_size": 30, "align": "center"}])
        para = shapes_of(slide, MSO_SHAPE_TYPE.TEXT_BOX)[0].text_frame.paragraphs[0]
        assert [r.text for r in para.runs] == ["bold", " and ", "it"]
        assert para.runs[0].font.bold and para.runs[2].font.italic
        assert para.runs[0].font.size.pt == 30
        assert para.alignment is not None and para.alignment.name == "CENTER"

    def test_image_keeps_its_aspect_within_the_box(self):
        _, slide = build([{"kind": "image", "source": PNG_DATA_URI, "x": 1, "y": 1, "w": 4, "h": 1}])
        pic = shapes_of(slide, MSO_SHAPE_TYPE.PICTURE)[0]
        assert pic.left == Inches(1) and pic.top == Inches(1)
        assert pic.width <= Inches(4) and pic.height <= Inches(1)
        assert pic.width == pic.height        # a square fixture stays square

    def test_a_broken_image_draws_a_placeholder_and_reports(self):
        pres, slide = build([{"kind": "image", "source": "data:image/png;base64,AAAA", "x": 1, "y": 1, "w": 2}])
        assert shapes_of(slide, MSO_SHAPE_TYPE.PICTURE) == []
        assert any("could not be loaded" in w for w in pres.warnings)

    @pytest.mark.parametrize("name", ["rectangle", "rounded_rectangle", "ellipse", "chevron", "arrow"])
    def test_every_shape_name_draws_an_autoshape(self, name):
        _, slide = build([{"kind": "shape", "shape": name, "x": 1, "y": 1, "w": 2, "h": 1}])
        assert len(shapes_of(slide, MSO_SHAPE_TYPE.AUTO_SHAPE)) == 1

    def test_shape_fill_hex_and_theme_name(self):
        _, slide = build([
            {"kind": "shape", "x": 0, "y": 0, "w": 1, "h": 1, "fill": "#FF0000"},
            {"kind": "shape", "x": 2, "y": 0, "w": 1, "h": 1, "fill": "accent2"},
        ])
        hex_shape, theme_shape = shapes_of(slide, MSO_SHAPE_TYPE.AUTO_SHAPE)
        assert str(hex_shape.fill.fore_color.rgb) == "FF0000"
        assert theme_shape.fill.fore_color.theme_color.name == "ACCENT_2"

    def test_shape_text_is_centred_inside_it(self):
        _, slide = build([{"kind": "shape", "x": 0, "y": 0, "w": 3, "h": 1, "text": "Go"}])
        shape = shapes_of(slide, MSO_SHAPE_TYPE.AUTO_SHAPE)[0]
        assert shape.text_frame.text == "Go"
        assert shape.text_frame.paragraphs[0].alignment.name == "CENTER"

    def test_shape_with_no_fill_keeps_the_theme_default(self):
        _, slide = build([{"kind": "shape", "x": 0, "y": 0, "w": 1, "h": 1}])
        shape = shapes_of(slide, MSO_SHAPE_TYPE.AUTO_SHAPE)[0]
        assert shape.fill.type is None or shape.fill.type.name != "SOLID"


class TestIntegration:

    def test_blank_is_a_registered_slide_type(self):
        from pptx_tools.schema import SLIDE_TYPES
        assert "blank" in SLIDE_TYPES

    def test_blank_lands_on_the_blank_layout(self):
        _, slide = build([{"kind": "text", "text": "x", "x": 0, "y": 0, "w": 1}])
        assert slide.slide_layout.name == "Prázdný"
