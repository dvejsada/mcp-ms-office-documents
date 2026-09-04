"""Tests for the slide types and inline features added in review Phase 3.

KPI rows, timelines, agendas and closing slides draw their own shapes rather
than filling a placeholder, so these assert on what actually lands on the
slide — shape counts, text, sizes and colours — not merely that a file was
produced.
"""

import base64
import sys
from pathlib import Path

project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

import pytest
from pptx import Presentation as PptxReader
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches

from pptx_tools.schema import coerce_slides
from pptx_tools.slide_builder import PowerpointPresentation

PNG_DATA_URI = "data:image/png;base64," + base64.b64encode(base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)).decode()


def build(slides, fmt="16:9", **kwargs):
    return PowerpointPresentation(slides, fmt, **kwargs)


def reload_presentation(pres):
    return PptxReader(pres.save())


def all_text(slide):
    return "\n".join(
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    )


# =============================================================================
# KPI
# =============================================================================

class TestKpi:

    def test_one_box_per_figure(self):
        pres = build([{"type": "kpi", "title": "At a glance", "items": [
            {"value": "€4.2M", "label": "ARR", "delta": "+12% vs Q2"},
            {"value": "18%", "label": "Churn"},
        ]}])
        slide = reload_presentation(pres).slides[0]
        text = all_text(slide)

        assert "€4.2M" in text and "ARR" in text and "+12% vs Q2" in text
        assert "18%" in text and "Churn" in text

    def test_value_is_larger_than_its_label(self):
        """The whole point of a KPI row is the typographic contrast."""
        pres = build([{"type": "kpi", "items": [{"value": "42", "label": "Widgets"}]}])
        box = [s for s in reload_presentation(pres).slides[0].shapes
               if s.has_text_frame and "42" in s.text_frame.text][0]
        value, label = box.text_frame.paragraphs[0], box.text_frame.paragraphs[1]
        assert value.font.size > label.font.size
        assert value.font.bold

    def test_boxes_do_not_overlap(self):
        pres = build([{"type": "kpi", "items": [
            {"value": str(n), "label": f"L{n}"} for n in range(4)
        ]}])
        boxes = sorted(
            (s for s in reload_presentation(pres).slides[0].shapes
             if s.has_text_frame and s.text_frame.text.strip().startswith(("0", "1", "2", "3"))),
            key=lambda s: s.left,
        )
        assert len(boxes) == 4
        for earlier, later in zip(boxes, boxes[1:]):
            assert earlier.left + earlier.width <= later.left

    def test_too_many_figures_warns(self):
        pres = build([{"type": "kpi", "items": [
            {"value": str(n), "label": f"L{n}"} for n in range(6)
        ]}])
        assert any("cramped" in w for w in pres.warnings)

    def test_more_than_six_is_rejected(self):
        with pytest.raises(ValueError):
            coerce_slides([{"type": "kpi", "items": [
                {"value": str(n), "label": "x"} for n in range(7)
            ]}])

    def test_at_least_one_is_required(self):
        with pytest.raises(ValueError):
            coerce_slides([{"type": "kpi", "items": []}])


# =============================================================================
# Agenda
# =============================================================================

class TestAgenda:

    def test_derived_from_the_decks_own_sections(self):
        """Omitting items is the usual case; it cannot drift from the deck."""
        pres = build([
            {"type": "title", "title": "Deck"},
            {"type": "agenda"},
            {"type": "section", "title": "Results"},
            {"type": "content", "title": "Detail", "body": "- x"},
            {"type": "section", "title": "Outlook"},
        ])
        text = all_text(reload_presentation(pres).slides[1])
        assert "1.  Results" in text
        assert "2.  Outlook" in text
        assert "Detail" not in text

    def test_explicit_items_win(self):
        pres = build([
            {"type": "agenda", "items": ["First", "Second"]},
            {"type": "section", "title": "Ignored"},
        ])
        text = all_text(reload_presentation(pres).slides[0])
        assert "1.  First" in text and "2.  Second" in text
        assert "Ignored" not in text

    def test_default_title(self):
        pres = build([{"type": "agenda", "items": ["A"]}])
        assert "Agenda" in all_text(reload_presentation(pres).slides[0])

    def test_no_sections_and_no_items_warns(self):
        pres = build([{"type": "agenda"}])
        assert any("no items" in w for w in pres.warnings)


# =============================================================================
# Closing
# =============================================================================

class TestClosing:

    def test_default_title_and_contact_lines(self):
        pres = build([{"type": "closing", "subtitle": "Questions?",
                       "contact": ["dan@example.com", "+420 123"]}])
        text = all_text(reload_presentation(pres).slides[0])
        assert "Thank you" in text
        assert "Questions?" in text
        assert "dan@example.com" in text and "+420 123" in text

    def test_custom_title(self):
        pres = build([{"type": "closing", "title": "Díky!"}])
        assert "Díky!" in all_text(reload_presentation(pres).slides[0])


# =============================================================================
# Timeline
# =============================================================================

class TestTimeline:

    def test_one_shape_per_step(self):
        pres = build([{"type": "timeline", "title": "Roadmap", "steps": [
            {"label": "Discovery"}, {"label": "Build"}, {"label": "Launch"},
        ]}])
        slide = reload_presentation(pres).slides[0]
        autoshapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        assert len(autoshapes) == 3
        assert {s.text_frame.text for s in autoshapes} == {"Discovery", "Build", "Launch"}

    def test_steps_run_left_to_right_without_reordering(self):
        pres = build([{"type": "timeline", "steps": [
            {"label": "A"}, {"label": "B"}, {"label": "C"},
        ]}])
        shapes = [s for s in reload_presentation(pres).slides[0].shapes
                  if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        ordered = sorted(shapes, key=lambda s: s.left)
        assert [s.text_frame.text for s in ordered] == ["A", "B", "C"]

    def test_detail_is_placed_below_the_shape(self):
        """Detail outside the chevron, so a long line cannot burst it."""
        pres = build([{"type": "timeline", "steps": [
            {"label": "Discovery", "detail": "January to February"},
            {"label": "Build", "detail": "March onwards"},
        ]}])
        slide = reload_presentation(pres).slides[0]
        chevron = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE][0]
        detail = [s for s in slide.shapes
                  if s.has_text_frame and "January" in s.text_frame.text][0]
        assert detail.top >= chevron.top + chevron.height
        assert "January" not in chevron.text_frame.text

    def test_box_style_uses_a_different_shape(self):
        chevrons = build([{"type": "timeline", "steps": [{"label": "A"}, {"label": "B"}]}])
        boxes = build([{"type": "timeline", "style": "box",
                        "steps": [{"label": "A"}, {"label": "B"}]}])

        def kind(pres):
            shape = [s for s in reload_presentation(pres).slides[0].shapes
                     if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE][0]
            return shape.auto_shape_type

        assert kind(chevrons) != kind(boxes)

    def test_two_steps_minimum(self):
        with pytest.raises(ValueError):
            coerce_slides([{"type": "timeline", "steps": [{"label": "only"}]}])


# =============================================================================
# Text beside a picture or chart
# =============================================================================

class TestSideBySide:

    def test_image_with_body_is_left_aligned_not_centred(self):
        """Alone the picture centres on the slide; beside text it must not.

        Asserting on position rather than width: this fixture is a square
        image, so the height of the content area is what bounds it and the
        width would be identical either way. Confinement to the left half is
        covered by test_image_body_does_not_overlap_the_picture.
        """
        alone = build([{"type": "image", "title": "Chart", "source": PNG_DATA_URI}])
        split = build([{"type": "image", "title": "Chart", "source": PNG_DATA_URI,
                        "body": "- Up and to the right"}])

        def picture(pres):
            return [s for s in reload_presentation(pres).slides[0].shapes
                    if s.shape_type == MSO_SHAPE_TYPE.PICTURE][0]

        slide_width = alone.presentation.slide_width
        centred = picture(alone)
        assert abs(centred.left - (slide_width - centred.width) // 2) < Inches(0.05)

        assert picture(split).left < centred.left
        assert "Up and to the right" in all_text(reload_presentation(split).slides[0])

    def test_image_body_does_not_overlap_the_picture(self):
        pres = build([{"type": "image", "title": "Chart", "source": PNG_DATA_URI,
                       "body": "- Beside it"}])
        slide = reload_presentation(pres).slides[0]
        pic = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE][0]
        text = [s for s in slide.shapes
                if s.has_text_frame and "Beside it" in s.text_frame.text][0]
        assert pic.left + pic.width <= text.left

    def test_chart_with_body_splits_the_slide(self):
        chart_only = build([{"type": "chart", "title": "R", "chart_type": "column",
                             "categories": ["a"], "series": [{"name": "s", "values": [1]}]}])
        split = build([{"type": "chart", "title": "R", "chart_type": "column",
                        "categories": ["a"], "series": [{"name": "s", "values": [1]}],
                        "body": "- Doubled"}])

        def chart_width(pres):
            return [s for s in reload_presentation(pres).slides[0].shapes if s.has_chart][0].width

        assert chart_width(split) < chart_width(chart_only)
        assert "Doubled" in all_text(reload_presentation(split).slides[0])

    def test_caption_tracks_the_narrowed_picture(self):
        pres = build([{"type": "image", "title": "C", "source": PNG_DATA_URI,
                       "body": "- text", "caption": "Fig 1"}])
        slide = reload_presentation(pres).slides[0]
        pic = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE][0]
        caption = [s for s in slide.shapes
                   if s.has_text_frame and "Fig 1" in s.text_frame.text][0]
        assert caption.left == pic.left
        assert caption.width == pic.width


# =============================================================================
# Hyperlinks and table markdown
# =============================================================================

class TestInlineAdditions:

    def test_hyperlink_in_a_bullet(self):
        pres = build([{"type": "content", "title": "L",
                       "body": "- See [the docs](https://example.com/x) now"}])
        runs = reload_presentation(pres).slides[0].placeholders[1].text_frame.paragraphs[0].runs
        linked = [r for r in runs if r.hyperlink.address]
        assert len(linked) == 1
        assert linked[0].text == "the docs"
        assert linked[0].hyperlink.address == "https://example.com/x"
        assert "".join(r.text for r in runs) == "See the docs now"

    def test_hyperlink_label_keeps_its_own_formatting(self):
        pres = build([{"type": "content", "title": "L",
                       "body": "- **[bold link](https://b.co)**"}])
        runs = reload_presentation(pres).slides[0].placeholders[1].text_frame.paragraphs[0].runs
        assert runs[0].font.bold
        assert runs[0].hyperlink.address == "https://b.co"

    def test_bracket_text_that_is_not_a_link_is_left_alone(self):
        pres = build([{"type": "content", "title": "L", "body": "- see [x] (y) here"}])
        paragraph = reload_presentation(pres).slides[0].placeholders[1].text_frame.paragraphs[0]
        assert paragraph.text == "see [x] (y) here"
        assert not any(r.hyperlink.address for r in paragraph.runs)

    def test_markdown_inside_a_table_cell(self):
        pres = build([{"type": "table", "title": "T",
                       "rows": [["Plan", "Price"], ["Pro", "**€29**"]]}])
        table = [s for s in reload_presentation(pres).slides[0].shapes if s.has_table][0].table
        cell = table.cell(1, 1)
        assert cell.text == "€29"
        assert any(run.font.bold for run in cell.text_frame.paragraphs[0].runs)

    def test_header_styling_survives_cell_markdown(self):
        pres = build([{"type": "table", "title": "T",
                       "rows": [["Plan", "**Price**"], ["Pro", "€29"]]}])
        table = [s for s in reload_presentation(pres).slides[0].shapes if s.has_table][0].table
        header = table.cell(0, 1)
        assert header.text == "Price"
        assert header.text_frame.paragraphs[0].font.bold


# =============================================================================
# Untitled quote
# =============================================================================

class TestUntitledQuote:

    def test_no_empty_title_placeholder_is_left_behind(self):
        """It showed as "Click to add title" the moment anyone edited the deck."""
        pres = build([{"type": "quote", "text": "Be water"}])
        slide = reload_presentation(pres).slides[0]
        empty = [ph for ph in slide.placeholders if not ph.text_frame.text.strip()]
        assert empty == []
        assert "Be water" in all_text(slide)

    def test_a_titled_quote_still_gets_a_title(self):
        pres = build([{"type": "quote", "title": "On change", "text": "Be water"}])
        slide = reload_presentation(pres).slides[0]
        assert slide.shapes.title is not None
        assert slide.shapes.title.text == "On change"


# =============================================================================
# The whole deck together
# =============================================================================

def test_a_full_deck_of_every_type_builds_cleanly():
    deck = [
        {"type": "title", "title": "Q3 Review", "subtitle": "Board"},
        {"type": "agenda"},
        {"type": "section", "title": "Results"},
        {"type": "kpi", "title": "At a glance", "items": [
            {"value": "€4.2M", "label": "ARR", "delta": "+12%"},
            {"value": "18%", "label": "Churn"}]},
        {"type": "content", "title": "Highlights", "body": "- **Up** 12%\n  - EMEA"},
        {"type": "table", "title": "Pricing", "rows": [["Plan", "Price"], ["Pro", 29]]},
        {"type": "chart", "title": "Revenue", "chart_type": "column",
         "categories": ["Q1", "Q2"], "series": [{"name": "2026", "values": [1, 2]}],
         "body": "- Doubled"},
        {"type": "scatter", "title": "Dose", "series": [{"name": "t", "points": [[1, 2]]}]},
        {"type": "image", "title": "Fig", "source": PNG_DATA_URI, "caption": "Fig 1"},
        {"type": "two_column", "title": "Split",
         "left": {"heading": "Now", "body": "- a"}, "right": {"heading": "Next", "body": "- b"}},
        {"type": "quote", "text": "Be water"},
        {"type": "section", "title": "Outlook"},
        {"type": "timeline", "title": "Roadmap", "steps": [
            {"label": "Build", "detail": "Q1"}, {"label": "Launch", "detail": "Q2"}]},
        {"type": "closing", "subtitle": "Questions?"},
    ]

    pres = build(deck, language="en-GB", author="Tester")
    rebuilt = reload_presentation(pres)

    assert len(rebuilt.slides) == len(deck)
    assert pres.warnings == []
    # The agenda picked up both section slides.
    assert "1.  Results" in all_text(rebuilt.slides[1])
    assert "2.  Outlook" in all_text(rebuilt.slides[1])
