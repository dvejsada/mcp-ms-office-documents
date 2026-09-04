"""Admin UI support for PowerPoint templates (Phase 3b).

A presentation template is not the same kind of thing as a Word or email
template — it declares no arguments and becomes no tool of its own — so most of
what these tests assert is that the pptx path says and does something different
from the other two, rather than reusing machinery that would be meaningless.
"""
import asyncio
import io
import re
import sys
import tempfile
import zipfile
from pathlib import Path

project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

import pytest
from pptx import Presentation as PptxReader
from pptx.oxml.ns import qn as pptx_qn
from starlette.testclient import TestClient

import admin.store as store_mod
import metrics
import pptx_tools.templates as templates_mod
import template_utils as tu
from admin.analysis import PptxAnalysis, analyze, analyze_pptx
from admin.preview import SAMPLE_DECK, render_pptx_preview
from admin.store import KIND_PPTX, TemplateStoreError, kind_meta, validate_asset_filename
from config import Config

TEMPLATE_16_9 = project_root / "default_templates" / "default_pptx_template_16_9.pptx"
TEMPLATE_4_3 = project_root / "default_templates" / "default_pptx_template_4_3.pptx"


@pytest.fixture
def pptx_bytes():
    return TEMPLATE_16_9.read_bytes()


def _template_with_a_duplicate_role() -> bytes:
    """A template where two layouts classify as the same role.

    Neither shipped template has one: their eight detected layouts map one to
    one onto the eight roles, so first-wins and last-wins tie-breaking are
    indistinguishable on them. Turning a vertical-text layout horizontal gives
    it the TITLE+BODY signature of a section header — which is precisely the
    collision ``_is_vertical`` exists to prevent — and so produces the tie a
    tie-breaking test needs.
    """
    prs = PptxReader(str(TEMPLATE_16_9))
    for placeholder in prs.slide_layouts[9].placeholders:      # 'Nadpis a svislý text'
        body_pr = placeholder.text_frame._txBody.find(pptx_qn("a:bodyPr"))
        if body_pr is not None and body_pr.get("vert"):
            body_pr.set("vert", "horz")
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _as_potx(data: bytes) -> bytes:
    """Rewrite a .pptx into the .potx content type, as PowerPoint would."""
    src = io.BytesIO(data)
    out = io.BytesIO()
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            blob = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                blob = blob.decode("utf-8").replace(
                    "presentationml.presentation.main+xml",
                    "presentationml.template.main+xml",
                ).encode("utf-8")
            zout.writestr(item, blob)
    return out.getvalue()


# =============================================================================
# Store
# =============================================================================

class TestStoreKind:

    def test_pptx_is_a_known_kind(self):
        meta = kind_meta(KIND_PPTX)
        assert meta["subdir"] == "pptx_templates.d"
        assert meta["path_key"] == "pptx_path"

    def test_potx_is_accepted_alongside_pptx(self):
        """A designer's brand deck routinely arrives as a .potx."""
        assert validate_asset_filename("brand.pptx", KIND_PPTX) == "brand.pptx"
        assert validate_asset_filename("brand.potx", KIND_PPTX) == "brand.potx"

    def test_a_wrong_extension_is_still_refused(self):
        with pytest.raises(TemplateStoreError):
            validate_asset_filename("brand.docx", KIND_PPTX)

    def test_a_pptx_is_not_accepted_for_the_docx_kind(self):
        with pytest.raises(TemplateStoreError):
            validate_asset_filename("brand.pptx", "docx")

    def test_directory_traversal_is_refused(self):
        with pytest.raises(TemplateStoreError):
            validate_asset_filename("../../etc/brand.pptx", KIND_PPTX)


# =============================================================================
# Analysis
# =============================================================================

class TestAnalyzePptx:

    def test_reports_aspect_and_slide_size(self, pptx_bytes):
        a = analyze_pptx(pptx_bytes)
        assert a.aspect == "16:9"
        assert "13.33" in a.slide_size

    def test_four_three_template_reads_as_four_three(self):
        assert analyze_pptx(TEMPLATE_4_3.read_bytes()).aspect == "4:3"

    def test_every_layout_is_listed_with_its_placeholders(self, pptx_bytes):
        a = analyze_pptx(pptx_bytes)
        assert len(a.layouts) == len(PptxReader(str(TEMPLATE_16_9)).slide_layouts)
        title = [layout for layout in a.layouts if layout.role == "title"][0]
        assert "CENTER_TITLE" in title.placeholders
        assert "SUBTITLE" in title.placeholders

    def test_roles_are_detected_for_the_shipped_template(self, pptx_bytes):
        """The shipped deck covers every role, so nothing falls back by position."""
        a = analyze_pptx(pptx_bytes)
        assert a.missing_roles == []
        assert set(a.role_map) >= {"title", "section", "content", "two_column", "blank"}

    def test_theme_fonts_and_colours_are_read(self, pptx_bytes):
        """Regression: the theme is a generic Part with no parsed element.

        Reaching for ``_element`` returned None and silently produced an empty
        palette — a feature that looks present and does nothing.
        """
        a = analyze_pptx(pptx_bytes)
        assert a.theme_fonts.get("body")
        assert a.theme_colors.get("accent1", "").startswith("#")
        assert len(a.theme_colors["accent1"]) == 7

    def test_system_colours_resolve_through_last_clr(self, pptx_bytes):
        """dk1/lt1 are sysClr, not srgbClr; without lastClr they read as blank."""
        colors = analyze_pptx(pptx_bytes).theme_colors
        assert colors.get("dk1") == "#000000"
        assert colors.get("lt1") == "#FFFFFF"

    def test_unrecognised_layouts_are_reported_not_guessed(self, pptx_bytes):
        """Vertical-text and content-with-caption are deliberate refusals."""
        a = analyze_pptx(pptx_bytes)
        unrecognised = [layout.name for layout in a.layouts if layout.role is None]
        assert unrecognised
        assert any("match no known role" in w for w in a.warnings)

    def test_a_potx_is_analysed_like_a_pptx(self, pptx_bytes):
        assert analyze_pptx(_as_potx(pptx_bytes)).aspect == "16:9"

    def test_a_file_that_is_not_a_presentation_warns_rather_than_raising(self):
        a = analyze_pptx(b"this is not a presentation")
        assert any("Could not open" in w for w in a.warnings)
        assert a.layouts == []

    def test_dispatch_by_kind_returns_the_pptx_shape(self, pptx_bytes):
        assert isinstance(analyze("pptx", pptx_bytes), PptxAnalysis)

    def test_layout_names_drive_the_role_dropdowns(self, pptx_bytes):
        a = analyze_pptx(pptx_bytes)
        assert a.layout_names == [layout.name for layout in a.layouts]

    @pytest.mark.parametrize("path", [TEMPLATE_16_9, TEMPLATE_4_3])
    def test_the_reported_role_map_is_what_the_builder_will_actually_use(self, path):
        """The UI's claim must match the resolver's behaviour, not merely resemble it.

        Both pick the first layout classifying as each role, but that is two
        separate implementations of the same rule in two modules. If either
        side's tie-breaking drifts, the admin page would report a mapping the
        tool does not honour — worse than reporting nothing, because it looks
        authoritative.
        """
        from pptx_tools.layouts import ROLES, LayoutResolver
        from pptx_tools.templates import open_template

        reported = analyze_pptx(path.read_bytes()).role_map
        resolver = LayoutResolver(open_template(path))

        for role in ROLES:
            layout, warning = resolver.resolve(role)
            if role in reported:
                assert reported[role] == layout.name, f"{role} disagrees"
                assert warning is None, f"{role} resolved with a warning: {warning}"
            else:
                # A role the UI reports as missing must be one the resolver
                # cannot satisfy either — it falls back by position and says so.
                assert warning is not None, f"{role} was omitted but resolves cleanly"

    def test_a_contested_role_reports_the_layout_the_resolver_picks(self):
        """The agreement test above cannot see tie-breaking; this one can.

        On the shipped templates every detected role has exactly one layout, so
        first-wins and last-wins agree and a divergence would go unnoticed. With
        two layouts claiming ``section``, only the earlier one is correct.
        """
        from pptx_tools.layouts import LayoutResolver
        from pptx_tools.templates import open_template

        data = _template_with_a_duplicate_role()
        analysis = analyze_pptx(data)

        claimants = [layout.name for layout in analysis.layouts if layout.role == "section"]
        assert len(claimants) == 2, "fixture no longer produces a contested role"

        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "tie.pptx"
            path.write_bytes(data)
            layout, warning = LayoutResolver(open_template(path)).resolve("section")

        assert analysis.role_map["section"] == claimants[0]
        assert analysis.role_map["section"] == layout.name
        assert warning is None


# =============================================================================
# Preview
# =============================================================================

class TestPptxPreview:

    def test_builds_the_sample_deck_on_the_template(self):
        data, warnings = render_pptx_preview(TEMPLATE_16_9, {"name": "brand"})
        doc = PptxReader(io.BytesIO(data))
        assert len(doc.slides) == len(SAMPLE_DECK)
        assert warnings == []

    def test_each_slide_lands_on_the_layout_its_role_asks_for(self):
        data, _ = render_pptx_preview(TEMPLATE_16_9, {"name": "brand"})
        doc = PptxReader(io.BytesIO(data))
        # title, section, content: the three roles with an unambiguous layout.
        assert doc.slides[0].slide_layout.name == "Úvodní snímek"
        assert doc.slides[1].slide_layout.name == "Záhlaví oddílu"
        assert doc.slides[2].slide_layout.name == "Nadpis a obsah"

    def test_an_unsaved_layout_override_is_honoured(self):
        """The whole point of previewing before saving."""
        data, _ = render_pptx_preview(
            TEMPLATE_16_9, {"name": "brand", "layouts": {"content": "Obsah s titulkem"}}
        )
        doc = PptxReader(io.BytesIO(data))
        assert doc.slides[2].slide_layout.name == "Obsah s titulkem"

    def test_defaults_are_applied_to_the_preview(self):
        data, _ = render_pptx_preview(
            TEMPLATE_16_9,
            {"name": "brand", "defaults": {"footer_text": "ACME · Confidential"}},
        )
        doc = PptxReader(io.BytesIO(data))
        text = "\n".join(
            shape.text_frame.text
            for shape in doc.slides[2].shapes if shape.has_text_frame
        )
        assert "ACME · Confidential" in text

    def test_the_preview_follows_the_templates_own_aspect(self):
        data, _ = render_pptx_preview(TEMPLATE_4_3, {"name": "legacy"})
        doc = PptxReader(io.BytesIO(data))
        assert doc.slide_width < doc.slide_height * 1.45

    def test_template_sample_slides_do_not_leak_into_the_preview(self):
        data, _ = render_pptx_preview(TEMPLATE_16_9, {"name": "brand"})
        assert len(PptxReader(io.BytesIO(data)).slides) == len(SAMPLE_DECK)


# =============================================================================
# Admin routes
# =============================================================================

@pytest.fixture
def admin_client(tmp_path, monkeypatch):
    """A logged-in client whose store *and* pptx registry live under tmp_path."""
    custom = tmp_path / "custom"
    cfg = tmp_path / "config"
    custom.mkdir()
    cfg.mkdir()

    monkeypatch.setattr(store_mod, "_APP_CUSTOM_DIR", tmp_path / "nope_custom")
    monkeypatch.setattr(store_mod, "_APP_CONFIG_DIR", tmp_path / "nope_config")
    monkeypatch.setattr(store_mod, "_LOCAL_CUSTOM_DIR", custom)
    monkeypatch.setattr(store_mod, "_LOCAL_CONFIG_DIR", cfg)
    monkeypatch.setattr(tu, "APP_CUSTOM_DIR", tmp_path / "nope_custom")
    monkeypatch.setattr(tu, "LOCAL_CUSTOM_DIR", custom)
    # The pptx registry resolves its config dir independently of the store.
    monkeypatch.setattr(templates_mod, "_APP_CONFIG_DIR", tmp_path / "nope_config")
    monkeypatch.setattr(templates_mod, "_LOCAL_CONFIG_DIR", cfg)
    templates_mod.clear_cache()

    monkeypatch.setenv("ADMIN_ENABLED", "true")
    monkeypatch.setenv("ADMIN_PASSWORD", "pw")
    monkeypatch.delenv("API_KEY", raising=False)

    from fastmcp import FastMCP
    from admin.app import build_combined_app

    mcp = FastMCP("test-admin-pptx")
    client = TestClient(build_combined_app(mcp, Config.from_env()))
    metrics.reset()
    client.__enter__()
    client.post("/admin/login", data={"password": "pw"})
    yield client, cfg, custom, mcp
    client.__exit__(None, None, None)
    metrics.reset()
    templates_mod.clear_cache()


def _tool_names(mcp) -> list:
    """The MCP tools currently registered, from a synchronous test body."""
    return sorted(t.name for t in asyncio.run(mcp.list_tools()))


def _csrf(client) -> str:
    html = client.get("/admin/new/pptx").text
    tag = re.search(r'<input[^>]*name="csrf"[^>]*>', html)
    assert tag, "no CSRF input found"
    return re.search(r'value="([^"]*)"', tag.group(0)).group(1)


def _post(client, url, data=None, files=None):
    payload = dict(data or {})
    payload["csrf"] = _csrf(client)
    return client.post(url, data=payload, files=files)


class TestAdminRoutes:

    def test_the_index_lists_a_powerpoint_section(self, admin_client):
        client, _cfg, _custom, _mcp = admin_client
        html = client.get("/admin/").text
        assert "PowerPoint templates" in html

    def test_the_new_page_explains_it_is_a_design_not_a_tool(self, admin_client):
        """The distinction the whole phase turns on; say it where it matters."""
        client, _cfg, _custom, _mcp = admin_client
        html = client.get("/admin/new/pptx").text
        assert "design" in html
        assert ".pptx,.potx" in html

    def test_upload_reports_the_layouts_and_offers_role_overrides(self, admin_client, pptx_bytes):
        client, _cfg, custom, _mcp = admin_client
        r = _post(client, "/admin/pptx/draft", data={"name": "brand"},
                  files={"file": ("brand.pptx", pptx_bytes,
                                  "application/vnd.openxmlformats-officedocument"
                                  ".presentationml.presentation")})
        assert r.status_code == 200
        assert "Úvodní snímek" in r.text          # a real layout name
        assert "layout_title" in r.text            # the role dropdown
        assert "Aptos" in r.text                   # theme font
        assert (custom / "brand.pptx").is_file()

    def test_the_upload_form_offers_no_arguments(self, admin_client, pptx_bytes):
        """A presentation template takes none; offering the editor would lie."""
        client, _cfg, _custom, _mcp = admin_client
        r = _post(client, "/admin/pptx/draft", data={"name": "brand"},
                  files={"file": ("brand.pptx", pptx_bytes, "application/vnd.ms-powerpoint")})
        assert "Add argument" not in r.text
        assert 'name="arg_name"' not in r.text

    def test_a_potx_upload_keeps_its_extension(self, admin_client, pptx_bytes):
        client, _cfg, custom, _mcp = admin_client
        r = _post(client, "/admin/pptx/draft", data={"name": "brand"},
                  files={"file": ("brand.potx", _as_potx(pptx_bytes),
                                  "application/vnd.openxmlformats-officedocument"
                                  ".presentationml.template")})
        assert r.status_code == 200
        assert (custom / "brand.potx").is_file()
        assert not (custom / "brand.pptx").exists()

    def test_a_file_that_is_not_a_presentation_is_rejected(self, admin_client):
        client, _cfg, custom, _mcp = admin_client
        r = _post(client, "/admin/pptx/draft", data={"name": "brand"},
                  files={"file": ("brand.pptx", b"not a deck", "application/octet-stream")})
        assert "Could not open" in r.text
        assert not (custom / "brand.pptx").exists()

    def test_saving_writes_a_spec_the_registry_picks_up(self, admin_client, pptx_bytes):
        client, cfg, custom, _mcp = admin_client
        (custom / "brand.pptx").write_bytes(pptx_bytes)

        r = _post(client, "/admin/pptx/save", data={
            "name": "brand", "asset_filename": "brand.pptx",
            "description": "Brand deck", "is_default": "1", "strip_slides": "1",
            "layout_content": "Obsah s titulkem",
            "default_footer_text": "ACME · Confidential",
            "default_language": "cs-CZ", "default_slide_numbers": "1",
        })
        assert r.status_code == 200

        spec_file = cfg / "pptx_templates.d" / "brand.yaml"
        assert spec_file.is_file()
        import yaml
        spec = yaml.safe_load(spec_file.read_text())
        assert spec["pptx_path"] == "brand.pptx"
        assert spec["default"] is True
        assert spec["layouts"] == {"content": "Obsah s titulkem"}
        assert spec["defaults"] == {
            "footer_text": "ACME · Confidential",
            "language": "cs-CZ",
            "show_slide_numbers": True,
        }

    def test_a_saved_template_becomes_usable_without_a_restart(self, admin_client, pptx_bytes):
        """"Live" for pptx means the registry re-read it, not that a tool exists."""
        client, _cfg, custom, _mcp = admin_client
        (custom / "brand.pptx").write_bytes(pptx_bytes)
        assert "brand" not in templates_mod.template_names()

        _post(client, "/admin/pptx/save",
              data={"name": "brand", "asset_filename": "brand.pptx", "strip_slides": "1"})

        assert "brand" in templates_mod.template_names()

    def test_saving_does_not_create_an_mcp_tool(self, admin_client, pptx_bytes):
        """One presentation tool exists; templates are its argument, not tools.

        A Word template of the same name would appear in the tool list. This
        asserts against the live MCP tool registry, not against page text.
        """
        client, _cfg, custom, mcp = admin_client
        (custom / "brand.pptx").write_bytes(pptx_bytes)

        names_before = _tool_names(mcp)
        _post(client, "/admin/pptx/save",
              data={"name": "brand", "asset_filename": "brand.pptx", "strip_slides": "1"})

        assert _tool_names(mcp) == names_before
        assert "brand" not in _tool_names(mcp)

    def test_the_save_message_does_not_promise_a_tool(self, admin_client, pptx_bytes):
        client, _cfg, custom, _mcp = admin_client
        (custom / "brand.pptx").write_bytes(pptx_bytes)
        r = _post(client, "/admin/pptx/save",
                  data={"name": "brand", "asset_filename": "brand.pptx", "strip_slides": "1"})

        assert "one of the templates the presentation tool can build on" in r.text
        assert "is now live and ready for the AI to use" not in r.text

    def test_saving_with_no_uploaded_file_is_refused_outright(self, admin_client):
        """The store refuses before the registry is ever consulted."""
        client, cfg, _custom, _mcp = admin_client
        r = _post(client, "/admin/pptx/save",
                  data={"name": "ghost", "asset_filename": "ghost.pptx", "strip_slides": "1"})

        assert "Save failed" in r.text
        assert "does not exist yet" in r.text
        assert not (cfg / "pptx_templates.d" / "ghost.yaml").exists()

    def test_a_spec_the_registry_cannot_resolve_is_not_called_live(self):
        """The store and the registry search different directory sets.

        A file present to the store but absent from the template search path is
        a misconfiguration, not a success — the admin must not be told the
        template is usable when the tool will not find it.
        """
        from admin.app import AdminContext

        templates_mod.clear_cache()
        assert AdminContext._refresh_pptx_registry("definitely_not_registered") is False

    def test_unticked_checkboxes_are_recorded_as_off(self, admin_client, pptx_bytes):
        """An HTML checkbox sends nothing when unticked; absence must mean False."""
        client, cfg, custom, _mcp = admin_client
        (custom / "brand.pptx").write_bytes(pptx_bytes)
        _post(client, "/admin/pptx/save",
              data={"name": "brand", "asset_filename": "brand.pptx"})

        import yaml
        spec = yaml.safe_load((cfg / "pptx_templates.d" / "brand.yaml").read_text())
        assert spec["strip_slides"] is False
        assert "default" not in spec
        assert "defaults" not in spec

    def test_preview_returns_a_deck(self, admin_client, pptx_bytes):
        client, _cfg, custom, _mcp = admin_client
        (custom / "brand.pptx").write_bytes(pptx_bytes)
        r = _post(client, "/admin/pptx/preview",
                  data={"name": "brand", "asset_filename": "brand.pptx", "strip_slides": "1"})

        assert r.status_code == 200
        assert "presentationml.presentation" in r.headers["content-type"]
        assert len(PptxReader(io.BytesIO(r.content)).slides) == len(SAMPLE_DECK)

    def test_editing_shows_what_was_saved(self, admin_client, pptx_bytes):
        client, _cfg, custom, _mcp = admin_client
        (custom / "brand.pptx").write_bytes(pptx_bytes)
        _post(client, "/admin/pptx/save", data={
            "name": "brand", "asset_filename": "brand.pptx",
            "description": "Brand deck", "layout_content": "Obsah s titulkem",
            "default_language": "cs-CZ", "strip_slides": "1",
        })
        html = client.get("/admin/pptx/brand/edit").text
        assert "Brand deck" in html
        assert "cs-CZ" in html
        assert 'value="Obsah s titulkem" selected' in html

    def test_deleting_removes_it_from_the_registry(self, admin_client, pptx_bytes):
        client, _cfg, custom, _mcp = admin_client
        (custom / "brand.pptx").write_bytes(pptx_bytes)
        _post(client, "/admin/pptx/save",
              data={"name": "brand", "asset_filename": "brand.pptx", "strip_slides": "1"})
        assert "brand" in templates_mod.template_names()

        _post(client, "/admin/pptx/brand/delete")

        assert "brand" not in templates_mod.template_names()
        # The source file is kept, matching the other kinds.
        assert (custom / "brand.pptx").is_file()
