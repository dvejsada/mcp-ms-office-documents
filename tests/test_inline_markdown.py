"""The inline-markdown grammar shared by the Word and PowerPoint renderers.

Before this module existed each renderer had its own regex and the two had
drifted, each carrying a bug the other had fixed. These tests pin the grammar
once, prove the two renderers now tokenise identically on everything both can
draw, and pin the behaviour changes each side received.
"""
import io
import sys
from pathlib import Path

project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

import pytest
from docx import Document
from pptx import Presentation as PptxReader

from inline_markdown import ESCAPE_RE, LINK_RE, build_inline_pattern

BS = "\\"


def tokens(pattern, text):
    return [t for t in pattern.split(text) if t]


FULL = build_inline_pattern(highlight=True, superscript=True, subscript=True)
PPTX = build_inline_pattern(superscript=True, subscript=True)


# =============================================================================
# Grammar
# =============================================================================

class TestFlanking:
    """CommonMark's rule: a marker next to whitespace is not a marker."""

    @pytest.mark.parametrize("text", [
        "5 * 3 * 2 = 30", "10 ~ 20 ~ 30", "* a * b *", "a * b", "x ~ y",
        "== not hl ==", "~~ x ~~", "** **", "2 * 3",
    ])
    def test_arithmetic_and_spaced_markers_are_plain_text(self, text):
        assert tokens(FULL, text) == [text]
        assert FULL.search(text) is None

    @pytest.mark.parametrize("text", ["**b**", "*i*", "~~s~~", "__u__", "`c`", "^s^", "~s~", "==h=="])
    def test_flanked_spans_match(self, text):
        assert FULL.search(text) is not None
        assert tokens(FULL, text) == [text]


class TestNesting:

    def test_bold_ending_in_italic_is_one_bold_span(self):
        """Was mis-parsed by the PowerPoint grammar, leaving a stray '*'."""
        assert tokens(FULL, "**a *b***") == ["**a *b***"]

    def test_italic_ending_in_bold_is_one_italic_span(self):
        """Was dropped by a closing lookbehind; the Phase 2 regression."""
        assert tokens(FULL, "*always **backup your data***") == ["*always **backup your data***"]

    def test_adjacent_bold_spans_stay_separate(self):
        assert tokens(FULL, "**a**b**c**") == ["**a**", "b", "**c**"]

    def test_triple_star_is_bold_italic_not_bold_plus_star(self):
        assert tokens(FULL, "***bi***") == ["***bi***"]


class TestEscapes:

    def test_punctuation_is_escapable(self):
        assert ESCAPE_RE.sub(lambda m: m.group(1), "price " + BS + "* qty") == "price * qty"
        assert ESCAPE_RE.sub(lambda m: m.group(1), "1" + BS + ".5") == "1.5"

    @pytest.mark.parametrize("text", ["C:" + BS + "new", "a" + BS + "nb", BS + "t", "x" + BS + "y"])
    def test_a_backslash_before_a_letter_is_not_an_escape(self, text):
        """The PowerPoint escape swallowed any character: C:\\new became C:new."""
        assert ESCAPE_RE.sub(lambda m: m.group(1), text) == text


class TestLinks:

    def test_link_token_and_anchored_parse_agree(self):
        assert tokens(FULL, "see [docs](https://x.y/z) now") == ["see ", "[docs](https://x.y/z)", " now"]
        assert LINK_RE.match("[docs](https://x.y/z)").groups() == ("docs", "https://x.y/z")

    @pytest.mark.parametrize("text", ["[](https://x)", "[a](b c)"])
    def test_degenerate_links_are_plain_text(self, text):
        assert tokens(FULL, text) == [text]


class TestPerTargetFlags:

    def test_highlight_is_only_matched_when_asked_for(self):
        """The split output cannot show this — a single unmatched token splits
        to itself — so it is asserted with search."""
        assert FULL.search("==hl==") is not None
        assert PPTX.search("==hl==") is None

    def test_the_two_renderers_agree_on_everything_both_draw(self):
        corpus = [
            "**bold**", "*it*", "***bi***", "~~s~~", "__u__", "`c`", "[a](http://x)",
            "5 * 3 * 2 = 30", "10 ~ 20 ~ 30", "__init__", "**bold *it* end**",
            "*it **b** end*", "*always **backup your data***", "^sup^", "~sub~",
            "**a**b**c**", "* not italic *", "**", "***", "a**b", "[**b**](http://x)",
            "**[l](http://x)**", "* a * b *", "_single_", "*a*b*c*", "**a *b***",
            "x^2^", "H~2~O",
        ]
        for text in corpus:
            assert tokens(FULL, text) == tokens(PPTX, text), text


# =============================================================================
# End to end, both renderers
# =============================================================================

class TestRendered:

    @staticmethod
    def _pptx_runs(text):
        from pptx_tools.slide_builder import PowerpointPresentation
        pres = PowerpointPresentation([{"type": "content", "title": "T", "body": "- " + text}], "16:9")
        para = PptxReader(pres.save()).slides[0].placeholders[1].text_frame.paragraphs[0]
        return para.runs

    def test_pptx_superscript_and_subscript_set_the_baseline(self):
        """Asserted by run order: both scripts here read "2", so keying by
        text would collapse them into one."""
        runs = self._pptx_runs("E = mc^2^ and H~2~O")
        rpr = "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
        baselines = [(r.text, (r._r.find(rpr).get("baseline") if r._r.find(rpr) is not None else None))
                     for r in runs]
        assert baselines == [
            ("E = mc", None), ("2", "30000"), (" and H", None), ("2", "-25000"), ("O", None),
        ]

    def test_pptx_keeps_a_windows_path_intact(self):
        runs = self._pptx_runs("saved to C:" + BS + "new" + BS + "temp")
        assert "".join(r.text for r in runs) == "saved to C:" + BS + "new" + BS + "temp"

    def test_docx_no_longer_italicises_arithmetic(self):
        """Behaviour change for Word, deliberate: it inherits the flanking rules."""
        from docx_tools.inline_formatting import parse_inline_formatting

        doc = Document()
        para = doc.add_paragraph()
        parse_inline_formatting("5 * 3 * 2 = 30", para)
        assert not any(r.italic for r in para.runs)
        assert "".join(r.text for r in para.runs) == "5 * 3 * 2 = 30"

    def test_docx_subscript_still_renders_when_flanked(self):
        from docx_tools.inline_formatting import parse_inline_formatting

        doc = Document()
        para = doc.add_paragraph()
        parse_inline_formatting("H~2~O", para)
        assert [r.text for r in para.runs] == ["H", "2", "O"]
        assert para.runs[1].font.subscript is True


# =============================================================================
# Package boundaries
# =============================================================================

class TestSharedModules:

    def test_neither_document_package_imports_the_other(self):
        """image_utils moved to the root so docx_tools stopped importing pptx_tools."""
        for package, other in (("docx_tools", "pptx_tools"), ("pptx_tools", "docx_tools")):
            offenders = [
                str(path.relative_to(project_root))
                for path in (project_root / package).glob("*.py")
                if f"from {other}" in path.read_text(encoding="utf-8")
                or f"import {other}" in path.read_text(encoding="utf-8")
            ]
            assert offenders == [], f"{package} imports {other}: {offenders}"

    def test_image_utils_is_importable_from_the_root(self):
        import image_utils
        assert callable(image_utils.load_image)
